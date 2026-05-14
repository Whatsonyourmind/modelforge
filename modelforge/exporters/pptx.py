"""PowerPoint exporter — executive committee deck from a ModelForge workbook.

Generates a 5-slide deck:
    1. Cover (deal name, date, version, author)
    2. Key Assumptions (top-10 named-range drivers, with source citations)
    3. Sources (the Sources sheet, rendered as a citation table)
    4. Key Outputs (IRR / MOIC / DSCR / Leverage — point estimate + scenario)
    5. QC Pass (12-check gate result)

Bulge-tier color convention preserved:
    Blue   = hardcoded input
    Black  = formula
    Green  = cross-sheet link
    Red    = warning / fail
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    raise ImportError(
        "pptx exporter requires `pip install modelforge[export]` "
        "(adds python-pptx)."
    ) from e

from openpyxl import load_workbook


# Bulge-tier color palette
BLUE_INPUT = RGBColor(0x00, 0x00, 0xFF)
BLACK_FORMULA = RGBColor(0x00, 0x00, 0x00)
GREEN_LINK = RGBColor(0x00, 0x80, 0x00)
RED_WARN = RGBColor(0xC0, 0x00, 0x00)
GREY_LABEL = RGBColor(0x55, 0x55, 0x55)
NAVY_HEADER = RGBColor(0x10, 0x2C, 0x57)


def _read_workbook_meta(xlsx_path: Path) -> dict[str, Any]:
    """Extract the metadata needed to populate the deck.

    Pulls (in order of preference):
        - Cover sheet — deal name, version, author, date
        - Assumptions sheet — top-10 driver values
        - Sources sheet — list of S-001..S-NNN entries
        - QC sheet — pass/fail status per check
        - Returns sheet — IRR / MOIC / DSCR if present
    """
    wb = load_workbook(xlsx_path, data_only=True)
    meta: dict[str, Any] = {
        "deal_name": xlsx_path.stem,
        "version": "v0.8.9",
        "author": "ModelForge",
        "date": "",
        "assumptions": [],
        "sources": [],
        "qc_checks": [],
        "outputs": {},
    }

    # Cover
    if "Cover" in wb.sheetnames:
        cov = wb["Cover"]
        for row in cov.iter_rows(min_row=1, max_row=15, values_only=True):
            if not row or row[0] is None:
                continue
            label = str(row[0]).strip().lower()
            value = row[1] if len(row) > 1 else None
            if "deal" in label and value:
                meta["deal_name"] = str(value)
            elif "version" in label and value:
                meta["version"] = str(value)
            elif "author" in label and value:
                meta["author"] = str(value)
            elif "date" in label and value:
                meta["date"] = str(value)

    # Assumptions
    if "Assumptions" in wb.sheetnames:
        asm = wb["Assumptions"]
        for row in asm.iter_rows(min_row=2, max_row=30, values_only=True):
            if not row or row[0] is None:
                continue
            name = str(row[0])
            value = row[1] if len(row) > 1 else None
            source = row[2] if len(row) > 2 else None
            if name and value is not None:
                meta["assumptions"].append({
                    "name": name,
                    "value": value,
                    "source": source,
                })
            if len(meta["assumptions"]) >= 10:
                break

    # Sources
    if "Sources" in wb.sheetnames:
        srcs = wb["Sources"]
        for row in srcs.iter_rows(min_row=2, max_row=50, values_only=True):
            if not row or row[0] is None:
                continue
            sid, doc, page, publisher, *rest = list(row) + [None] * 8
            if sid and doc:
                meta["sources"].append({
                    "id": str(sid),
                    "doc": str(doc),
                    "page": page,
                    "publisher": publisher,
                })

    # QC
    if "QC" in wb.sheetnames:
        qc = wb["QC"]
        for row in qc.iter_rows(min_row=2, max_row=30, values_only=True):
            if not row or row[0] is None:
                continue
            check, result = row[0], row[1] if len(row) > 1 else None
            meta["qc_checks"].append({"check": str(check), "result": result})

    # Returns / Outputs (look at common sheet names)
    for sheet_name in ("Returns", "Outputs", "Summary", "KPIs"):
        if sheet_name in wb.sheetnames:
            sh = wb[sheet_name]
            for row in sh.iter_rows(min_row=1, max_row=20, values_only=True):
                if not row or row[0] is None:
                    continue
                label = str(row[0]).strip()
                value = row[1] if len(row) > 1 else None
                if label and value is not None and isinstance(value, (int, float)):
                    meta["outputs"][label] = value
            if meta["outputs"]:
                break

    return meta


def _add_title_slide(prs: Presentation, meta: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Deal name (large)
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(9.0), Inches(1.5))
    tf = tx.text_frame
    tf.text = meta["deal_name"]
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY_HEADER

    # Subtitle
    tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(9.0), Inches(0.8))
    tf2 = tx2.text_frame
    tf2.text = "Bulge-Tier Financial Model · Committee Deck"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(20)
    p2.font.color.rgb = GREY_LABEL

    # Footer
    tx3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(9.0), Inches(0.5))
    tf3 = tx3.text_frame
    tf3.text = f"{meta['version']}  ·  {meta['date']}  ·  {meta['author']}  ·  Generated by ModelForge"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(11)
    p3.font.color.rgb = GREY_LABEL


def _add_section_slide(prs: Presentation, title: str, rows: list[tuple[str, str]],
                       cell_color: RGBColor = BLUE_INPUT) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    tf = tx.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY_HEADER

    # Body (table)
    if not rows:
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.5))
        body.text_frame.text = "(no entries)"
        body.text_frame.paragraphs[0].font.color.rgb = GREY_LABEL
        return

    n = min(len(rows), 12)
    table = slide.shapes.add_table(
        rows=n + 1, cols=2,
        left=Inches(0.5), top=Inches(1.2),
        width=Inches(9.0), height=Inches(5.5),
    ).table

    # Header
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Value"
    for c in range(2):
        for p in table.cell(0, c).text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = NAVY_HEADER

    for i, (label, value) in enumerate(rows[:n], start=1):
        table.cell(i, 0).text = str(label)
        table.cell(i, 1).text = str(value)
        # Color the values per the bulge convention (blue = input)
        for p in table.cell(i, 1).text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = cell_color
        for p in table.cell(i, 0).text_frame.paragraphs:
            p.font.size = Pt(11)


def _add_qc_slide(prs: Presentation, qc_checks: list[dict[str, Any]]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    tf = tx.text_frame
    pass_count = sum(1 for c in qc_checks if str(c.get("result", "")).upper() == "PASS")
    total = len(qc_checks)
    tf.text = f"QC Gate — {pass_count}/{total} checks PASS"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY_HEADER if pass_count == total else RED_WARN

    if not qc_checks:
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.5))
        body.text_frame.text = "(no QC sheet found in workbook)"
        body.text_frame.paragraphs[0].font.color.rgb = GREY_LABEL
        return

    n = min(len(qc_checks), 12)
    table = slide.shapes.add_table(
        rows=n + 1, cols=2,
        left=Inches(0.5), top=Inches(1.2),
        width=Inches(9.0), height=Inches(5.5),
    ).table

    table.cell(0, 0).text = "Check"
    table.cell(0, 1).text = "Result"
    for c in range(2):
        for p in table.cell(0, c).text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = NAVY_HEADER

    for i, chk in enumerate(qc_checks[:n], start=1):
        table.cell(i, 0).text = str(chk.get("check", "—"))
        result = str(chk.get("result", "—"))
        table.cell(i, 1).text = result
        color = NAVY_HEADER if result.upper() == "PASS" else RED_WARN
        for p in table.cell(i, 1).text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = color
        for p in table.cell(i, 0).text_frame.paragraphs:
            p.font.size = Pt(11)


def build_committee_deck(xlsx_path: Path, output_path: Path) -> Path:
    """Build a 5-slide committee deck from a ModelForge workbook.

    Args:
        xlsx_path: Source ModelForge workbook (.xlsx).
        output_path: Output .pptx path.

    Returns:
        The output .pptx path.
    """
    xlsx_path = Path(xlsx_path)
    output_path = Path(output_path)

    meta = _read_workbook_meta(xlsx_path)
    prs = Presentation()

    # 1. Cover
    _add_title_slide(prs, meta)

    # 2. Key Assumptions
    rows = [
        (a["name"],
         f"{a['value']:,.2f}" if isinstance(a["value"], (int, float)) else str(a["value"]))
        for a in meta["assumptions"][:10]
    ]
    _add_section_slide(prs, "Key Assumptions", rows, cell_color=BLUE_INPUT)

    # 3. Sources
    src_rows = [
        (s["id"], f'{s["doc"]} (p. {s["page"]})' if s.get("page") else s["doc"])
        for s in meta["sources"][:10]
    ]
    _add_section_slide(prs, "Sources", src_rows, cell_color=GREEN_LINK)

    # 4. Key Outputs
    out_rows = [
        (k, f"{v:,.2f}" if isinstance(v, (int, float)) else str(v))
        for k, v in list(meta["outputs"].items())[:10]
    ]
    _add_section_slide(prs, "Key Outputs", out_rows, cell_color=BLACK_FORMULA)

    # 5. QC Pass
    _add_qc_slide(prs, meta["qc_checks"])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
