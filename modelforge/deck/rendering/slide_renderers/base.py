"""Base finance slide renderer -- abstract base class with shared helpers."""

from __future__ import annotations

import logging
import re
import statistics
from abc import ABC, abstractmethod
from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from modelforge.deck.finance.conditional import ConditionalFormatter
from modelforge.deck.finance.formatter import FinancialFormatter
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import Position
    from modelforge.deck.ir.slides.base import BaseSlide
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)

# Header keywords for auto-detecting column format types.
_MULTIPLE_KEYWORDS = re.compile(
    r"(ev/ebitda|p/e|p/b|ev/revenue|ev/sales|multiple|moic|coc)",
    re.IGNORECASE,
)
_CURRENCY_KEYWORDS = re.compile(
    r"(\$|revenue|market\s*cap|value|price|ebitda|income|earnings|cash\s*flow|debt|equity\s*value|amount|ev\b)",
    re.IGNORECASE,
)
_PERCENTAGE_KEYWORDS = re.compile(
    r"(%|margin|growth|yield|rate|return|irr|cagr|wacc|% of)",
    re.IGNORECASE,
)


def _infer_format(header: str) -> str | None:
    """Infer a NumberFormat string from a column header."""
    if _MULTIPLE_KEYWORDS.search(header):
        return "multiple"
    if _CURRENCY_KEYWORDS.search(header):
        return "currency"
    if _PERCENTAGE_KEYWORDS.search(header):
        return "percentage"
    return None


class BaseFinanceSlideRenderer(ABC):
    """Abstract base class for finance slide renderers.

    Provides shared helpers for title, table creation, median computation,
    and traffic-light indicators used across all 9 finance slide types.
    """

    @abstractmethod
    def render(
        self, slide: Slide, ir_slide: BaseSlide, theme: ResolvedTheme
    ) -> None:
        """Render a finance slide onto a python-pptx slide."""

    # ── Shared helpers ────────────────────────────────────────────────────

    def _add_title(
        self,
        slide: Slide,
        title: str,
        theme: ResolvedTheme,
        position: Position,
    ) -> None:
        """Add a title text box at the given position."""
        left = Inches(position.x or 0)
        top = Inches(position.y or 0)
        width = Inches(position.width or 10)
        height = Inches(position.height or 0.8)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.text = title
        para.alignment = PP_ALIGN.LEFT

        font_name = resolve_font_name(theme.typography.heading_family)
        font_size = Pt(theme.typography.scale.get("h2", 36))

        for run in para.runs:
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)

    def _add_table(
        self,
        slide: Slide,
        headers: list[str],
        rows: list[list],
        theme: ResolvedTheme,
        position: Position,
        column_formats: list[str | None] | None = None,
        highlight_rows: list[int] | None = None,
        footer_row: list | None = None,
        cell_colors: dict[tuple[int, int], str] | None = None,
    ) -> None:
        """Create a table shape with financial formatting.

        Args:
            slide: python-pptx Slide.
            headers: Column header strings.
            rows: List of row data lists.
            theme: Resolved theme.
            position: Table position in inches.
            column_formats: Per-column format string (e.g., "currency", "multiple") or None.
            highlight_rows: Row indices (0-based in body rows) to highlight.
            footer_row: Optional footer row data.
            cell_colors: Optional dict mapping (row_idx, col_idx) to hex color for cell bg.
        """
        num_cols = len(headers)
        num_rows = 1 + len(rows)  # header + body
        if footer_row:
            num_rows += 1

        left = Inches(position.x or 0)
        top = Inches(position.y or 0)
        width = Inches(position.width or 10)
        height = Inches(position.height or 4)

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        font_size = Pt(theme.typography.scale.get("caption", 14))
        font_name = resolve_font_name(theme.typography.body_family)
        highlight_rows = highlight_rows or []
        column_formats = column_formats or [None] * num_cols
        cell_colors = cell_colors or {}

        # Distribute column widths proportionally
        total_len = sum(max(len(str(h)), 3) for h in headers)
        for col_idx, header in enumerate(headers):
            proportion = max(len(str(header)), 3) / total_len
            table.columns[col_idx].width = int(width * proportion)

        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            self._style_cell(
                cell,
                font_size=font_size,
                font_name=font_name,
                bold=True,
                text_color=hex_to_rgb("#FFFFFF"),
                fill_color=hex_to_rgb(theme.colors.primary),
            )

        # Body rows
        for row_idx, row_data in enumerate(rows):
            table_row = row_idx + 1
            is_highlighted = row_idx in highlight_rows

            if is_highlighted:
                default_bg = hex_to_rgb(theme.colors.accent)
            elif row_idx % 2 == 0:
                default_bg = hex_to_rgb(theme.colors.surface)
            else:
                default_bg = hex_to_rgb(theme.colors.background)

            for col_idx in range(num_cols):
                raw_value = row_data[col_idx] if col_idx < len(row_data) else ""
                cell = table.cell(table_row, col_idx)

                # Format value
                fmt = column_formats[col_idx] if col_idx < len(column_formats) else None
                if fmt and isinstance(raw_value, (int, float)):
                    cell.text = FinancialFormatter.auto_format(raw_value, fmt)
                else:
                    cell.text = str(raw_value) if raw_value is not None else ""

                # Determine alignment
                is_numeric = isinstance(raw_value, (int, float)) or (fmt is not None)
                alignment = PP_ALIGN.RIGHT if is_numeric else PP_ALIGN.LEFT

                # Cell-specific color override
                if (row_idx, col_idx) in cell_colors:
                    bg_color = hex_to_rgb(cell_colors[(row_idx, col_idx)])
                else:
                    bg_color = default_bg

                self._style_cell(
                    cell,
                    font_size=font_size,
                    font_name=font_name,
                    bold=is_highlighted,
                    text_color=hex_to_rgb(theme.colors.text_primary),
                    fill_color=bg_color,
                    alignment=alignment,
                )

        # Footer row
        if footer_row:
            footer_table_row = num_rows - 1
            for col_idx in range(num_cols):
                raw_value = footer_row[col_idx] if col_idx < len(footer_row) else ""
                cell = table.cell(footer_table_row, col_idx)

                fmt = column_formats[col_idx] if col_idx < len(column_formats) else None
                if fmt and isinstance(raw_value, (int, float)):
                    cell.text = FinancialFormatter.auto_format(raw_value, fmt)
                else:
                    cell.text = str(raw_value) if raw_value is not None else ""

                is_numeric = isinstance(raw_value, (int, float)) or (fmt is not None)
                alignment = PP_ALIGN.RIGHT if is_numeric else PP_ALIGN.LEFT

                self._style_cell(
                    cell,
                    font_size=font_size,
                    font_name=font_name,
                    bold=True,
                    text_color=hex_to_rgb(theme.colors.text_primary),
                    fill_color=hex_to_rgb(theme.colors.accent),
                    alignment=alignment,
                )

    def _compute_median_row(
        self, rows: list[list], numeric_cols: list[int]
    ) -> list:
        """Compute a median footer row for given numeric column indices."""
        result: list = ["Median"]
        for col_idx in range(1, max(numeric_cols or [0]) + 1):
            if col_idx in numeric_cols:
                values = []
                for row in rows:
                    if col_idx < len(row) and isinstance(row[col_idx], (int, float)):
                        values.append(float(row[col_idx]))
                if values:
                    result.append(statistics.median(values))
                else:
                    result.append("")
            else:
                result.append("")
        return result

    def _add_shape_indicator(
        self,
        slide: Slide,
        status: str,
        position: Position,
    ) -> None:
        """Draw a small circle with traffic-light color."""
        color = ConditionalFormatter.traffic_light(status)
        left = Inches(position.x or 0)
        top = Inches(position.y or 0)
        size = Inches(min(position.width or 0.3, position.height or 0.3, 0.3))

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color)
        shape.line.fill.background()  # No border

    @staticmethod
    def _style_cell(
        cell,
        font_size,
        font_name: str,
        bold: bool,
        text_color,
        fill_color,
        alignment=None,
    ) -> None:
        """Apply styling to a table cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

        for paragraph in cell.text_frame.paragraphs:
            if alignment:
                paragraph.alignment = alignment
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.name = font_name
                run.font.bold = bold
                run.font.color.rgb = text_color

    @staticmethod
    def _find_table_element(elements: list) -> object | None:
        """Find the first TableElement in elements list."""
        for elem in elements:
            elem_type = getattr(elem, "type", None)
            if hasattr(elem_type, "value"):
                elem_type = elem_type.value
            if elem_type == "table":
                return elem
        return None

    @staticmethod
    def _find_table_elements(elements: list) -> list:
        """Find all TableElement instances in elements list."""
        tables = []
        for elem in elements:
            elem_type = getattr(elem, "type", None)
            if hasattr(elem_type, "value"):
                elem_type = elem_type.value
            if elem_type == "table":
                tables.append(elem)
        return tables

    @staticmethod
    def _find_chart_element(elements: list, chart_type: str | None = None) -> object | None:
        """Find the first ChartElement, optionally filtered by chart_type."""
        for elem in elements:
            elem_type = getattr(elem, "type", None)
            if hasattr(elem_type, "value"):
                elem_type = elem_type.value
            if elem_type == "chart":
                if chart_type is None:
                    return elem
                ct = getattr(getattr(elem, "chart_data", None), "chart_type", None)
                if ct == chart_type:
                    return elem
        return None

    @staticmethod
    def _find_heading(elements: list) -> str:
        """Extract heading text from elements."""
        for elem in elements:
            elem_type = getattr(elem, "type", None)
            if hasattr(elem_type, "value"):
                elem_type = elem_type.value
            if elem_type == "heading":
                return getattr(getattr(elem, "content", None), "text", "")
        return ""

    @staticmethod
    def _find_text_elements(elements: list) -> list:
        """Find all body text and bullet list elements."""
        results = []
        for elem in elements:
            elem_type = getattr(elem, "type", None)
            if hasattr(elem_type, "value"):
                elem_type = elem_type.value
            if elem_type in ("body_text", "bullet_list", "numbered_list"):
                results.append(elem)
        return results


__all__ = ["BaseFinanceSlideRenderer"]
