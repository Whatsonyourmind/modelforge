"""Market landscape slide renderer -- TAM/SAM/SOM shapes and market data table."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from modelforge.deck.finance.formatter import FinancialFormatter
from modelforge.deck.ir.elements.base import Position
from modelforge.deck.rendering.slide_renderers.base import BaseFinanceSlideRenderer, _infer_format
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.slides.base import BaseSlide
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)

# TAM/SAM/SOM keywords for detection.
_TAM_KEYWORDS = frozenset({"tam", "total addressable market"})
_SAM_KEYWORDS = frozenset({"sam", "serviceable addressable market", "serviceable available market"})
_SOM_KEYWORDS = frozenset({"som", "serviceable obtainable market"})


class MarketLandscapeRenderer(BaseFinanceSlideRenderer):
    """Renders market landscape with TAM/SAM/SOM nested shapes and data tables."""

    def render(self, slide: Slide, ir_slide: BaseSlide, theme: ResolvedTheme) -> None:
        elements = ir_slide.elements

        # Title
        title = self._find_heading(elements) or "Market Landscape"
        title_pos = Position(x=0.75, y=0.4, width=11.8, height=0.8)
        self._add_title(slide, title, theme, title_pos)

        # Find table elements
        tables = self._find_table_elements(elements)

        # Classify: TAM/SAM/SOM table vs competitors/market data table
        tam_table = None
        data_table = None

        for tbl in tables:
            content = tbl.content
            if self._is_tam_table(content):
                tam_table = tbl
            elif data_table is None:
                data_table = tbl

        # Render TAM/SAM/SOM as nested rounded rectangles (left half)
        if tam_table:
            self._render_tam_shapes(slide, tam_table, theme)

        # Render market data table (right half)
        if data_table:
            content = data_table.content
            col_formats = [_infer_format(h) for h in content.headers]
            d_pos = Position(x=7.0, y=1.5, width=5.5, height=3.5)
            self._add_table(
                slide,
                headers=content.headers,
                rows=[list(r) for r in content.rows],
                theme=theme,
                position=d_pos,
                column_formats=col_formats,
            )

    def _render_tam_shapes(
        self, slide: Slide, tam_table, theme: ResolvedTheme
    ) -> None:
        """Render TAM/SAM/SOM as nested rounded rectangles."""
        content = tam_table.content
        rows = content.rows

        # Build TAM/SAM/SOM data from rows
        segments = []
        for row in rows:
            label = str(row[0]) if row else ""
            value = row[1] if len(row) > 1 else 0
            segments.append((label, value))

        # Sizes (largest to smallest) -- three nested rectangles
        sizes = [
            (0.75, 1.5, 5.5, 4.5),   # Outer (TAM)
            (1.25, 2.0, 4.5, 3.5),   # Middle (SAM)
            (1.75, 2.5, 3.5, 2.5),   # Inner (SOM)
        ]

        # Color shades -- lighter for TAM (outer), darker for SOM (inner)
        from modelforge.deck.finance.conditional import _lighten
        colors = [
            _lighten(theme.colors.primary, 0.75),  # Lightest for TAM
            _lighten(theme.colors.primary, 0.50),   # Medium for SAM
            _lighten(theme.colors.primary, 0.25),   # Darkest for SOM
        ]

        font_name = resolve_font_name(theme.typography.body_family)

        for idx, (label, value) in enumerate(segments[:3]):
            if idx >= len(sizes):
                break

            x, y, w, h = sizes[idx]
            left = Inches(x)
            top = Inches(y)
            width = Inches(w)
            height = Inches(h)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(colors[idx])

            # Add label text
            tf = shape.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER

            # Format value as currency if numeric
            if isinstance(value, (int, float)):
                display = f"{label}\n{FinancialFormatter.currency(value)}"
            else:
                display = f"{label}\n{value}"

            para.text = display
            for run in para.runs:
                run.font.name = font_name
                run.font.size = Pt(14 - idx * 2)  # Slightly smaller for inner
                run.font.bold = True
                run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)

    @staticmethod
    def _is_tam_table(content) -> bool:
        """Check if a table contains TAM/SAM/SOM data."""
        for row in content.rows:
            if row:
                label = str(row[0]).lower()
                if any(kw in label for kw in ("tam", "sam", "som", "addressable")):
                    return True
        return False


__all__ = ["MarketLandscapeRenderer"]
