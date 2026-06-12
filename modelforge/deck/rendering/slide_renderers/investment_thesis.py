"""Investment thesis slide renderer -- numbered thesis points with risk/reward."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from modelforge.deck.ir.elements.base import Position
from modelforge.deck.rendering.slide_renderers.base import BaseFinanceSlideRenderer, _infer_format
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.slides.base import BaseSlide
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)


class InvestmentThesisRenderer(BaseFinanceSlideRenderer):
    """Renders investment thesis with numbered points and optional risk table."""

    def render(self, slide: Slide, ir_slide: BaseSlide, theme: ResolvedTheme) -> None:
        elements = ir_slide.elements

        # Title
        title = self._find_heading(elements) or "Investment Thesis"
        title_pos = Position(x=0.75, y=0.4, width=11.8, height=0.8)
        self._add_title(slide, title, theme, title_pos)

        # Collect thesis points from body_text, bullet_list elements
        text_elements = self._find_text_elements(elements)

        # Render numbered thesis points
        font_name = resolve_font_name(theme.typography.body_family)
        heading_font = resolve_font_name(theme.typography.heading_family)
        y_offset = 1.5
        element_gap = theme.spacing.element_gap

        for idx, elem in enumerate(text_elements):
            number = idx + 1
            elem_type = getattr(elem, "type", None)
            if hasattr(elem_type, "value"):
                elem_type = elem_type.value

            # Extract text content
            if elem_type == "body_text":
                text = getattr(getattr(elem, "content", None), "text", "")
            elif elem_type == "bullet_list":
                items = getattr(getattr(elem, "content", None), "items", [])
                text = "; ".join(items)
            elif elem_type == "numbered_list":
                items = getattr(getattr(elem, "content", None), "items", [])
                text = "; ".join(items)
            else:
                text = str(elem)

            # Create text box with numbered point
            left = Inches(0.75)
            top = Inches(y_offset)
            width = Inches(11.8)
            height = Inches(0.7)

            txbox = slide.shapes.add_textbox(left, top, width, height)
            tf = txbox.text_frame
            tf.word_wrap = True

            # Number + text
            para = tf.paragraphs[0]
            para.text = f"{number}. {text}"
            para.alignment = PP_ALIGN.LEFT

            for run in para.runs:
                run.font.name = font_name
                run.font.size = Pt(theme.typography.scale.get("body", 18))
                run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)

            y_offset += 0.7 + element_gap

        # Optional: risk/reward table at bottom
        table_elem = self._find_table_element(elements)
        if table_elem:
            content = table_elem.content
            col_formats = [_infer_format(h) for h in content.headers]
            table_pos = Position(x=0.75, y=max(y_offset, 5.0), width=11.8, height=2.0)
            self._add_table(
                slide,
                headers=content.headers,
                rows=[list(r) for r in content.rows],
                theme=theme,
                position=table_pos,
                column_formats=col_formats,
            )


__all__ = ["InvestmentThesisRenderer"]
