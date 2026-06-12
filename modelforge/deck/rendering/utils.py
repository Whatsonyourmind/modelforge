"""Rendering utility functions -- color conversion, font mapping, transitions, backgrounds."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap, qn

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PptxPresentation
    from pptx.slide import Slide

    from modelforge.deck.ir.enums import Transition

logger = logging.getLogger(__name__)

# Safe PowerPoint fonts that are universally available
_SAFE_FONTS = frozenset({
    "Arial",
    "Calibri",
    "Cambria",
    "Candara",
    "Century Gothic",
    "Comic Sans MS",
    "Consolas",
    "Constantia",
    "Corbel",
    "Courier New",
    "Franklin Gothic Medium",
    "Garamond",
    "Georgia",
    "Impact",
    "Liberation Sans",
    "Liberation Serif",
    "Liberation Mono",
    "Lucida Console",
    "Lucida Sans Unicode",
    "Open Sans",
    "Palatino Linotype",
    "Segoe UI",
    "Tahoma",
    "Times New Roman",
    "Trebuchet MS",
    "Verdana",
})

# Mapping of popular Google/non-system fonts to their closest safe equivalents.
# Themes reference these fonts, and this mapping ensures graceful degradation
# when the exact font is not installed on the rendering host.
_FONT_SUBSTITUTES: dict[str, str] = {
    # Sans-serif: geometric/modern -> closest system match
    "Montserrat": "Open Sans",
    "Open Sans": "Open Sans",
    "Inter": "Calibri",
    "Roboto": "Arial",
    "Lato": "Calibri",
    "Poppins": "Calibri",
    "Nunito": "Calibri",
    "Work Sans": "Calibri",
    "DM Sans": "Calibri",
    "Source Sans Pro": "Calibri",
    "Source Sans 3": "Calibri",
    "Raleway": "Century Gothic",
    "Space Grotesk": "Calibri",
    "IBM Plex Sans": "Segoe UI",
    # Sans-serif: condensed/display
    "Oswald": "Franklin Gothic Medium",
    # Serif fonts
    "Merriweather": "Georgia",
    "Playfair Display": "Georgia",
    "Libre Baskerville": "Palatino Linotype",
    # Monospace fonts
    "JetBrains Mono": "Consolas",
    "Roboto Mono": "Consolas",
    "Fira Code": "Consolas",
    "Fira Mono": "Consolas",
    "IBM Plex Mono": "Consolas",
    "Source Code Pro": "Consolas",
    "Ubuntu Mono": "Consolas",
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string to a python-pptx RGBColor.

    Args:
        hex_color: Hex color string, with or without '#' prefix.

    Returns:
        RGBColor instance.
    """
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def resolve_font_name(requested: str, fallback: str = "Calibri") -> str:
    """Map a requested font name to a safe PowerPoint font.

    First checks if the font is already safe (universally available in
    PowerPoint). If not, looks up the best substitute from the
    ``_FONT_SUBSTITUTES`` mapping. Falls back to ``fallback`` if no
    known substitute exists.

    Args:
        requested: The font name to look up.
        fallback: The fallback font if no substitute is known.

    Returns:
        A safe font name.
    """
    if requested in _SAFE_FONTS:
        return requested
    substitute = _FONT_SUBSTITUTES.get(requested)
    if substitute:
        logger.debug("Font substitute: %s -> %s", requested, substitute)
        return substitute
    logger.debug("Unknown font '%s', using fallback '%s'", requested, fallback)
    return fallback


def set_slide_background(slide: Slide, hex_color: str) -> None:
    """Set a solid fill background on a slide.

    Args:
        slide: python-pptx Slide object.
        hex_color: Hex color string for the background.
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def set_transition(slide: Slide, transition: Transition) -> None:
    """Apply a slide transition via Open XML manipulation.

    Args:
        slide: python-pptx Slide object.
        transition: IR Transition enum value (NONE, FADE, SLIDE, PUSH).
    """
    from modelforge.deck.ir.enums import Transition as TransitionEnum

    slide_elem = slide._element

    # Remove any existing transition element
    existing = slide_elem.findall(qn("p:transition"))
    for elem in existing:
        slide_elem.remove(elem)

    if transition == TransitionEnum.NONE:
        return

    # Build the transition element
    # The p:transition element is a child of p:cSld's parent (the slide element)
    nsmap_p = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

    trans_elem = etree.SubElement(slide_elem, qn("p:transition"))
    trans_elem.set("spd", "med")

    if transition == TransitionEnum.FADE:
        etree.SubElement(trans_elem, qn("p:fade"))
    elif transition == TransitionEnum.SLIDE:
        push = etree.SubElement(trans_elem, qn("p:push"))
        push.set("dir", "l")
    elif transition == TransitionEnum.PUSH:
        push = etree.SubElement(trans_elem, qn("p:push"))
        push.set("dir", "l")


def get_blank_layout(prs: PptxPresentation):
    """Find the layout with the fewest placeholders (ideally blank).

    Args:
        prs: python-pptx Presentation object.

    Returns:
        A slide layout object.
    """
    best_layout = None
    min_placeholders = float("inf")

    for layout in prs.slide_layouts:
        ph_count = len(layout.placeholders)
        if ph_count < min_placeholders:
            min_placeholders = ph_count
            best_layout = layout
        if ph_count == 0:
            break

    return best_layout
