"""PptxRenderer -- orchestrates rendering of a complete PPTX from IR + layout + theme."""

from __future__ import annotations

import io
import logging
from typing import TYPE_CHECKING

from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from modelforge.deck.rendering.element_renderers import render_element
from modelforge.deck.rendering.slide_renderers import render_finance_slide
from modelforge.deck.rendering.utils import (
    get_blank_layout,
    set_slide_background,
    set_transition,
)

if TYPE_CHECKING:
    from modelforge.deck.ir.presentation import Presentation
    from modelforge.deck.layout.types import LayoutResult
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)


class DeckRenderError(RuntimeError):
    """Raised when slide content would be silently lost during rendering.

    Either an element reached the renderer without a layout position
    (composer/pattern contract breach) or an element renderer failed.
    Silent content loss is never acceptable in a certified deck.
    """


# Element types that legitimately carry no renderable content of their own
# (handled at slide level or structural only) — exempt from the
# missing-position failure.
_POSITION_EXEMPT_TYPES = frozenset(
    {"background", "icon", "container", "column", "row", "grid_cell"}
)


def _pattern_name_for(slide_type: str) -> str:
    """Best-effort name of the layout pattern that owns a slide type."""
    try:
        from modelforge.deck.layout.patterns import PATTERN_REGISTRY

        pattern_cls = PATTERN_REGISTRY.get(slide_type)
        return pattern_cls.__name__ if pattern_cls is not None else "<unknown>"
    except Exception:  # pragma: no cover - defensive
        return "<unknown>"


class PptxRenderer:
    """Renders a complete PPTX file from IR presentation + layout results + theme.

    The renderer creates one PPTX slide per LayoutResult, applies backgrounds,
    transitions, and speaker notes, then dispatches element rendering to the
    ELEMENT_RENDERERS registry.

    Content-loss contract: any element that reaches the renderer without a
    position, or whose element renderer raises, is collected and reported via
    DeckRenderError at the end of the render — content is never silently
    dropped.
    """

    def render(
        self,
        presentation: Presentation,
        layout_results: list[LayoutResult],
        theme: ResolvedTheme,
    ) -> bytes:
        """Render a complete PPTX from IR + layout results + theme.

        Args:
            presentation: The IR Presentation model.
            layout_results: List of LayoutResults from the layout engine.
            theme: Resolved theme for styling.

        Returns:
            Raw bytes of the generated .pptx file.

        Raises:
            DeckRenderError: If any element would be silently dropped
                (no position assigned by the layout pattern) or fails to
                render.
        """
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = get_blank_layout(prs)

        content_loss: list[str] = []

        for slide_no, layout_result in enumerate(layout_results, start=1):
            slide = prs.slides.add_slide(blank_layout)
            ir_slide = layout_result.slide

            # 1. Set background color
            self._apply_background(slide, ir_slide, theme)

            # 2. Render elements
            self._render_elements(slide, ir_slide, theme, slide_no, content_loss)

            # 3. Apply transition
            self._apply_transition(slide, ir_slide)

            # 4. Apply speaker notes
            self._apply_speaker_notes(slide, ir_slide)

        if content_loss:
            raise DeckRenderError(
                "Deck rendering would silently lose content on "
                f"{len(content_loss)} element(s): " + "; ".join(content_loss)
            )

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    def _apply_background(self, slide, ir_slide, theme: ResolvedTheme) -> None:
        """Set slide background from theme slide master or default theme background."""
        slide_type = ir_slide.slide_type
        if hasattr(slide_type, "value"):
            slide_type = slide_type.value

        # Look up slide master for this slide type
        master = theme.slide_masters.get(slide_type)
        if master:
            bg_color = master.background
        else:
            bg_color = theme.colors.background

        set_slide_background(slide, bg_color)

    def _render_elements(
        self,
        slide,
        ir_slide,
        theme: ResolvedTheme,
        slide_no: int,
        content_loss: list[str],
    ) -> None:
        """Render all elements on the slide.

        Finance slide types are dispatched to FINANCE_SLIDE_RENDERERS for
        full-slide rendering. Non-finance slides use the element-by-element path.

        Any element that cannot be rendered (no position, or renderer raised)
        is recorded in ``content_loss`` so the caller can fail loudly.
        """
        slide_type = ir_slide.slide_type
        if hasattr(slide_type, "value"):
            slide_type = slide_type.value

        # Check if this is a finance slide type -- if so, the finance renderer
        # handles the entire slide (title, tables, charts, positioning).
        if render_finance_slide(slide, ir_slide, theme):
            return

        for element in ir_slide.elements:
            element_type = element.type
            if hasattr(element_type, "value"):
                element_type = element_type.value

            # Slide-level / structural elements carry no own content
            if element_type in _POSITION_EXEMPT_TYPES:
                continue

            position = element.position
            if position is None:
                content_loss.append(
                    f"element {element_type} on slide {slide_no} has no "
                    f"position assigned by pattern "
                    f"{_pattern_name_for(slide_type)}"
                )
                continue

            try:
                render_element(slide, element, position, theme)
            except Exception as e:
                logger.exception(
                    "Failed to render element type=%s",
                    getattr(element, "type", "unknown"),
                )
                content_loss.append(
                    f"element {element_type} on slide {slide_no} failed to "
                    f"render: {e!r}"
                )

    def _apply_transition(self, slide, ir_slide) -> None:
        """Apply slide transition if specified in the IR."""
        from modelforge.deck.ir.enums import Transition

        transition = ir_slide.transition
        if transition is not None and transition != Transition.NONE:
            set_transition(slide, transition)

    def _apply_speaker_notes(self, slide, ir_slide) -> None:
        """Add speaker notes to the slide if present in the IR."""
        notes = ir_slide.speaker_notes
        if notes:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes


__all__ = ["DeckRenderError", "PptxRenderer"]
