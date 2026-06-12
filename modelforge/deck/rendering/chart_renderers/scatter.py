"""Scatter and bubble chart renderers — XY and bubble native charts."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.chart.data import BubbleChartData as PptxBubbleChartData
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from modelforge.deck.rendering.chart_renderers.base import BaseChartRenderer

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import Position
    from modelforge.deck.themes.types import ResolvedTheme


class ScatterChartRenderer(BaseChartRenderer):
    """Renders native editable XY_SCATTER charts.

    Uses XyChartData to pair x_values with each series' y-values
    as individual data points. The chart is fully editable in PowerPoint.
    """

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        cd = XyChartData()

        x_values = chart_data.x_values  # type: ignore[attr-defined]
        for series in chart_data.series:  # type: ignore[attr-defined]
            xy_series = cd.add_series(series.name)
            for i, x_val in enumerate(x_values):
                y_val = series.values[i] if i < len(series.values) else 0
                if y_val is None:
                    y_val = 0
                xy_series.add_data_point(x_val, y_val)

        x = Inches(position.x or 0)
        y = Inches(position.y or 0)
        w = Inches(position.width or 10)
        h = Inches(position.height or 5)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, x, y, w, h, cd
        )
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._apply_theme_formatting(chart, theme, title=title)


class BubbleChartRenderer(BaseChartRenderer):
    """Renders native editable BUBBLE charts.

    Uses BubbleChartData with (x, y, size) tuples for each data point.
    The chart is fully editable in PowerPoint.
    """

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        cd = PptxBubbleChartData()

        x_values = chart_data.x_values  # type: ignore[attr-defined]
        sizes = chart_data.sizes  # type: ignore[attr-defined]

        for series in chart_data.series:  # type: ignore[attr-defined]
            bubble_series = cd.add_series(series.name)
            for i, x_val in enumerate(x_values):
                y_val = series.values[i] if i < len(series.values) else 0
                if y_val is None:
                    y_val = 0
                size_val = sizes[i] if i < len(sizes) else 1
                bubble_series.add_data_point(x_val, y_val, size_val)

        x = Inches(position.x or 0)
        y = Inches(position.y or 0)
        w = Inches(position.width or 10)
        h = Inches(position.height or 5)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BUBBLE, x, y, w, h, cd
        )
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._apply_theme_formatting(chart, theme, title=title)
