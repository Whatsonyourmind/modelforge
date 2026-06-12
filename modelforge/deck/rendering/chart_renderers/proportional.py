"""Proportional chart renderers — pie and donut charts."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from modelforge.deck.rendering.chart_renderers.base import BaseChartRenderer

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import Position
    from modelforge.deck.themes.types import ResolvedTheme


class PieChartRenderer(BaseChartRenderer):
    """Renders native editable PIE charts."""

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        cd = CategoryChartData()
        cd.categories = chart_data.labels  # type: ignore[attr-defined]
        cd.add_series("Values", chart_data.values)  # type: ignore[attr-defined]

        x = Inches(position.x or 0)
        y = Inches(position.y or 0)
        w = Inches(position.width or 10)
        h = Inches(position.height or 5)

        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, w, h, cd)
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._apply_theme_formatting(chart, theme, title=title)

        # Color each slice using theme.chart_colors
        colors = theme.chart_colors
        if colors:
            plot = chart.plots[0]
            for idx, point in enumerate(plot.series[0].points):
                color_hex = colors[idx % len(colors)].lstrip("#")
                try:
                    fill = point.format.fill
                    fill.solid()
                    from pptx.dml.color import RGBColor

                    fill.fore_color.rgb = RGBColor.from_string(color_hex)
                except Exception:
                    pass

        # Percentage data labels
        show_percentages = getattr(chart_data, "show_percentages", True)
        if show_percentages:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.show_category_name = False


class DonutChartRenderer(BaseChartRenderer):
    """Renders native editable DOUGHNUT charts."""

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        cd = CategoryChartData()
        cd.categories = chart_data.labels  # type: ignore[attr-defined]
        cd.add_series("Values", chart_data.values)  # type: ignore[attr-defined]

        x = Inches(position.x or 0)
        y = Inches(position.y or 0)
        w = Inches(position.width or 10)
        h = Inches(position.height or 5)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd
        )
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._apply_theme_formatting(chart, theme, title=title)

        # Color each slice using theme.chart_colors
        colors = theme.chart_colors
        if colors:
            plot = chart.plots[0]
            for idx, point in enumerate(plot.series[0].points):
                color_hex = colors[idx % len(colors)].lstrip("#")
                try:
                    fill = point.format.fill
                    fill.solid()
                    from pptx.dml.color import RGBColor

                    fill.fore_color.rgb = RGBColor.from_string(color_hex)
                except Exception:
                    pass

        # Set hole size from inner_radius via XML
        inner_radius = getattr(chart_data, "inner_radius", 0.5)
        hole_size = int(inner_radius * 100)
        try:
            from lxml import etree

            chart_space = chart._chartSpace
            ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
            # Find the doughnut plot element
            doughnut_plots = chart_space.findall(f".//{{{ns}}}doughnutChart")
            for dp in doughnut_plots:
                hole_elem = dp.find(f"{{{ns}}}holeSize")
                if hole_elem is None:
                    hole_elem = etree.SubElement(dp, f"{{{ns}}}holeSize")
                hole_elem.set("val", str(hole_size))
        except Exception:
            pass

        # Percentage data labels
        show_percentages = getattr(chart_data, "show_percentages", True)
        if show_percentages:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.show_category_name = False
