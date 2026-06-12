"""Category-based chart renderers — bar, line, area (and their variants)."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from modelforge.deck.rendering.chart_renderers.base import BaseChartRenderer

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import Position
    from modelforge.deck.themes.types import ResolvedTheme


# ── Mapping from chart_type to XL_CHART_TYPE ─────────────────────────────────

_BAR_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "grouped_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
}

_LINE_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "multi_line": XL_CHART_TYPE.LINE_MARKERS,
}

_AREA_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "area": XL_CHART_TYPE.AREA,
    "stacked_area": XL_CHART_TYPE.AREA_STACKED,
}


# ── Helpers ───────────────────────────────────────────────────────────────────


def _build_category_chart_data(chart_data: object) -> CategoryChartData:
    """Build CategoryChartData from any chart model with categories + series."""
    cd = CategoryChartData()
    cd.categories = chart_data.categories  # type: ignore[attr-defined]
    for series in chart_data.series:  # type: ignore[attr-defined]
        # Replace None values with 0 for pptx compatibility
        values = [v if v is not None else 0 for v in series.values]
        cd.add_series(series.name, values)
    return cd


# ── Bar Chart Renderer ────────────────────────────────────────────────────────


class BarChartRenderer(BaseChartRenderer):
    """Renders bar, stacked_bar, grouped_bar, and horizontal_bar charts."""

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        chart_type_str = getattr(chart_data, "chart_type", "bar")
        xl_type = _BAR_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        cd = _build_category_chart_data(chart_data)

        x = Inches(position.x or 0)
        y = Inches(position.y or 0)
        w = Inches(position.width or 10)
        h = Inches(position.height or 5)

        graphic_frame = slide.shapes.add_chart(xl_type, x, y, w, h, cd)
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._apply_theme_formatting(chart, theme, title=title)

        # Data labels
        show_values = getattr(chart_data, "show_values", False)
        if show_values:
            for plot in chart.plots:
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.show_value = True


# ── Line Chart Renderer ──────────────────────────────────────────────────────


class LineChartRenderer(BaseChartRenderer):
    """Renders line and multi_line charts."""

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        chart_type_str = getattr(chart_data, "chart_type", "line")
        show_markers = getattr(chart_data, "show_markers", True)

        if show_markers:
            xl_type = _LINE_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.LINE_MARKERS)
        else:
            xl_type = XL_CHART_TYPE.LINE

        cd = _build_category_chart_data(chart_data)

        x = Inches(position.x or 0)
        y = Inches(position.y or 0)
        w = Inches(position.width or 10)
        h = Inches(position.height or 5)

        graphic_frame = slide.shapes.add_chart(xl_type, x, y, w, h, cd)
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._apply_theme_formatting(chart, theme, title=title)


# ── Area Chart Renderer ──────────────────────────────────────────────────────


class AreaChartRenderer(BaseChartRenderer):
    """Renders area and stacked_area charts."""

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        chart_type_str = getattr(chart_data, "chart_type", "area")
        xl_type = _AREA_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.AREA)

        cd = _build_category_chart_data(chart_data)

        x = Inches(position.x or 0)
        y = Inches(position.y or 0)
        w = Inches(position.width or 10)
        h = Inches(position.height or 5)

        graphic_frame = slide.shapes.add_chart(xl_type, x, y, w, h, cd)
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._apply_theme_formatting(chart, theme, title=title)
