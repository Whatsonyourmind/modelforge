"""NATIVE python-pptx rewrites of the signature IB-style chart renderers.

The original DeckForge implementations of these chart types were static
Plotly/kaleido PNG exporters (``static_base.py`` + one module per chart
type). plotly/kaleido are intentionally NOT modelforge dependencies, so those
modules were EXCLUDED from the E2 extraction. This module provides native,
fully editable python-pptx replacements -- no Plotly, no kaleido, no PNG, no
randomness, no timestamps: rendering the same IR twice yields byte-identical
chart XML parts.

Implemented here (registered in ``CHART_RENDERERS``, see ``__init__.py``):

    tornado          -- horizontal stacked bars centered on a base value
                        (invisible-base + downside/upside spread series; the
                        same noFill trick ``builder/sheets/fairness_football.py``
                        uses for the Excel football field).
    football_field   -- floating low->high horizontal valuation-range bars
                        (invisible low base + visible range series).
    waterfall        -- vertical floating bricks (stacked column, invisible
                        running base; positive / negative / total colored;
                        connector-free v1).
    sensitivity      -- NOT a chart: a native PPTX table with theme-driven
                        conditional shading from
                        ``modelforge.deck.finance.conditional``.

Still pending (registry maps them to ``PendingNativeChartRenderer``):
    funnel, treemap, heatmap, sankey, gantt, sunburst
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from modelforge.deck.finance.conditional import ConditionalFormatter
from modelforge.deck.rendering.chart_renderers.base import BaseChartRenderer
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.chart.chart import Chart
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import Position
    from modelforge.deck.themes.types import ResolvedTheme

__all__ = [
    "FootballFieldChartRenderer",
    "SensitivityTableRenderer",
    "TornadoChartRenderer",
    "WaterfallChartRenderer",
]


# Name of the invisible offset series used by the floating-bar trick. The
# leading underscore signals "internal plumbing, not data" to anyone editing
# the chart data in PowerPoint.
_BASE_SERIES_NAME = "_base"

# Deterministic styling constants (do NOT derive from anything environmental).
_GAP_WIDTH = 60  # percent -- bar/column gap for all IB charts in this module


# ── Shared helpers ────────────────────────────────────────────────────────────


def _position_box(position: Position) -> tuple:
    """Resolve a Position into (x, y, w, h) EMU values with category.py defaults."""
    return (
        Inches(position.x or 0),
        Inches(position.y or 0),
        Inches(position.width or 10),
        Inches(position.height or 5),
    )


def _fmt_num(value: float | int) -> str:
    """Deterministic, human-friendly number formatting for table cells."""
    f = float(value)
    if f.is_integer():
        return f"{int(f):,}"
    return f"{f:,.2f}"


class _FloatingBarRendererMixin(BaseChartRenderer):
    """Shared finishing pass for the invisible-base stacked-bar renderers."""

    def _finish_floating_chart(
        self,
        chart: Chart,
        theme: ResolvedTheme,
        title: str | None,
        visible_series_colors: list[str],
    ) -> None:
        """Apply theme fonts/title, hide the base series, color the visible ones.

        ``visible_series_colors`` maps (in order) onto ``chart.series[1:]``.
        The legend is disabled: the invisible base series would otherwise
        appear as an empty legend swatch.
        """
        self._apply_theme_formatting(chart, theme, title=title)

        # The invisible-base trick makes the legend misleading -> IB style is
        # axis labels only.
        chart.has_legend = False

        series_list = list(chart.series)

        # Series 0: the invisible base (same noFill trick as the Excel
        # football field in builder/sheets/fairness_football.py).
        base_series = series_list[0]
        base_series.format.fill.background()
        base_series.format.line.fill.background()

        # Visible series: explicit theme-driven colors.
        for series, color_hex in zip(series_list[1:], visible_series_colors):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(color_hex)
            series.format.line.fill.background()

        # Deterministic bar geometry.
        plot = chart.plots[0]
        plot.gap_width = _GAP_WIDTH


# ── Tornado ───────────────────────────────────────────────────────────────────


class TornadoChartRenderer(_FloatingBarRendererMixin):
    """Tornado / sensitivity-driver chart -- horizontal bars centered on base.

    Built as BAR_STACKED with three series:
        ``_base``    invisible offset from 0 to the left edge of each bar
        ``Downside`` the low->base spread (theme negative color)
        ``Upside``   the base->high spread (theme positive color)

    Categories are sorted by spread (widest at the TOP -- python-pptx BAR
    charts plot the first category at the bottom, so the sort is ascending).
    The sort is stable, so equal spreads keep their input order
    (deterministic). If ``base_value`` is missing, the midpoint of the overall
    low/high envelope is used.
    """

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        categories = list(chart_data.categories)  # type: ignore[attr-defined]
        lows = [float(v) for v in chart_data.low_values]  # type: ignore[attr-defined]
        highs = [float(v) for v in chart_data.high_values]  # type: ignore[attr-defined]

        base_value = getattr(chart_data, "base_value", None)
        if base_value is None:
            envelope = lows + highs
            base = (min(envelope) + max(envelope)) / 2.0
        else:
            base = float(base_value)

        # Normalize each row so lo <= hi, then sort ascending by spread
        # (stable) so the widest driver lands at the top of the chart.
        rows = []
        for cat, lo, hi in zip(categories, lows, highs):
            if lo > hi:
                lo, hi = hi, lo
            rows.append((cat, lo, hi))
        rows.sort(key=lambda r: r[2] - r[1])

        invisible = [min(lo, base) for _, lo, _ in rows]
        downside = [max(min(hi, base) - lo, 0.0) for _, lo, hi in rows]
        upside = [max(hi - max(lo, base), 0.0) for _, lo, hi in rows]

        cd = CategoryChartData()
        cd.categories = [cat for cat, _, _ in rows]
        cd.add_series(_BASE_SERIES_NAME, invisible)
        cd.add_series("Downside", downside)
        cd.add_series("Upside", upside)

        x, y, w, h = _position_box(position)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, x, y, w, h, cd)
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._finish_floating_chart(
            chart,
            theme,
            title,
            [theme.colors.negative, theme.colors.positive],
        )


# ── Football field ────────────────────────────────────────────────────────────


class FootballFieldChartRenderer(_FloatingBarRendererMixin):
    """Football field -- floating low->high valuation-range bars per method.

    Built as BAR_STACKED with two series:
        ``_base``  invisible offset from 0 to each method's low value
        ``Range``  the visible low->high spread (theme primary color)

    Category order is preserved from the IR (methodology order is the
    author's narrative choice). ``mid_values`` are accepted but not drawn in
    v1 (no midpoint tick marks -- documented limitation).
    """

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        categories = list(chart_data.categories)  # type: ignore[attr-defined]
        lows = [float(v) for v in chart_data.low_values]  # type: ignore[attr-defined]
        highs = [float(v) for v in chart_data.high_values]  # type: ignore[attr-defined]

        invisible: list[float] = []
        spreads: list[float] = []
        for lo, hi in zip(lows, highs):
            if lo > hi:
                lo, hi = hi, lo
            invisible.append(lo)
            spreads.append(hi - lo)

        cd = CategoryChartData()
        cd.categories = categories
        cd.add_series(_BASE_SERIES_NAME, invisible)
        cd.add_series("Range", spreads)

        x, y, w, h = _position_box(position)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, x, y, w, h, cd)
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._finish_floating_chart(chart, theme, title, [theme.colors.primary])


# ── Waterfall ─────────────────────────────────────────────────────────────────


class WaterfallChartRenderer(_FloatingBarRendererMixin):
    """Waterfall bridge -- vertical floating bricks (connector-free v1).

    Built as COLUMN_STACKED with four series:
        ``_base``     invisible offset from 0 to the bottom of each brick
        ``Increase``  positive flows (theme positive color)
        ``Decrease``  negative flows (theme negative color)
        ``Total``     total/subtotal bricks anchored at 0 (theme primary color)

    Total detection (deterministic): the first and last categories are
    treated as totals, and any middle value exactly equal to the running
    cumulative is treated as a subtotal (e.g. "Gross Profit" after
    Revenue - COGS). All other values are flows added to the running total.
    """

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        categories = list(chart_data.categories)  # type: ignore[attr-defined]
        values = [float(v) for v in chart_data.values]  # type: ignore[attr-defined]

        invisible: list[float] = []
        increases: list[float] = []
        decreases: list[float] = []
        totals: list[float] = []

        cumulative = 0.0
        last_idx = len(values) - 1
        for idx, value in enumerate(values):
            is_total = idx == 0 or idx == last_idx or value == cumulative
            if is_total:
                # Brick anchored at 0 (covers [min(0, v), max(0, v)]).
                invisible.append(min(0.0, value))
                increases.append(0.0)
                decreases.append(0.0)
                totals.append(abs(value))
                cumulative = value
            else:
                new_cumulative = cumulative + value
                invisible.append(min(cumulative, new_cumulative))
                increases.append(value if value > 0 else 0.0)
                decreases.append(-value if value < 0 else 0.0)
                totals.append(0.0)
                cumulative = new_cumulative

        cd = CategoryChartData()
        cd.categories = categories
        cd.add_series(_BASE_SERIES_NAME, invisible)
        cd.add_series("Increase", increases)
        cd.add_series("Decrease", decreases)
        cd.add_series("Total", totals)

        x, y, w, h = _position_box(position)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, w, h, cd)
        chart = graphic_frame.chart

        title = getattr(chart_data, "title", None)
        self._finish_floating_chart(
            chart,
            theme,
            title,
            [theme.colors.positive, theme.colors.negative, theme.colors.primary],
        )


# ── Sensitivity table ─────────────────────────────────────────────────────────


class SensitivityTableRenderer(BaseChartRenderer):
    """Two-variable sensitivity matrix -- a NATIVE PPTX table, not a chart.

    Layout: (1 + len(row_values)) rows x (1 + len(col_values)) columns.
    The corner cell shows ``"{row_header} \\ {col_header}"``; the header row
    carries ``col_values`` and the first column carries ``row_values``. Data
    cells get theme-driven conditional shading via
    ``ConditionalFormatter.heatmap_gradient`` (negative theme color at the
    matrix minimum -> positive theme color at the maximum). An optional title
    is rendered as a text box above the table.
    """

    _TITLE_HEIGHT_IN = 0.5  # inches reserved for the optional title text box

    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        row_header = chart_data.row_header  # type: ignore[attr-defined]
        col_header = chart_data.col_header  # type: ignore[attr-defined]
        row_values = list(chart_data.row_values)  # type: ignore[attr-defined]
        col_values = list(chart_data.col_values)  # type: ignore[attr-defined]
        data = [list(row) for row in chart_data.data]  # type: ignore[attr-defined]
        title = getattr(chart_data, "title", None)

        x_in = position.x or 0
        y_in = position.y or 0
        w_in = position.width or 10
        h_in = position.height or 5

        typo = theme.typography
        body_font = resolve_font_name(typo.body_family)
        caption_size = Pt(typo.scale.get("caption", 14))

        # Optional title text box above the table.
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(x_in),
                Inches(y_in),
                Inches(w_in),
                Inches(self._TITLE_HEIGHT_IN),
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = title
            run.font.name = resolve_font_name(typo.heading_family)
            run.font.size = Pt(typo.scale.get("subtitle", 24))
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)
            y_in += self._TITLE_HEIGHT_IN
            h_in = max(h_in - self._TITLE_HEIGHT_IN, 0.5)

        n_rows = 1 + len(row_values)
        n_cols = 1 + len(col_values)
        table_shape = slide.shapes.add_table(
            n_rows,
            n_cols,
            Inches(x_in),
            Inches(y_in),
            Inches(w_in),
            Inches(h_in),
        )
        table = table_shape.table

        header_fill = hex_to_rgb(theme.colors.primary)
        header_text = hex_to_rgb("#FFFFFF")
        body_text = hex_to_rgb(theme.colors.text_primary)

        # Flattened matrix bounds for the conditional gradient (deterministic).
        flat = [float(v) for row in data for v in row]
        min_val = min(flat)
        max_val = max(flat)

        # Corner cell.
        corner = table.cell(0, 0)
        corner.text = f"{row_header} \\ {col_header}"
        self._style_cell(corner, caption_size, body_font, bold=True,
                         text_color=header_text, fill_color=header_fill,
                         alignment=PP_ALIGN.CENTER)

        # Header row (column variable values).
        for c_idx, col_val in enumerate(col_values):
            cell = table.cell(0, c_idx + 1)
            cell.text = _fmt_num(col_val)
            self._style_cell(cell, caption_size, body_font, bold=True,
                             text_color=header_text, fill_color=header_fill,
                             alignment=PP_ALIGN.CENTER)

        # Body rows: row-variable label + conditionally shaded data cells.
        for r_idx, row_val in enumerate(row_values):
            label_cell = table.cell(r_idx + 1, 0)
            label_cell.text = _fmt_num(row_val)
            self._style_cell(label_cell, caption_size, body_font, bold=True,
                             text_color=header_text, fill_color=header_fill,
                             alignment=PP_ALIGN.CENTER)

            for c_idx in range(len(col_values)):
                value = float(data[r_idx][c_idx])
                cell = table.cell(r_idx + 1, c_idx + 1)
                cell.text = _fmt_num(value)
                shade = ConditionalFormatter.heatmap_gradient(
                    value, min_val, max_val, theme
                )
                self._style_cell(cell, caption_size, body_font, bold=False,
                                 text_color=body_text,
                                 fill_color=hex_to_rgb(shade),
                                 alignment=PP_ALIGN.RIGHT)

    @staticmethod
    def _style_cell(
        cell,
        font_size,
        font_family: str,
        bold: bool,
        text_color,
        fill_color,
        alignment=None,
    ) -> None:
        """Apply deterministic styling to a table cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
        for paragraph in cell.text_frame.paragraphs:
            if alignment is not None:
                paragraph.alignment = alignment
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.name = font_family
                run.font.bold = bold
                run.font.color.rgb = text_color
