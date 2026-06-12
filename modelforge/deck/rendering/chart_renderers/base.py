"""Base chart renderer — abstract interface and shared theme formatting."""

from __future__ import annotations

from abc import ABC, abstractmethod
from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Pt

if TYPE_CHECKING:
    from pptx.chart.chart import Chart
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import Position
    from modelforge.deck.themes.types import ResolvedTheme


class BaseChartRenderer(ABC):
    """Abstract base class for all chart renderers.

    Subclasses implement `render()` to create a native python-pptx chart
    on a slide. The base class provides shared theme formatting helpers.
    """

    @abstractmethod
    def render(
        self,
        slide: Slide,
        chart_data: object,
        position: Position,
        theme: ResolvedTheme,
    ) -> None:
        """Render a chart onto a python-pptx slide at the given position."""

    def _apply_theme_formatting(
        self,
        chart: Chart,
        theme: ResolvedTheme,
        *,
        title: str | None = None,
    ) -> None:
        """Apply theme colors, fonts, and styling to a chart.

        - Sets chart series colors from theme.chart_colors (cycling if more series than colors)
        - Sets tick label font to theme.typography.body_family, caption scale size
        - Sets legend font to body_family, caption size
        - Positions legend at bottom
        - Sets chart title font to heading_family, h3 scale if title exists
        """
        colors = theme.chart_colors
        typo = theme.typography

        # -- Series colors --
        if colors:
            for idx, series in enumerate(chart.series):
                color_hex = colors[idx % len(colors)].lstrip("#")
                try:
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(color_hex)
                except Exception:
                    pass
                # Also set line color for line/area charts
                try:
                    line = series.format.line
                    line.color.rgb = RGBColor.from_string(color_hex)
                except Exception:
                    pass

        # -- Title --
        if title:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = title
            h3_size = typo.scale.get("h3", 28)
            for paragraph in chart.chart_title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = typo.heading_family
                    run.font.size = Pt(h3_size)
        else:
            chart.has_title = False

        # -- Legend --
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        caption_size = typo.scale.get("caption", 14)
        try:
            chart.legend.font.name = typo.body_family
            chart.legend.font.size = Pt(caption_size)
        except Exception:
            pass

        # -- Axis tick labels --
        # Pie/doughnut charts raise ValueError for axis access, so wrap in try
        for axis_attr in ("category_axis", "value_axis"):
            try:
                axis = getattr(chart, axis_attr, None)
            except (ValueError, AttributeError):
                continue
            if axis is not None:
                try:
                    axis.tick_labels.font.name = typo.body_family
                    axis.tick_labels.font.size = Pt(caption_size)
                except Exception:
                    pass

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """Convert '#RRGGBB' to RGBColor."""
        return RGBColor.from_string(hex_color.lstrip("#"))
