"""Shape element renderers -- shape, divider, spacer, logo."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from modelforge.deck.rendering.element_renderers.base import BaseElementRenderer
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import BaseElement, Position
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)

# Mapping from IR shape names to MSO_SHAPE enum values
_SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "circle": MSO_SHAPE.OVAL,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
}


class ShapeRenderer(BaseElementRenderer):
    """Renders shape elements (rectangle, circle, rounded_rect, arrow, line)."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        shape_type = content.shape
        fill_color = content.fill
        stroke_color = content.stroke

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        if shape_type == "line":
            # Use connector for lines
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT (enum value 1)
                left,
                top + height // 2,
                left + width,
                top + height // 2,
            )
            if stroke_color:
                connector.line.color.rgb = hex_to_rgb(stroke_color)
            else:
                connector.line.color.rgb = hex_to_rgb(theme.colors.text_muted)
            connector.line.width = Pt(2)
        else:
            mso_shape = _SHAPE_MAP.get(shape_type, MSO_SHAPE.RECTANGLE)
            shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            else:
                shape.fill.background()

            if stroke_color:
                shape.line.color.rgb = hex_to_rgb(stroke_color)
                shape.line.width = Pt(1)
            else:
                shape.line.fill.background()


class DividerRenderer(BaseElementRenderer):
    """Renders a thin horizontal divider line."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        left = Inches(position.x)
        top = Inches(position.y + position.height / 2)
        right = Inches(position.x + position.width)

        connector = slide.shapes.add_connector(
            1,  # STRAIGHT
            left,
            top,
            right,
            top,
        )
        connector.line.color.rgb = hex_to_rgb(theme.colors.text_muted)
        connector.line.width = Pt(2)


class SpacerRenderer(BaseElementRenderer):
    """No-op renderer -- spacers are layout-only elements."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        # Spacers affect layout only, nothing to render
        pass


class LogoRenderer(BaseElementRenderer):
    """Renders logo elements as images."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        from modelforge.deck.rendering.element_renderers.image import ImageRenderer

        # Treat logo as an image element
        content = element.content
        url = content.url

        if url:
            # Create a mock image element for the ImageRenderer
            from modelforge.deck.ir.elements.visual import ImageContent, ImageElement

            image_element = ImageElement(
                content=ImageContent(url=url, alt_text="Logo", fit="contain")
            )
            ImageRenderer().render(slide, image_element, position, theme)
        else:
            # No logo URL -- render a small placeholder
            logger.debug("Logo element has no URL, skipping render")


__all__ = [
    "DividerRenderer",
    "LogoRenderer",
    "ShapeRenderer",
    "SpacerRenderer",
]
