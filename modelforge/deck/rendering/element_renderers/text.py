"""Text element renderers -- heading, body, bullet list, numbered list, callout, quote."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from modelforge.deck.rendering.element_renderers.base import BaseElementRenderer
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import BaseElement, Position
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)

# Mapping from element type to theme typography scale key and font family source
_TYPE_SCALE_MAP = {
    "heading": "h2",
    "subheading": "subtitle",
    "body_text": "body",
    "footnote": "footnote",
    "label": "caption",
}


def _get_font_size(element_type: str, theme: ResolvedTheme, content=None) -> int:
    """Get the font size in points for an element type from theme scale."""
    if element_type == "heading" and content and hasattr(content, "level"):
        level = content.level
        if hasattr(level, "value"):
            level = level.value
        return theme.typography.scale.get(level, 36)
    scale_key = _TYPE_SCALE_MAP.get(element_type, "body")
    return theme.typography.scale.get(scale_key, 18)


def _get_font_family(element_type: str, theme: ResolvedTheme) -> str:
    """Get the font family for an element type."""
    if element_type in ("heading", "subheading"):
        return resolve_font_name(theme.typography.heading_family)
    return resolve_font_name(theme.typography.body_family)


def _is_bold(element_type: str) -> bool:
    """Determine if an element type should be bold."""
    return element_type in ("heading", "subheading")


class TextRenderer(BaseElementRenderer):
    """Renders heading, subheading, body_text, footnote, and label elements as text boxes."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        text = content.text
        element_type = element.type
        if hasattr(element_type, "value"):
            element_type = element_type.value

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text

        font_size = _get_font_size(element_type, theme, content)
        font_family = _get_font_family(element_type, theme)
        bold = _is_bold(element_type)

        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs:
            run.text = text
            p.text = ""

        run.font.size = Pt(font_size)
        run.font.name = font_family
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)


class BulletListRenderer(BaseElementRenderer):
    """Renders bullet list elements with one paragraph per item."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        items = content.items
        style = content.style  # disc, dash, arrow

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        font_size = Pt(theme.typography.scale.get("body", 18))
        font_family = resolve_font_name(theme.typography.body_family)

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            if style == "disc":
                # Use bullet character
                prefix = "\u2022 "
            elif style == "dash":
                prefix = "- "
            elif style == "arrow":
                prefix = "\u25b8 "
            else:
                prefix = "\u2022 "

            p.text = f"{prefix}{item}"
            p.level = 0

            for run in p.runs:
                run.font.size = font_size
                run.font.name = font_family
                run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)


class NumberedListRenderer(BaseElementRenderer):
    """Renders numbered list elements with sequential numbering."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        items = content.items
        start = content.start

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        font_size = Pt(theme.typography.scale.get("body", 18))
        font_family = resolve_font_name(theme.typography.body_family)

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{start + i}. {item}"

            for run in p.runs:
                run.font.size = font_size
                run.font.name = font_family
                run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)


class CalloutBoxRenderer(BaseElementRenderer):
    """Renders callout box elements as rounded rectangles with text overlay."""

    _STYLE_COLOR_MAP = {
        "info": "accent",
        "warning": "warning",
        "success": "positive",
        "error": "negative",
    }

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        content = element.content
        style = content.style
        text = content.text

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        # Create rounded rectangle shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )

        # Set fill color based on callout style
        color_attr = self._STYLE_COLOR_MAP.get(style, "accent")
        fill_color = getattr(theme.colors, color_attr, theme.colors.accent)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)

        # Set line to no outline
        shape.line.fill.background()

        # Add text inside the shape
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text

        font_size = Pt(theme.typography.scale.get("body", 18))
        font_family = resolve_font_name(theme.typography.body_family)

        for run in tf.paragraphs[0].runs:
            run.font.size = font_size
            run.font.name = font_family
            run.font.color.rgb = hex_to_rgb("#FFFFFF")


class PullQuoteRenderer(BaseElementRenderer):
    """Renders pull quote elements with italic text and optional attribution."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        text = content.text
        attribution = content.attribution

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        # Quote text - italic
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"\u201c{text}\u201d"
        run.font.size = Pt(theme.typography.scale.get("h3", 28))
        run.font.name = resolve_font_name(theme.typography.heading_family)
        run.font.italic = True
        run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)

        # Attribution line
        if attribution:
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = f"\u2014 {attribution}"
            run2.font.size = Pt(theme.typography.scale.get("caption", 14))
            run2.font.name = resolve_font_name(theme.typography.body_family)
            run2.font.italic = False
            run2.font.color.rgb = hex_to_rgb(theme.colors.text_muted)


__all__ = [
    "BulletListRenderer",
    "CalloutBoxRenderer",
    "NumberedListRenderer",
    "PullQuoteRenderer",
    "TextRenderer",
]
