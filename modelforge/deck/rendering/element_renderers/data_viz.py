"""Data visualization element renderers -- KPI card, metric group, progress bar, gauge, sparkline."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from modelforge.deck.rendering.element_renderers.base import BaseElementRenderer
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import BaseElement, Position
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)


class KpiCardRenderer(BaseElementRenderer):
    """Renders KPI card with large value and small label, optional change indicator."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        label = content.label
        value = str(content.value)

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        # Value - large, bold, primary color
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = value
        run.font.size = Pt(theme.typography.scale.get("h2", 36))
        run.font.name = resolve_font_name(theme.typography.heading_family)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(theme.colors.primary)

        # Label - small, muted
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = label
        run2.font.size = Pt(theme.typography.scale.get("caption", 14))
        run2.font.name = resolve_font_name(theme.typography.body_family)
        run2.font.color.rgb = hex_to_rgb(theme.colors.text_muted)

        # Change indicator
        if content.change is not None:
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            run3 = p3.add_run()

            direction = content.change_direction or ("up" if content.change >= 0 else "down")
            arrow = "\u25b2" if direction == "up" else "\u25bc" if direction == "down" else "\u25b6"
            run3.text = f"{arrow} {abs(content.change):.1f}%"
            run3.font.size = Pt(theme.typography.scale.get("caption", 14))
            run3.font.name = resolve_font_name(theme.typography.body_family)

            if direction == "up":
                run3.font.color.rgb = hex_to_rgb(theme.colors.positive)
            elif direction == "down":
                run3.font.color.rgb = hex_to_rgb(theme.colors.negative)
            else:
                run3.font.color.rgb = hex_to_rgb(theme.colors.text_muted)


class MetricGroupRenderer(BaseElementRenderer):
    """Renders a group of KPI cards side by side."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        metrics = content.metrics

        if not metrics:
            return

        card_width = position.width / len(metrics)
        kpi_renderer = KpiCardRenderer()

        for i, metric in enumerate(metrics):
            from modelforge.deck.ir.elements.base import Position as Pos
            from modelforge.deck.ir.elements.data import KpiCardElement

            card_pos = Pos(
                x=position.x + i * card_width,
                y=position.y,
                width=card_width,
                height=position.height,
            )
            card_element = KpiCardElement(content=metric)
            kpi_renderer.render(slide, card_element, card_pos, theme)


class ProgressBarRenderer(BaseElementRenderer):
    """Renders a progress bar with track and fill shapes plus label."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        label = content.label
        value = content.value
        max_value = content.max_value
        ratio = min(value / max_value, 1.0) if max_value > 0 else 0.0

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        bar_height = min(height // 3, Inches(0.4))
        bar_top = top + (height - bar_height) // 2

        # Track (full width, surface color)
        track = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, bar_top, width, bar_height,
        )
        track.fill.solid()
        track.fill.fore_color.rgb = hex_to_rgb(theme.colors.surface)
        track.line.fill.background()

        # Fill (proportional width, primary color)
        fill_width = int(width * ratio)
        if fill_width > 0:
            fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, bar_top, fill_width, bar_height,
            )
            fill.fill.solid()
            fill.fill.fore_color.rgb = hex_to_rgb(theme.colors.primary)
            fill.line.fill.background()

        # Label text above the bar
        label_height = Inches(0.3)
        txbox = slide.shapes.add_textbox(left, top, width, label_height)
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}: {value:.0f}/{max_value:.0f} ({ratio * 100:.0f}%)"
        run.font.size = Pt(theme.typography.scale.get("caption", 14))
        run.font.name = resolve_font_name(theme.typography.body_family)
        run.font.color.rgb = hex_to_rgb(theme.colors.text_primary)


class GaugeRenderer(BaseElementRenderer):
    """Renders gauge as a progress bar (simplified from circular arc)."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        label = content.label
        value = content.value
        min_val = content.min_value
        max_val = content.max_value

        # Normalize to 0-100 scale for progress bar rendering
        range_val = max_val - min_val
        normalized = ((value - min_val) / range_val * 100) if range_val > 0 else 0

        # Reuse ProgressBarRenderer internally
        from modelforge.deck.ir.elements.base import Position as Pos
        from modelforge.deck.ir.elements.data import ProgressBarContent, ProgressBarElement

        pb_element = ProgressBarElement(
            content=ProgressBarContent(
                label=label,
                value=normalized,
                max_value=100,
            )
        )
        ProgressBarRenderer().render(slide, pb_element, position, theme)


class SparklineRenderer(BaseElementRenderer):
    """Renders sparkline as a series of connected line segments."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        values = content.values

        if not values or len(values) < 2:
            return

        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        min_val = min(values)
        max_val = max(values)
        val_range = max_val - min_val if max_val != min_val else 1.0

        n = len(values)
        step_x = width / (n - 1) if n > 1 else width

        # Draw line segments between consecutive points
        for i in range(n - 1):
            x1 = left + int(i * step_x)
            y1 = top + int(height * (1 - (values[i] - min_val) / val_range))
            x2 = left + int((i + 1) * step_x)
            y2 = top + int(height * (1 - (values[i + 1] - min_val) / val_range))

            connector = slide.shapes.add_connector(
                1,  # STRAIGHT
                x1, y1, x2, y2,
            )
            connector.line.color.rgb = hex_to_rgb(theme.colors.primary)
            connector.line.width = Pt(2)

        # Optional label
        if content.label:
            txbox = slide.shapes.add_textbox(left, top + height, width, Inches(0.3))
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = content.label
            run.font.size = Pt(theme.typography.scale.get("footnote", 10))
            run.font.name = resolve_font_name(theme.typography.body_family)
            run.font.color.rgb = hex_to_rgb(theme.colors.text_muted)


__all__ = [
    "GaugeRenderer",
    "KpiCardRenderer",
    "MetricGroupRenderer",
    "ProgressBarRenderer",
    "SparklineRenderer",
]
