"""Image element renderer -- renders images from URL, base64, or placeholder."""

from __future__ import annotations

import base64
import io
import logging
from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from modelforge.deck.rendering.element_renderers.base import BaseElementRenderer
from modelforge.deck.rendering.utils import hex_to_rgb, resolve_font_name

if TYPE_CHECKING:
    from pptx.slide import Slide

    from modelforge.deck.ir.elements.base import BaseElement, Position
    from modelforge.deck.themes.types import ResolvedTheme

logger = logging.getLogger(__name__)


class ImageRenderer(BaseElementRenderer):
    """Renders image elements from base64, URL, or as a placeholder."""

    def render(self, slide: Slide, element: BaseElement, position: Position, theme: ResolvedTheme) -> None:
        content = element.content
        left = Inches(position.x)
        top = Inches(position.y)
        width = Inches(position.width)
        height = Inches(position.height)

        image_bytes = None

        # Try base64 first
        if content.base64:
            try:
                image_bytes = base64.b64decode(content.base64)
            except Exception:
                logger.warning("Failed to decode base64 image, using placeholder")

        # Try URL download (SSRF-guarded: public http(s) hosts only, no
        # redirects; connected-IP re-checked to close the DNS-rebind TOCTOU)
        if image_bytes is None and content.url:
            try:
                import httpx

                from modelforge.deck.security.url_guard import (
                    assert_connected_ip_public,
                    validate_public_url,
                )

                safe_url = validate_public_url(content.url)
                with httpx.Client(
                    timeout=10.0,
                    event_hooks={"response": [assert_connected_ip_public]},
                ) as client:
                    response = client.get(safe_url, follow_redirects=False)
                    response.raise_for_status()
                    image_bytes = response.content
            except Exception:
                logger.warning("Failed to download image from URL: %s", content.url)

        if image_bytes is not None:
            # Render actual image
            image_stream = io.BytesIO(image_bytes)

            fit = content.fit
            if fit == "fill":
                # Use position dimensions directly
                slide.shapes.add_picture(image_stream, left, top, width, height)
            elif fit == "contain":
                # Calculate aspect-ratio-preserving dimensions
                self._add_contain(slide, image_stream, left, top, width, height)
            elif fit == "cover":
                # For cover mode, just use full dimensions (cropping not easily done)
                slide.shapes.add_picture(image_stream, left, top, width, height)
            else:
                slide.shapes.add_picture(image_stream, left, top, width, height)
        else:
            # Create placeholder rectangle with alt_text
            self._add_placeholder(slide, element, left, top, width, height, theme)

    def _add_contain(self, slide, image_stream, left, top, width, height) -> None:
        """Add image with contain fit (aspect ratio preserved, centered)."""
        from PIL import Image

        # Read image to get dimensions
        image_stream.seek(0)
        img = Image.open(image_stream)
        img_width, img_height = img.size
        image_stream.seek(0)

        # Calculate aspect-ratio-preserving dimensions
        pos_width_px = width
        pos_height_px = height

        aspect = img_width / img_height
        pos_aspect = pos_width_px / pos_height_px

        if aspect > pos_aspect:
            # Image is wider - fit to width
            new_width = width
            new_height = int(width / aspect)
        else:
            # Image is taller - fit to height
            new_height = height
            new_width = int(height * aspect)

        # Center within bounds
        x_offset = left + (width - new_width) // 2
        y_offset = top + (height - new_height) // 2

        slide.shapes.add_picture(image_stream, x_offset, y_offset, new_width, new_height)

    def _add_placeholder(self, slide, element, left, top, width, height, theme) -> None:
        """Add a placeholder rectangle with alt text when no image source available."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(theme.colors.surface)

        # Add alt_text as centered text
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        alt_text = element.content.alt_text or "Image"
        run.text = alt_text
        run.font.size = Pt(theme.typography.scale.get("caption", 14))
        run.font.name = resolve_font_name(theme.typography.body_family)
        run.font.color.rgb = hex_to_rgb(theme.colors.text_muted)


__all__ = ["ImageRenderer"]
