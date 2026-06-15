"""Ground-up development RE deck adapter tests.

Exercises the third deck-mappable template (``development_re``) end-to-end:

    development_re spec → build (trust/moat + deterministic finishing + manifest)
                        → adapt (fail-closed: manifest verify + CERTIFIED audit)
                        → extract dev facts (DevSchedule TDC / forward-NOI exit +
                          Returns unlevered + levered equity IRR / MoIC / LP pref
                          + GP promote, every fact source-cited)
                        → compose ic_memo / teaser (no hollow slides)
                        → render → deterministic stamp (real SHAs on the
                          Certification slide).

One development_re workbook is built ONCE per module and shared by every test
(template build + trust/moat + recalc audits are session-expensive).
"""

from __future__ import annotations

import hashlib
import shutil
from pathlib import Path
from types import SimpleNamespace

import pytest
import yaml

REPO_ROOT = Path(__file__).resolve().parents[2]
SPEC_PATH = REPO_ROOT / "examples" / "development_pbsa_genericcity.yaml"


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _slide_texts(slide) -> list[str]:
    """All non-empty text fragments on a rendered pptx slide (+ chart marker)."""
    out: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            out.append(shape.text_frame.text.strip())
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        out.append(cell.text.strip())
        if getattr(shape, "has_chart", False):
            out.append("[CHART]")
    return out


def _assert_no_hollow_slides(pptx_path: Path) -> None:
    """Every slide must carry content beyond a lone heading (or a chart)."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    for i, slide in enumerate(prs.slides):
        texts = _slide_texts(slide)
        content = [t for t in texts if t != "[CHART]"]
        has_chart = "[CHART]" in texts
        assert (len(content) >= 2) or has_chart, (
            f"slide {i} is hollow (only {content!r}, chart={has_chart})"
        )


# ─────────────────────────────────────────────────────────────────────────────
# Fixtures
# ─────────────────────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def built_dev_workbook(tmp_path_factory) -> Path:
    """Build the development_re example exactly as `modelforge build` ships it."""
    out_dir = tmp_path_factory.mktemp("deck_dev")
    from modelforge.cli import _inject_trust_moat_and_finish, _load_spec_class
    from modelforge.templates import build_model

    spec_bytes = SPEC_PATH.read_bytes()
    raw = yaml.safe_load(spec_bytes)
    spec_class = _load_spec_class(raw["model_type"])
    spec = spec_class.model_validate(raw)

    xlsx_out = out_dir / "development_pbsa_genericcity.xlsx"
    xlsx, _graph = build_model(
        spec, xlsx_out,
        spec_source_bytes=spec_bytes,
        spec_source_path=SPEC_PATH,
    )
    _inject_trust_moat_and_finish(xlsx, spec, spec_bytes, SPEC_PATH, quiet=True)
    xlsx = Path(xlsx)
    assert xlsx.exists()
    assert xlsx.with_suffix(".manifest.json").exists()
    return xlsx


@pytest.fixture(scope="module")
def dev_facts(built_dev_workbook: Path):
    from modelforge.deck.adapter import adapt_workbook

    return adapt_workbook(built_dev_workbook)


@pytest.fixture(scope="module")
def dev_ic_deck(built_dev_workbook: Path):
    from modelforge.deck.pipeline import build_deck_from_workbook

    return build_deck_from_workbook(
        built_dev_workbook,
        deck_type="ic_memo",
        out_path=built_dev_workbook.with_name("dev_ic.pptx"),
    )


# ─────────────────────────────────────────────────────────────────────────────
# (a) development_re is recognised + adapts to DealFacts-compatible facts
# ─────────────────────────────────────────────────────────────────────────────


class TestDevelopmentREAdapts:
    def test_template_in_supported_set(self):
        from modelforge.deck.adapter import SUPPORTED_DECK_TEMPLATES

        assert "development_re" in SUPPORTED_DECK_TEMPLATES

    def test_detect_template_from_redflags_line(self, built_dev_workbook):
        from openpyxl import load_workbook

        from modelforge.deck.adapter import detect_template

        wb = load_workbook(built_dev_workbook, data_only=False, keep_links=True)
        assert detect_template(wb) == "development_re"

    def test_detect_template_from_devschedule_sheet(self, tmp_path):
        """No RedFlags Template line → discriminated by the DevSchedule sheet."""
        from openpyxl import Workbook, load_workbook

        from modelforge.deck.adapter import detect_template

        xlsx = tmp_path / "probe_dev.xlsx"
        wb = Workbook()
        wb.active.title = "Cover"
        wb.create_sheet("DevSchedule")
        wb.create_sheet("Returns")
        wb.save(xlsx)
        assert detect_template(load_workbook(xlsx)) == "development_re"

    def test_adapt_returns_dev_facts_with_lineage(self, dev_facts):
        assert dev_facts.template == "development_re"
        assert dev_facts.audit_verdict == "CERTIFIED"
        refs = dev_facts.source_refs
        for key in ("levered_equity_irr", "unlevered_irr", "equity_moic",
                    "total_dev_cost", "exit_value_gross", "forward_noi",
                    "peak_debt", "net_exit_proceeds", "lp_pref",
                    "total_equity_contrib", "exit_cap_rate"):
            assert key in refs, f"missing dev fact {key}"
            assert "!" in refs[key], f"{key} ref {refs[key]!r} is not sheet!cell"

    def test_levered_irr_reads_from_returns(self, dev_facts):
        # The levered equity IRR is the headline development return — sourced
        # from the Returns sheet's levered equity IRR cell, recomputed (never
        # cached).
        irr = dev_facts.facts["levered_equity_irr"]
        assert irr.ref.startswith("Returns!")
        assert 0.0 < float(irr.value) < 1.0

    def test_unlevered_irr_is_distinct_and_lower(self, dev_facts):
        # A geared development has a levered equity IRR materially above the
        # unlevered project IRR (positive leverage). Both are live + distinct.
        lev = float(dev_facts.facts["levered_equity_irr"].value)
        unlev = float(dev_facts.facts["unlevered_irr"].value)
        assert 0.0 < unlev < 1.0
        assert lev > unlev

    def test_equity_invested_aggregated_from_schedule_row(self, dev_facts):
        # Equity invested is summed cell-by-cell from the schedule equity-CF
        # row (the isolated Returns SUMIF cell is not engine-recomputable), so
        # the lineage ref is a composite over the DevSchedule equity-CF row.
        cf = dev_facts.facts["total_equity_contrib"]
        assert cf.ref.startswith("DevSchedule!")
        assert float(cf.value) > 0

    def test_deal_facts_populates_dev_specific_fields(self, dev_facts):
        deal = dev_facts.deal_facts()
        assert deal.vertical == "re"
        assert deal.levered_irr_pct is not None
        assert deal.unlevered_irr_pct is not None      # development-specific
        assert deal.levered_irr_pct > deal.unlevered_irr_pct
        assert deal.total_size_eur_m > 0               # total development cost
        assert deal.equity_required_eur_m > 0
        assert deal.debt_eur_m > 0                      # peak senior debt
        # Capital stack split into Senior dev debt + LP + GP equity.
        assert any("Senior" in t for t in deal.capital_stack_tranches)
        assert any("LP" in t for t in deal.capital_stack_tranches)
        assert any("GP" in t for t in deal.capital_stack_tranches)
        # development_re carries no exit-multiple comps / sensitivity grid →
        # omitted (no-hollow-slides invariant).
        assert deal.comparable_transactions == []
        assert deal.sensitivity_values == []
        # sector metrics carry the development economics.
        assert "total_development_cost" in deal.sector_metrics
        assert "loan_to_cost" in deal.sector_metrics

    def test_teaser_facts_omits_revenue_ebitda(self, dev_facts):
        """A ground-up development has no LTM revenue/EBITDA → those teaser
        fields are None so the snapshot KPI cards drop, not render empty."""
        teaser = dev_facts.teaser_facts()
        assert teaser.vertical == "re"
        assert teaser.revenue_eur_m is None
        assert teaser.ebitda_eur_m is None
        assert teaser.enterprise_value_eur_m is not None  # total dev cost
        assert teaser.anonymized is True


# ─────────────────────────────────────────────────────────────────────────────
# (b) ic_memo + teaser build with NO hollow slides + real SHAs
# ─────────────────────────────────────────────────────────────────────────────


class TestDevelopmentREDecks:
    def test_ic_memo_builds_certified(self, dev_ic_deck):
        assert dev_ic_deck.pptx_path.exists()
        assert dev_ic_deck.pptx_path.stat().st_size > 10_000
        assert dev_ic_deck.template == "development_re"
        assert dev_ic_deck.audit_verdict == "CERTIFIED"
        # 12 composer slides − comps − sensitivity (no data → omitted)
        # + mandatory certification slide.
        assert dev_ic_deck.slide_count == 11

    def test_ic_memo_no_hollow_slides(self, dev_ic_deck):
        _assert_no_hollow_slides(dev_ic_deck.pptx_path)

    def test_ic_memo_certification_slide_has_real_shas(self, dev_ic_deck):
        from pptx import Presentation

        prs = Presentation(str(dev_ic_deck.pptx_path))
        last = prs.slides[-1]
        blob = " | ".join(_slide_texts(last))
        assert "Certification & Red Flags" in blob
        assert "CERTIFIED" in blob
        assert dev_ic_deck.workbook_sha256 in blob
        assert dev_ic_deck.spec_sha256 in blob
        assert len(dev_ic_deck.workbook_sha256) == 64
        assert len(dev_ic_deck.spec_sha256) == 64

    def test_teaser_builds_no_hollow_slides(self, built_dev_workbook):
        from modelforge.deck.determinism import read_pptx_stamp
        from modelforge.deck.pipeline import build_deck_from_workbook

        res = build_deck_from_workbook(
            built_dev_workbook,
            deck_type="teaser",
            out_path=built_dev_workbook.with_name("dev_teaser.pptx"),
        )
        assert res.template == "development_re"
        assert res.slide_count == 7  # 6 composer + certification
        _assert_no_hollow_slides(res.pptx_path)
        stamp = read_pptx_stamp(res.pptx_path)
        assert stamp["workbook_sha256"] == res.workbook_sha256
        assert stamp["spec_sha256"] == res.spec_sha256


# ─────────────────────────────────────────────────────────────────────────────
# (c) byte-stability: same workbook → two dev decks → identical SHA-256
# ─────────────────────────────────────────────────────────────────────────────


class TestDevelopmentREByteStability:
    def test_same_workbook_two_decks_byte_identical(
        self, built_dev_workbook, dev_ic_deck
    ):
        from modelforge.deck.pipeline import build_deck_from_workbook

        second = build_deck_from_workbook(
            built_dev_workbook,
            deck_type="ic_memo",
            out_path=built_dev_workbook.with_name("dev_ic_b.pptx"),
        )
        assert _sha256(dev_ic_deck.pptx_path) == _sha256(second.pptx_path)


# ─────────────────────────────────────────────────────────────────────────────
# (d) fail-closed: an uncertified / tampered dev workbook is refused
# ─────────────────────────────────────────────────────────────────────────────


class TestDevelopmentREFailClosed:
    def test_missing_manifest_refused(self, built_dev_workbook, tmp_path):
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        orphan = tmp_path / "orphan_dev.xlsx"
        shutil.copy(built_dev_workbook, orphan)  # no manifest sidecar
        with pytest.raises(DeckAdapterError, match="[Mm]anifest"):
            build_deck_from_workbook(orphan)

    def test_tampered_workbook_refused(self, built_dev_workbook, tmp_path):
        from openpyxl import load_workbook

        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        tampered = tmp_path / "tampered_dev.xlsx"
        shutil.copy(built_dev_workbook, tampered)
        shutil.copy(
            built_dev_workbook.with_suffix(".manifest.json"),
            tampered.with_suffix(".manifest.json"),
        )
        # Flip a workbook cell — bytes no longer match the manifest hash.
        wb = load_workbook(tampered)
        wb["Assumptions"]["I6"] = 99.0
        wb.save(tampered)
        with pytest.raises(DeckAdapterError,
                           match="verification FAILED|not the bytes"):
            build_deck_from_workbook(tampered)

    def test_uncertified_workbook_refused(self, built_dev_workbook, tmp_path):
        """Manifest verifies but the audit is not CERTIFIED (cached #REF!
        literal) → refused at the certification step, before extraction."""
        from openpyxl import load_workbook

        from modelforge.analytics.manifest import write_manifest
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        bad = tmp_path / "uncertified_dev.xlsx"
        shutil.copy(built_dev_workbook, bad)
        wb = load_workbook(bad)
        wb["DevSchedule"]["Z2"] = "#REF!"  # cached error literal → audit FAIL
        wb.save(bad)
        # Re-manifest so the manifest gate PASSES and the refusal provably
        # comes from the certification audit, not a hash mismatch.
        write_manifest(bad, SimpleNamespace(), spec_source_bytes=b"tampered")
        with pytest.raises(DeckAdapterError, match="not\\s+CERTIFIED"):
            build_deck_from_workbook(bad)
