"""Real-estate deck adapter tests.

Exercises the second deck-mappable template (``real_estate``) end-to-end:

    real_estate spec → build (trust/moat + deterministic finishing + manifest)
                     → adapt (fail-closed: manifest verify + CERTIFIED audit)
                     → extract RE facts (DCF NOI/exit + Financing equity IRR /
                       MoIC / LP pref + GP promote, every fact source-cited)
                     → compose ic_memo / teaser (no hollow slides)
                     → render → deterministic stamp (real SHAs on the
                       Certification slide).

One real_estate workbook is built ONCE per module and shared by every test
(template build + trust/moat + recalc audits are session-expensive).
"""

from __future__ import annotations

import hashlib
import shutil
from pathlib import Path
from types import SimpleNamespace

import pytest
import yaml

REPO_ROOT = Path(__file__).resolve().parents[2]
SPEC_PATH = REPO_ROOT / "examples" / "real_estate_pbsa.yaml"


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _slide_texts(slide) -> list[str]:
    """All non-empty text fragments on a rendered pptx slide (+ chart marker)."""
    out: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            out.append(shape.text_frame.text.strip())
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        out.append(cell.text.strip())
        if getattr(shape, "has_chart", False):
            out.append("[CHART]")
    return out


def _assert_no_hollow_slides(pptx_path: Path) -> None:
    """Every slide must carry content beyond a lone heading (or a chart)."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    for i, slide in enumerate(prs.slides):
        texts = _slide_texts(slide)
        content = [t for t in texts if t != "[CHART]"]
        has_chart = "[CHART]" in texts
        assert (len(content) >= 2) or has_chart, (
            f"slide {i} is hollow (only {content!r}, chart={has_chart})"
        )


# ─────────────────────────────────────────────────────────────────────────────
# Fixtures
# ─────────────────────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def built_re_workbook(tmp_path_factory) -> Path:
    """Build the real_estate example exactly as `modelforge build` ships it."""
    out_dir = tmp_path_factory.mktemp("deck_re")
    from modelforge.cli import _inject_trust_moat_and_finish, _load_spec_class
    from modelforge.templates import build_model

    spec_bytes = SPEC_PATH.read_bytes()
    raw = yaml.safe_load(spec_bytes)
    spec_class = _load_spec_class(raw["model_type"])
    spec = spec_class.model_validate(raw)

    xlsx_out = out_dir / "real_estate_pbsa.xlsx"
    xlsx, _graph = build_model(
        spec, xlsx_out,
        spec_source_bytes=spec_bytes,
        spec_source_path=SPEC_PATH,
    )
    _inject_trust_moat_and_finish(xlsx, spec, spec_bytes, SPEC_PATH, quiet=True)
    xlsx = Path(xlsx)
    assert xlsx.exists()
    assert xlsx.with_suffix(".manifest.json").exists()
    return xlsx


@pytest.fixture(scope="module")
def re_facts(built_re_workbook: Path):
    from modelforge.deck.adapter import adapt_workbook

    return adapt_workbook(built_re_workbook)


@pytest.fixture(scope="module")
def re_ic_deck(built_re_workbook: Path):
    from modelforge.deck.pipeline import build_deck_from_workbook

    return build_deck_from_workbook(
        built_re_workbook,
        deck_type="ic_memo",
        out_path=built_re_workbook.with_name("re_ic.pptx"),
    )


# ─────────────────────────────────────────────────────────────────────────────
# (a) real_estate is recognised + adapts to DealFacts-compatible facts
# ─────────────────────────────────────────────────────────────────────────────


class TestRealEstateAdapts:
    def test_template_in_supported_set(self):
        from modelforge.deck.adapter import SUPPORTED_DECK_TEMPLATES

        assert "real_estate" in SUPPORTED_DECK_TEMPLATES

    def test_detect_template_from_redflags_line(self, built_re_workbook):
        from openpyxl import load_workbook

        from modelforge.deck.adapter import detect_template

        wb = load_workbook(built_re_workbook, data_only=False, keep_links=True)
        assert detect_template(wb) == "real_estate"

    def test_detect_template_from_dcf_financing_sheets(self, tmp_path):
        """No RedFlags Template line → discriminated by {DCF, Financing}."""
        from openpyxl import Workbook, load_workbook

        from modelforge.deck.adapter import detect_template

        xlsx = tmp_path / "probe_re.xlsx"
        wb = Workbook()
        wb.active.title = "Cover"
        wb.create_sheet("DCF")
        wb.create_sheet("Financing")
        wb.save(xlsx)
        assert detect_template(load_workbook(xlsx)) == "real_estate"

    def test_adapt_returns_re_facts_with_lineage(self, re_facts):
        assert re_facts.template == "real_estate"
        assert re_facts.audit_verdict == "CERTIFIED"
        refs = re_facts.source_refs
        for key in ("levered_equity_irr", "equity_moic", "acquisition_price",
                    "loan_amount", "exit_value_gross", "exit_noi",
                    "exit_cap_rate", "lp_pref", "total_equity_contrib"):
            assert key in refs, f"missing RE fact {key}"
            assert "!" in refs[key], f"{key} ref {refs[key]!r} is not sheet!cell"

    def test_levered_irr_reads_from_financing(self, re_facts):
        # The levered equity IRR is the headline RE return — sourced from the
        # Financing sheet's Equity IRR cell, recomputed (never cached).
        irr = re_facts.facts["levered_equity_irr"]
        assert irr.ref.startswith("Financing!")
        assert 0.0 < float(irr.value) < 1.0

    def test_deal_facts_populates_re_specific_fields(self, re_facts):
        deal = re_facts.deal_facts()
        assert deal.vertical == "re"
        assert deal.levered_irr_pct is not None
        assert deal.yield_on_cost_pct is not None  # going-in NOI / cost
        assert deal.total_size_eur_m > 0  # going-in asset value
        assert deal.equity_required_eur_m > 0
        assert deal.debt_eur_m > 0
        # Capital stack split into Senior + LP + GP equity
        assert "Senior Mortgage" in deal.capital_stack_tranches
        assert any("LP" in t for t in deal.capital_stack_tranches)
        # RE carries no exit-multiple comps / sensitivity grid → omitted.
        assert deal.comparable_transactions == []
        assert deal.sensitivity_values == []

    def test_teaser_facts_omits_revenue_ebitda(self, re_facts):
        """RE has no LTM revenue/EBITDA → those teaser fields are None so the
        snapshot KPI cards drop, not render empty."""
        teaser = re_facts.teaser_facts()
        assert teaser.vertical == "re"
        assert teaser.revenue_eur_m is None
        assert teaser.ebitda_eur_m is None
        assert teaser.enterprise_value_eur_m is not None  # going-in value
        assert teaser.anonymized is True


# ─────────────────────────────────────────────────────────────────────────────
# (b) ic_memo + teaser build with NO hollow slides + real SHAs
# ─────────────────────────────────────────────────────────────────────────────


class TestRealEstateDecks:
    def test_ic_memo_builds_certified(self, re_ic_deck):
        assert re_ic_deck.pptx_path.exists()
        assert re_ic_deck.pptx_path.stat().st_size > 10_000
        assert re_ic_deck.template == "real_estate"
        assert re_ic_deck.audit_verdict == "CERTIFIED"
        # 12 composer slides − comps − sensitivity (no data → omitted)
        # + mandatory certification slide.
        assert re_ic_deck.slide_count == 11

    def test_ic_memo_no_hollow_slides(self, re_ic_deck):
        _assert_no_hollow_slides(re_ic_deck.pptx_path)

    def test_ic_memo_certification_slide_has_real_shas(self, re_ic_deck):
        from pptx import Presentation

        prs = Presentation(str(re_ic_deck.pptx_path))
        last = prs.slides[-1]
        blob = " | ".join(_slide_texts(last))
        assert "Certification & Red Flags" in blob
        assert "CERTIFIED" in blob
        assert re_ic_deck.workbook_sha256 in blob
        assert re_ic_deck.spec_sha256 in blob
        assert len(re_ic_deck.workbook_sha256) == 64
        assert len(re_ic_deck.spec_sha256) == 64

    def test_teaser_builds_no_hollow_slides(self, built_re_workbook):
        from modelforge.deck.determinism import read_pptx_stamp
        from modelforge.deck.pipeline import build_deck_from_workbook

        res = build_deck_from_workbook(
            built_re_workbook,
            deck_type="teaser",
            out_path=built_re_workbook.with_name("re_teaser.pptx"),
        )
        assert res.template == "real_estate"
        assert res.slide_count == 7  # 6 composer + certification
        _assert_no_hollow_slides(res.pptx_path)
        stamp = read_pptx_stamp(res.pptx_path)
        assert stamp["workbook_sha256"] == res.workbook_sha256
        assert stamp["spec_sha256"] == res.spec_sha256


# ─────────────────────────────────────────────────────────────────────────────
# (c) byte-stability: same workbook → two RE decks → identical SHA-256
# ─────────────────────────────────────────────────────────────────────────────


class TestRealEstateByteStability:
    def test_same_workbook_two_decks_byte_identical(
        self, built_re_workbook, re_ic_deck
    ):
        from modelforge.deck.pipeline import build_deck_from_workbook

        second = build_deck_from_workbook(
            built_re_workbook,
            deck_type="ic_memo",
            out_path=built_re_workbook.with_name("re_ic_b.pptx"),
        )
        assert _sha256(re_ic_deck.pptx_path) == _sha256(second.pptx_path)


# ─────────────────────────────────────────────────────────────────────────────
# (d) fail-closed: an uncertified / tampered RE workbook is refused
# ─────────────────────────────────────────────────────────────────────────────


class TestRealEstateFailClosed:
    def test_missing_manifest_refused(self, built_re_workbook, tmp_path):
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        orphan = tmp_path / "orphan_re.xlsx"
        shutil.copy(built_re_workbook, orphan)  # no manifest sidecar
        with pytest.raises(DeckAdapterError, match="[Mm]anifest"):
            build_deck_from_workbook(orphan)

    def test_tampered_workbook_refused(self, built_re_workbook, tmp_path):
        from openpyxl import load_workbook

        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        tampered = tmp_path / "tampered_re.xlsx"
        shutil.copy(built_re_workbook, tampered)
        shutil.copy(
            built_re_workbook.with_suffix(".manifest.json"),
            tampered.with_suffix(".manifest.json"),
        )
        # Flip a workbook cell — bytes no longer match the manifest hash.
        wb = load_workbook(tampered)
        wb["Assumptions"]["I6"] = 99.0
        wb.save(tampered)
        with pytest.raises(DeckAdapterError,
                           match="verification FAILED|not the bytes"):
            build_deck_from_workbook(tampered)

    def test_uncertified_workbook_refused(self, built_re_workbook, tmp_path):
        """Manifest verifies but the audit is not CERTIFIED (cached #REF!
        literal) → refused at the certification step, before extraction."""
        from openpyxl import load_workbook

        from modelforge.analytics.manifest import write_manifest
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        bad = tmp_path / "uncertified_re.xlsx"
        shutil.copy(built_re_workbook, bad)
        wb = load_workbook(bad)
        wb["DCF"]["Z2"] = "#REF!"  # cached error literal → audit FAIL
        wb.save(bad)
        # Re-manifest so the manifest gate PASSES and the refusal provably
        # comes from the certification audit, not a hash mismatch.
        write_manifest(bad, SimpleNamespace(), spec_source_bytes=b"tampered")
        with pytest.raises(DeckAdapterError, match="not\\s+CERTIFIED"):
            build_deck_from_workbook(bad)
