"""Unit tests for element renderers and utility functions."""

from __future__ import annotations

import io
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from modelforge.deck.ir.elements.base import Position
from modelforge.deck.ir.elements.data import (
    GaugeContent,
    GaugeElement,
    KpiCardContent,
    KpiCardElement,
    MetricGroupContent,
    MetricGroupElement,
    ProgressBarContent,
    ProgressBarElement,
    SparklineContent,
    SparklineElement,
    TableContent,
    TableElement,
)
from modelforge.deck.ir.elements.text import (
    BodyTextContent,
    BodyTextElement,
    BulletListContent,
    BulletListElement,
    CalloutBoxContent,
    CalloutBoxElement,
    FootnoteContent,
    FootnoteElement,
    HeadingContent,
    HeadingElement,
    LabelContent,
    LabelElement,
    NumberedListContent,
    NumberedListElement,
    PullQuoteContent,
    PullQuoteElement,
    SubheadingContent,
    SubheadingElement,
)
from modelforge.deck.ir.elements.visual import (
    BackgroundContent,
    BackgroundElement,
    DividerElement,
    ImageContent,
    ImageElement,
    LogoContent,
    LogoElement,
    ShapeContent,
    ShapeElement,
    SpacerElement,
)
from modelforge.deck.ir.enums import ElementType, Transition
from modelforge.deck.themes.types import (
    FooterDefaults,
    LogoDefaults,
    ResolvedTheme,
    ThemeColors,
    ThemeSpacing,
    ThemeTypography,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture()
def theme() -> ResolvedTheme:
    """Minimal resolved theme for testing."""
    return ResolvedTheme(
        name="test-theme",
        description="Test theme",
        colors=ThemeColors(
            primary="#0A1E3D",
            secondary="#1A3A5C",
            accent="#FF6B35",
            background="#FFFFFF",
            surface="#F5F5F5",
            text_primary="#1A1A1A",
            text_secondary="#4A4A4A",
            text_muted="#8A8A8A",
            positive="#28A745",
            negative="#DC3545",
            warning="#FFC107",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
        logo=LogoDefaults(),
        footer=FooterDefaults(),
    )


@pytest.fixture()
def position() -> Position:
    """Standard position for testing."""
    return Position(x=1.0, y=2.0, width=10.0, height=3.0)


@pytest.fixture()
def pptx_slide():
    """Create a real python-pptx slide for testing."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    return slide


# ── Utility Tests ─────────────────────────────────────────────────────────────


class TestHexToRgb:
    def test_with_hash(self):
        from modelforge.deck.rendering.utils import hex_to_rgb

        result = hex_to_rgb("#0A1E3D")
        assert result == RGBColor(0x0A, 0x1E, 0x3D)

    def test_without_hash(self):
        from modelforge.deck.rendering.utils import hex_to_rgb

        result = hex_to_rgb("0A1E3D")
        assert result == RGBColor(0x0A, 0x1E, 0x3D)

    def test_white(self):
        from modelforge.deck.rendering.utils import hex_to_rgb

        result = hex_to_rgb("#FFFFFF")
        assert result == RGBColor(0xFF, 0xFF, 0xFF)


class TestResolveFontName:
    def test_safe_font_passthrough(self):
        from modelforge.deck.rendering.utils import resolve_font_name

        assert resolve_font_name("Calibri") == "Calibri"

    def test_safe_font_arial(self):
        from modelforge.deck.rendering.utils import resolve_font_name

        assert resolve_font_name("Arial") == "Arial"

    def test_unknown_font_fallback(self):
        from modelforge.deck.rendering.utils import resolve_font_name

        assert resolve_font_name("Inter") == "Calibri"

    def test_known_substitute_font(self):
        from modelforge.deck.rendering.utils import resolve_font_name

        # Fira Code is a known monospace font that maps to Consolas
        assert resolve_font_name("Fira Code") == "Consolas"

    def test_truly_unknown_font_fallback(self):
        from modelforge.deck.rendering.utils import resolve_font_name

        # Completely unknown fonts should fall back to the default
        assert resolve_font_name("UnknownFont12345") == "Calibri"


class TestSetSlideBackground:
    def test_sets_solid_fill(self, pptx_slide):
        from modelforge.deck.rendering.utils import set_slide_background

        set_slide_background(pptx_slide, "#0A1E3D")
        bg = pptx_slide.background
        fill = bg.fill
        assert fill.type is not None


class TestSetTransition:
    def test_fade_transition(self, pptx_slide):
        from modelforge.deck.rendering.utils import set_transition

        set_transition(pptx_slide, Transition.FADE)
        # Check that transition XML was added to the slide
        slide_xml = pptx_slide._element.xml
        assert "mc:AlternateContent" in slide_xml or "transition" in slide_xml.lower() or "Transition" in slide_xml

    def test_none_transition(self, pptx_slide):
        from modelforge.deck.rendering.utils import set_transition

        # Should not raise
        set_transition(pptx_slide, Transition.NONE)

    def test_push_transition(self, pptx_slide):
        from modelforge.deck.rendering.utils import set_transition

        set_transition(pptx_slide, Transition.PUSH)
        slide_xml = pptx_slide._element.xml
        assert "transition" in slide_xml.lower() or "Transition" in slide_xml


# ── Text Renderer Tests ──────────────────────────────────────────────────────


class TestTextRenderer:
    def test_heading_creates_textbox(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import TextRenderer

        element = HeadingElement(content=HeadingContent(text="Hello World", level="h1"))
        renderer = TextRenderer()
        renderer.render(pptx_slide, element, position, theme)
        # Should have added at least one shape (textbox)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_body_text_creates_textbox(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import TextRenderer

        element = BodyTextElement(content=BodyTextContent(text="Body text here"))
        renderer = TextRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_subheading_creates_textbox(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import TextRenderer

        element = SubheadingElement(content=SubheadingContent(text="Subheading"))
        renderer = TextRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_footnote_creates_textbox(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import TextRenderer

        element = FootnoteElement(content=FootnoteContent(text="See note 1"))
        renderer = TextRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_label_creates_textbox(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import TextRenderer

        element = LabelElement(content=LabelContent(text="Label"))
        renderer = TextRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


class TestBulletListRenderer:
    def test_creates_textbox_with_bullets(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import BulletListRenderer

        element = BulletListElement(
            content=BulletListContent(items=["Item 1", "Item 2", "Item 3"])
        )
        renderer = BulletListRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_dash_style(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import BulletListRenderer

        element = BulletListElement(
            content=BulletListContent(items=["Item 1"], style="dash")
        )
        renderer = BulletListRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


class TestNumberedListRenderer:
    def test_creates_numbered_paragraphs(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import NumberedListRenderer

        element = NumberedListElement(
            content=NumberedListContent(items=["First", "Second", "Third"], start=1)
        )
        renderer = NumberedListRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_custom_start_number(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import NumberedListRenderer

        element = NumberedListElement(
            content=NumberedListContent(items=["Item A"], start=5)
        )
        renderer = NumberedListRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


class TestCalloutBoxRenderer:
    def test_creates_rounded_rect(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import CalloutBoxRenderer

        element = CalloutBoxElement(
            content=CalloutBoxContent(text="Important info", style="info")
        )
        renderer = CalloutBoxRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_warning_style(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import CalloutBoxRenderer

        element = CalloutBoxElement(
            content=CalloutBoxContent(text="Warning!", style="warning")
        )
        renderer = CalloutBoxRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


class TestPullQuoteRenderer:
    def test_creates_italic_textbox(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import PullQuoteRenderer

        element = PullQuoteElement(
            content=PullQuoteContent(text="To be or not to be")
        )
        renderer = PullQuoteRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_with_attribution(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.text import PullQuoteRenderer

        element = PullQuoteElement(
            content=PullQuoteContent(
                text="To be or not to be", attribution="Shakespeare"
            )
        )
        renderer = PullQuoteRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


# ── Table Renderer Tests ─────────────────────────────────────────────────────


class TestTableRenderer:
    def test_creates_table_shape(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.table import TableRenderer

        element = TableElement(
            content=TableContent(
                headers=["Name", "Age", "City"],
                rows=[["Alice", 30, "NYC"], ["Bob", 25, "LA"]],
            )
        )
        renderer = TableRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_table_with_footer(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.table import TableRenderer

        element = TableElement(
            content=TableContent(
                headers=["Metric", "Value"],
                rows=[["Revenue", 100], ["Cost", 60]],
                footer_row=["Total", 40],
            )
        )
        renderer = TableRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_table_with_highlight_rows(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.table import TableRenderer

        element = TableElement(
            content=TableContent(
                headers=["Item", "Count"],
                rows=[["A", 10], ["B", 20], ["C", 30]],
                highlight_rows=[1],
            )
        )
        renderer = TableRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


# ── Image Renderer Tests ─────────────────────────────────────────────────────


class TestImageRenderer:
    def test_base64_image(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.image import ImageRenderer

        # Minimal 1x1 pixel PNG in base64
        import base64

        # Create a tiny valid PNG
        from PIL import Image

        buf = io.BytesIO()
        img = Image.new("RGB", (10, 10), "red")
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode()

        element = ImageElement(
            content=ImageContent(base64=b64, alt_text="test image")
        )
        renderer = ImageRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_placeholder_when_no_source(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.image import ImageRenderer

        element = ImageElement(
            content=ImageContent(url=None, base64=None, alt_text="Missing image")
        )
        renderer = ImageRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


# ── Shape Renderer Tests ─────────────────────────────────────────────────────


class TestShapeRenderer:
    def test_rectangle(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.shape import ShapeRenderer

        element = ShapeElement(
            content=ShapeContent(shape="rectangle", fill="#FF0000")
        )
        renderer = ShapeRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_circle(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.shape import ShapeRenderer

        element = ShapeElement(
            content=ShapeContent(shape="circle", fill="#00FF00")
        )
        renderer = ShapeRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_rounded_rect(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.shape import ShapeRenderer

        element = ShapeElement(
            content=ShapeContent(shape="rounded_rect", fill="#0000FF")
        )
        renderer = ShapeRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_arrow(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.shape import ShapeRenderer

        element = ShapeElement(content=ShapeContent(shape="arrow"))
        renderer = ShapeRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_line(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.shape import ShapeRenderer

        element = ShapeElement(content=ShapeContent(shape="line", stroke="#000000"))
        renderer = ShapeRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


class TestDividerRenderer:
    def test_creates_line(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.shape import DividerRenderer

        element = DividerElement()
        renderer = DividerRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


class TestSpacerRenderer:
    def test_noop(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.shape import SpacerRenderer

        element = SpacerElement()
        renderer = SpacerRenderer()
        initial_count = len(list(pptx_slide.shapes))
        renderer.render(pptx_slide, element, position, theme)
        final_count = len(list(pptx_slide.shapes))
        # Spacer is a no-op, shape count unchanged
        assert final_count == initial_count


# ── Data Viz Renderer Tests ──────────────────────────────────────────────────


class TestKpiCardRenderer:
    def test_creates_textbox(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.data_viz import KpiCardRenderer

        element = KpiCardElement(
            content=KpiCardContent(label="Revenue", value="$1.2M")
        )
        renderer = KpiCardRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1

    def test_with_change_indicator(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.data_viz import KpiCardRenderer

        element = KpiCardElement(
            content=KpiCardContent(
                label="Revenue", value="$1.2M", change=12.5, change_direction="up"
            )
        )
        renderer = KpiCardRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1


class TestProgressBarRenderer:
    def test_creates_shapes(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers.data_viz import ProgressBarRenderer

        element = ProgressBarElement(
            content=ProgressBarContent(label="Completion", value=75, max_value=100)
        )
        renderer = ProgressBarRenderer()
        renderer.render(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        # Should create at least track + fill shapes
        assert len(shapes) >= 2


# ── Registry Tests ────────────────────────────────────────────────────────────


class TestElementRegistry:
    def test_all_element_types_mapped(self):
        from modelforge.deck.rendering.element_renderers import ELEMENT_RENDERERS

        for et in ElementType:
            assert et.value in ELEMENT_RENDERERS, f"Missing renderer for {et.value}"

    def test_render_element_function(self, pptx_slide, position, theme):
        from modelforge.deck.rendering.element_renderers import render_element

        element = HeadingElement(content=HeadingContent(text="Test"))
        # Should not raise
        render_element(pptx_slide, element, position, theme)
        shapes = list(pptx_slide.shapes)
        assert len(shapes) >= 1
