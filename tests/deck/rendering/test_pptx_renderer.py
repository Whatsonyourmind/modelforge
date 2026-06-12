"""Integration tests for PptxRenderer producing valid PPTX files."""

from __future__ import annotations

import io
from typing import Any

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from modelforge.deck.ir.elements.base import Position
from modelforge.deck.ir.elements.data import (
    KpiCardContent,
    KpiCardElement,
    TableContent,
    TableElement,
)
from modelforge.deck.ir.elements.text import (
    BodyTextContent,
    BodyTextElement,
    BulletListContent,
    BulletListElement,
    HeadingContent,
    HeadingElement,
    PullQuoteContent,
    PullQuoteElement,
    SubheadingContent,
    SubheadingElement,
)
from modelforge.deck.ir.elements.visual import (
    ImageContent,
    ImageElement,
    ShapeContent,
    ShapeElement,
)
from modelforge.deck.ir.enums import SlideType, Transition
from modelforge.deck.ir.presentation import Presentation
from modelforge.deck.ir.slides.universal import (
    AgendaSlide,
    AppendixSlide,
    BulletPointsSlide,
    ChartSlideSlide,
    ComparisonSlide,
    FunnelSlide,
    IconGridSlide,
    ImageWithCaptionSlide,
    KeyMessageSlide,
    MapSlide,
    MatrixSlide,
    OrgChartSlide,
    ProcessFlowSlide,
    QAndASlide,
    QuoteSlide,
    SectionDividerSlide,
    StatsCalloutSlide,
    TableSlideSlide,
    TeamSlide,
    ThankYouSlide,
    TimelineSlide,
    TitleSlide,
    TwoColumnTextSlide,
)
from modelforge.deck.layout.types import LayoutResult, ResolvedPosition
from modelforge.deck.themes.types import (
    FooterDefaults,
    LogoDefaults,
    ResolvedTheme,
    SlideMaster,
    ThemeColors,
    ThemeSpacing,
    ThemeTypography,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture()
def theme() -> ResolvedTheme:
    """Minimal resolved theme for testing."""
    return ResolvedTheme(
        name="test-theme",
        description="Test theme",
        colors=ThemeColors(
            primary="#0A1E3D",
            secondary="#1A3A5C",
            accent="#FF6B35",
            background="#FFFFFF",
            surface="#F5F5F5",
            text_primary="#1A1A1A",
            text_secondary="#4A4A4A",
            text_muted="#8A8A8A",
            positive="#28A745",
            negative="#DC3545",
            warning="#FFC107",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
        slide_masters={
            "title_slide": SlideMaster(background="#0A1E3D"),
            "bullet_points": SlideMaster(background="#FFFFFF"),
        },
        logo=LogoDefaults(),
        footer=FooterDefaults(),
    )


def _make_layout_result(slide, elements_with_positions: list[tuple[Any, Position]] | None = None) -> LayoutResult:
    """Create a LayoutResult with positions assigned to elements."""
    lr = LayoutResult(slide=slide, positions={})
    if elements_with_positions:
        for elem, pos in elements_with_positions:
            elem.position = pos
    return lr


# ── Core Rendering Tests ─────────────────────────────────────────────────────


class TestPptxRendererBasic:
    def test_render_produces_valid_bytes(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slide = TitleSlide(elements=[
            HeadingElement(
                content=HeadingContent(text="Hello World", level="h1"),
                position=Position(x=1.0, y=2.0, width=11.0, height=2.0),
            )
        ])
        lr = LayoutResult(slide=slide, positions={"title": ResolvedPosition(x=1.0, y=2.0, width=11.0, height=2.0)})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)

        assert isinstance(result, bytes)
        assert len(result) > 0

        # Verify it's a valid PPTX by loading it
        prs = PptxPresentation(io.BytesIO(result))
        assert prs is not None

    def test_slide_count_matches_layout_results(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slides = []
        layout_results = []
        for i in range(3):
            s = TitleSlide(elements=[
                HeadingElement(
                    content=HeadingContent(text=f"Slide {i}", level="h1"),
                    position=Position(x=1.0, y=2.0, width=11.0, height=2.0),
                )
            ])
            slides.append(s)
            layout_results.append(LayoutResult(slide=s, positions={}))

        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [s.model_dump() for s in slides],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, layout_results, theme)
        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_slide_dimensions_16_9(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slide = TitleSlide(elements=[])
        lr = LayoutResult(slide=slide, positions={})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)
        prs = PptxPresentation(io.BytesIO(result))

        # 16:9 = 13.333" x 7.5"
        expected_width = Inches(13.333)
        expected_height = Inches(7.5)
        assert abs(prs.slide_width - expected_width) < Inches(0.01)
        assert abs(prs.slide_height - expected_height) < Inches(0.01)


class TestSpeakerNotes:
    def test_speaker_notes_present(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slide = TitleSlide(
            elements=[
                HeadingElement(
                    content=HeadingContent(text="Title", level="h1"),
                    position=Position(x=1.0, y=2.0, width=11.0, height=2.0),
                )
            ],
            speaker_notes="These are my notes",
        )
        lr = LayoutResult(slide=slide, positions={})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)
        prs = PptxPresentation(io.BytesIO(result))

        notes_text = prs.slides[0].notes_slide.notes_text_frame.text
        assert "These are my notes" in notes_text


class TestTransitions:
    def test_fade_transition_xml(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slide = TitleSlide(
            elements=[],
            transition=Transition.FADE,
        )
        lr = LayoutResult(slide=slide, positions={})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)
        prs = PptxPresentation(io.BytesIO(result))

        slide_xml = prs.slides[0]._element.xml
        assert "transition" in slide_xml.lower() or "Transition" in slide_xml


class TestBackgroundColor:
    def test_background_from_slide_master(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slide = TitleSlide(elements=[])
        lr = LayoutResult(slide=slide, positions={})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)
        prs = PptxPresentation(io.BytesIO(result))

        # The background should be set (solid fill)
        bg = prs.slides[0].background
        assert bg.fill.type is not None


class TestMultiSlidePresentation:
    def test_title_bullets_table(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        title = TitleSlide(elements=[
            HeadingElement(
                content=HeadingContent(text="Title Slide", level="h1"),
                position=Position(x=1.0, y=2.0, width=11.0, height=2.0),
            ),
        ])
        bullets = BulletPointsSlide(elements=[
            BulletListElement(
                content=BulletListContent(items=["Item 1", "Item 2", "Item 3"]),
                position=Position(x=1.0, y=1.5, width=11.0, height=4.0),
            ),
        ])
        table = TableSlideSlide(elements=[
            TableElement(
                content=TableContent(
                    headers=["Name", "Value"],
                    rows=[["A", 1], ["B", 2]],
                ),
                position=Position(x=1.0, y=1.5, width=11.0, height=4.0),
            ),
        ])

        layout_results = [
            LayoutResult(slide=title, positions={}),
            LayoutResult(slide=bullets, positions={}),
            LayoutResult(slide=table, positions={}),
        ]

        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Multi"},
            "slides": [title.model_dump(), bullets.model_dump(), table.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, layout_results, theme)
        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 3


class TestEmptySlide:
    def test_no_elements_produces_valid_slide(self, theme):
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slide = TitleSlide(elements=[])
        lr = LayoutResult(slide=slide, positions={})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)
        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 1


# ── Parametrized test: all 23 universal slide types ──────────────────────────


_UNIVERSAL_SLIDE_CLASSES = [
    TitleSlide,
    AgendaSlide,
    SectionDividerSlide,
    KeyMessageSlide,
    BulletPointsSlide,
    TwoColumnTextSlide,
    ComparisonSlide,
    TimelineSlide,
    ProcessFlowSlide,
    OrgChartSlide,
    TeamSlide,
    QuoteSlide,
    ImageWithCaptionSlide,
    IconGridSlide,
    StatsCalloutSlide,
    TableSlideSlide,
    ChartSlideSlide,
    MatrixSlide,
    FunnelSlide,
    MapSlide,
    ThankYouSlide,
    AppendixSlide,
    QAndASlide,
]


@pytest.mark.parametrize(
    "slide_cls",
    _UNIVERSAL_SLIDE_CLASSES,
    ids=[cls.__name__ for cls in _UNIVERSAL_SLIDE_CLASSES],
)
def test_all_23_universal_slide_types_render_without_error(slide_cls, theme):
    """Every universal slide type should render without raising."""
    from modelforge.deck.rendering.pptx_renderer import PptxRenderer

    # Create slide with a simple heading element
    slide = slide_cls(elements=[
        HeadingElement(
            content=HeadingContent(text="Test", level="h1"),
            position=Position(x=1.0, y=1.0, width=11.0, height=2.0),
        ),
    ])
    lr = LayoutResult(slide=slide, positions={})
    ir = Presentation.model_validate({
        "schema_version": "1.0",
        "metadata": {"title": "Test"},
        "slides": [slide.model_dump()],
    })

    renderer = PptxRenderer()
    result = renderer.render(ir, [lr], theme)

    # Should produce valid bytes loadable by python-pptx
    prs = PptxPresentation(io.BytesIO(result))
    assert len(prs.slides) == 1
