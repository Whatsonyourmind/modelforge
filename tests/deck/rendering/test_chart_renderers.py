"""Tests for chart renderers — native editable PPTX charts via python-pptx."""

from __future__ import annotations

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from modelforge.deck.ir.charts.types import (
    BarChartData,
    StackedBarChartData,
    GroupedBarChartData,
    HorizontalBarChartData,
    LineChartData,
    MultiLineChartData,
    AreaChartData,
    StackedAreaChartData,
    PieChartData,
    DonutChartData,
    ScatterChartData,
    BubbleChartData,
    ComboChartData,
    RadarChartData,
    ChartDataSeries,
    WaterfallChartData,
)
from modelforge.deck.ir.elements.base import Position
from modelforge.deck.themes.types import (
    ResolvedTheme,
    ThemeColors,
    ThemeTypography,
    ThemeSpacing,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture
def pptx_slide():
    """Create a blank PPTX presentation slide for chart testing."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return slide


@pytest.fixture
def theme():
    """Minimal resolved theme for chart formatting."""
    return ResolvedTheme(
        name="test-theme",
        description="Test theme",
        colors=ThemeColors(
            primary="#2E86AB",
            secondary="#A23B72",
            accent="#F18F01",
            background="#1A1A2E",
            surface="#16213E",
            text_primary="#EAEAEA",
            text_secondary="#B0B0B0",
            text_muted="#808080",
            positive="#4CAF50",
            negative="#F44336",
            warning="#FF9800",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
        chart_colors=["#2E86AB", "#A23B72", "#F18F01", "#4CAF50", "#FF9800", "#9C27B0"],
    )


@pytest.fixture
def position():
    """Standard chart position on a slide."""
    return Position(x=1.0, y=1.5, width=11.0, height=5.0)


# ── Registry Tests ────────────────────────────────────────────────────────────


class TestChartRegistry:
    def test_chart_renderers_dict_exists(self):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        assert isinstance(CHART_RENDERERS, dict)

    def test_registry_contains_all_native_types(self):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        native_types = [
            "bar", "stacked_bar", "grouped_bar", "horizontal_bar",
            "line", "multi_line", "area", "stacked_area",
            "pie", "donut",
        ]
        for chart_type in native_types:
            assert chart_type in CHART_RENDERERS, f"Missing: {chart_type}"

    def test_registry_contains_placeholder_types(self):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        placeholder_types = [
            "waterfall", "funnel", "treemap", "tornado",
            "football_field", "sensitivity_table",
            "heatmap", "sankey", "gantt", "sunburst",
        ]
        for chart_type in placeholder_types:
            assert chart_type in CHART_RENDERERS, f"Missing placeholder: {chart_type}"

    def test_render_chart_function_exists(self):
        from modelforge.deck.rendering.chart_renderers import render_chart

        assert callable(render_chart)

    def test_render_chart_dispatches_bar(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import render_chart

        chart_data = BarChartData(
            categories=["Q1", "Q2", "Q3"],
            series=[ChartDataSeries(name="Revenue", values=[100, 120, 140])],
            title="Revenue by Quarter",
        )
        render_chart(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) > 0


# ── Bar Chart Tests ───────────────────────────────────────────────────────────


class TestBarChartRenderer:
    def test_bar_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series=[ChartDataSeries(name="Revenue", values=[100, 120, 140, 160])],
            title="Revenue by Quarter",
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) == 1

    def test_bar_chart_has_chart_object(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1", "Q2"],
            series=[ChartDataSeries(name="Rev", values=[100, 200])],
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart

    def test_stacked_bar_creates_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = StackedBarChartData(
            categories=["A", "B"],
            series=[
                ChartDataSeries(name="S1", values=[10, 20]),
                ChartDataSeries(name="S2", values=[30, 40]),
            ],
        )
        renderer = CHART_RENDERERS["stacked_bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart

    def test_grouped_bar_creates_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = GroupedBarChartData(
            categories=["A", "B"],
            series=[ChartDataSeries(name="S1", values=[10, 20])],
        )
        renderer = CHART_RENDERERS["grouped_bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_horizontal_bar_creates_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = HorizontalBarChartData(
            categories=["X", "Y"],
            series=[ChartDataSeries(name="S1", values=[50, 60])],
        )
        renderer = CHART_RENDERERS["horizontal_bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_bar_chart_with_title(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1"],
            series=[ChartDataSeries(name="Rev", values=[100])],
            title="My Chart Title",
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "My Chart Title"

    def test_bar_chart_without_title(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BarChartData(
            categories=["Q1"],
            series=[ChartDataSeries(name="Rev", values=[100])],
        )
        renderer = CHART_RENDERERS["bar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert not chart.has_title


# ── Line Chart Tests ──────────────────────────────────────────────────────────


class TestLineChartRenderer:
    def test_line_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = LineChartData(
            categories=["Jan", "Feb", "Mar"],
            series=[ChartDataSeries(name="Sales", values=[50, 60, 70])],
            title="Monthly Sales",
        )
        renderer = CHART_RENDERERS["line"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_multi_line_creates_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = MultiLineChartData(
            categories=["Jan", "Feb"],
            series=[
                ChartDataSeries(name="S1", values=[10, 20]),
                ChartDataSeries(name="S2", values=[30, 40]),
            ],
        )
        renderer = CHART_RENDERERS["multi_line"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_line_chart_with_title(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = LineChartData(
            categories=["Q1"],
            series=[ChartDataSeries(name="Rev", values=[100])],
            title="Line Title",
        )
        renderer = CHART_RENDERERS["line"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title


# ── Area Chart Tests ──────────────────────────────────────────────────────────


class TestAreaChartRenderer:
    def test_area_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = AreaChartData(
            categories=["Jan", "Feb", "Mar"],
            series=[ChartDataSeries(name="S1", values=[10, 20, 30])],
        )
        renderer = CHART_RENDERERS["area"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_stacked_area_creates_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = StackedAreaChartData(
            categories=["Jan", "Feb"],
            series=[
                ChartDataSeries(name="S1", values=[10, 20]),
                ChartDataSeries(name="S2", values=[30, 40]),
            ],
        )
        renderer = CHART_RENDERERS["stacked_area"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart


# ── Pie Chart Tests ───────────────────────────────────────────────────────────


class TestPieChartRenderer:
    def test_pie_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = PieChartData(
            labels=["A", "B", "C"],
            values=[50, 30, 20],
            title="Market Share",
        )
        renderer = CHART_RENDERERS["pie"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_pie_chart_has_title(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = PieChartData(
            labels=["A", "B"],
            values=[60, 40],
            title="Pie Title",
        )
        renderer = CHART_RENDERERS["pie"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "Pie Title"


# ── Donut Chart Tests ─────────────────────────────────────────────────────────


class TestDonutChartRenderer:
    def test_donut_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = DonutChartData(
            labels=["X", "Y", "Z"],
            values=[40, 35, 25],
        )
        renderer = CHART_RENDERERS["donut"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_donut_chart_is_doughnut_type(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = DonutChartData(
            labels=["A", "B"],
            values=[60, 40],
        )
        renderer = CHART_RENDERERS["donut"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        # Verify it is a doughnut chart via plot class name
        plot = chart.plots[0]
        assert "Doughnut" in type(plot).__name__


# ── Placeholder Tests ─────────────────────────────────────────────────────────


class TestPlaceholderChartRenderer:
    def test_placeholder_creates_shape(self, pptx_slide, theme, position):
        """Test PlaceholderChartRenderer directly (no longer in registry for waterfall)."""
        from modelforge.deck.rendering.chart_renderers.placeholder import PlaceholderChartRenderer

        renderer = PlaceholderChartRenderer()
        chart_data = WaterfallChartData(
            categories=["Start", "Change", "End"],
            values=[100, -20, 80],
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        assert len(pptx_slide.shapes) > 0

    def test_placeholder_has_text_label(self, pptx_slide, theme, position):
        """Test PlaceholderChartRenderer directly (no longer in registry for waterfall)."""
        from modelforge.deck.rendering.chart_renderers.placeholder import PlaceholderChartRenderer

        renderer = PlaceholderChartRenderer()
        chart_data = WaterfallChartData(
            categories=["Start", "End"],
            values=[100, 80],
        )
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_text_frame
        assert "waterfall" in shape.text_frame.text.lower()


# ── Scatter Chart Tests ───────────────────────────────────────────────────────


class TestScatterChartRenderer:
    def test_scatter_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = ScatterChartData(
            x_values=[1.0, 2.0, 3.0, 4.0, 5.0],
            series=[ChartDataSeries(name="Points", values=[10.0, 20.0, 15.0, 25.0, 30.0])],
            title="Scatter Plot",
        )
        renderer = CHART_RENDERERS["scatter"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_scatter_chart_is_native(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = ScatterChartData(
            x_values=[1.0, 2.0, 3.0],
            series=[ChartDataSeries(name="S1", values=[10.0, 20.0, 30.0])],
        )
        renderer = CHART_RENDERERS["scatter"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        # Verify it is a native chart object (not an image)
        assert shape.has_chart
        chart = shape.chart
        assert len(chart.series) >= 1

    def test_scatter_chart_with_multiple_series(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = ScatterChartData(
            x_values=[1.0, 2.0, 3.0],
            series=[
                ChartDataSeries(name="S1", values=[10.0, 20.0, 30.0]),
                ChartDataSeries(name="S2", values=[15.0, 25.0, 35.0]),
            ],
        )
        renderer = CHART_RENDERERS["scatter"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert len(chart.series) == 2

    def test_scatter_chart_with_title(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = ScatterChartData(
            x_values=[1.0, 2.0],
            series=[ChartDataSeries(name="S1", values=[10.0, 20.0])],
            title="My Scatter",
        )
        renderer = CHART_RENDERERS["scatter"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "My Scatter"


# ── Bubble Chart Tests ────────────────────────────────────────────────────────


class TestBubbleChartRenderer:
    def test_bubble_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BubbleChartData(
            x_values=[1.0, 2.0, 3.0],
            series=[ChartDataSeries(name="Bubbles", values=[10.0, 20.0, 30.0])],
            sizes=[5.0, 10.0, 15.0],
        )
        renderer = CHART_RENDERERS["bubble"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_bubble_chart_is_native(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BubbleChartData(
            x_values=[10.0, 20.0, 30.0],
            series=[ChartDataSeries(name="B1", values=[1.0, 2.0, 3.0])],
            sizes=[5.0, 10.0, 15.0],
        )
        renderer = CHART_RENDERERS["bubble"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart
        assert len(shape.chart.series) >= 1

    def test_bubble_chart_with_title(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = BubbleChartData(
            x_values=[1.0, 2.0],
            series=[ChartDataSeries(name="B1", values=[10.0, 20.0])],
            sizes=[5.0, 10.0],
            title="Bubble Title",
        )
        renderer = CHART_RENDERERS["bubble"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title


# ── Combo Chart Tests ─────────────────────────────────────────────────────────


class TestComboChartRenderer:
    def test_combo_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = ComboChartData(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series=[
                ChartDataSeries(name="Revenue", values=[100, 120, 140, 160]),
                ChartDataSeries(name="Cost", values=[80, 90, 100, 110]),
            ],
            line_series=[
                ChartDataSeries(name="Margin", values=[20, 30, 40, 50]),
            ],
            title="Revenue vs Margin",
        )
        renderer = CHART_RENDERERS["combo"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_combo_chart_has_all_series(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = ComboChartData(
            categories=["A", "B"],
            series=[ChartDataSeries(name="Bar1", values=[10, 20])],
            line_series=[ChartDataSeries(name="Line1", values=[15, 25])],
        )
        renderer = CHART_RENDERERS["combo"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        # Should have at least 2 series total (1 bar + 1 line)
        total_series = sum(len(plot.series) for plot in chart.plots)
        assert total_series >= 2

    def test_combo_chart_is_native(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = ComboChartData(
            categories=["X", "Y", "Z"],
            series=[ChartDataSeries(name="Bars", values=[10, 20, 30])],
            line_series=[ChartDataSeries(name="Lines", values=[15, 25, 35])],
        )
        renderer = CHART_RENDERERS["combo"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart


# ── Radar Chart Tests ─────────────────────────────────────────────────────────


class TestRadarChartRenderer:
    def test_radar_chart_creates_shape(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = RadarChartData(
            axes=["Speed", "Power", "Range", "Defense"],
            series=[
                ChartDataSeries(name="Model A", values=[8, 6, 7, 9]),
                ChartDataSeries(name="Model B", values=[6, 8, 5, 7]),
            ],
            title="Performance Comparison",
        )
        renderer = CHART_RENDERERS["radar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_radar_chart_is_native(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = RadarChartData(
            axes=["A", "B", "C"],
            series=[ChartDataSeries(name="S1", values=[1, 2, 3])],
        )
        renderer = CHART_RENDERERS["radar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart
        assert len(shape.chart.series) >= 1

    def test_radar_chart_with_title(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = RadarChartData(
            axes=["X", "Y", "Z"],
            series=[ChartDataSeries(name="S1", values=[5, 6, 7])],
            title="Radar Title",
        )
        renderer = CHART_RENDERERS["radar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "Radar Title"

    def test_radar_chart_with_multiple_series(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        chart_data = RadarChartData(
            axes=["Speed", "Power", "Range", "Defense"],
            series=[
                ChartDataSeries(name="S1", values=[8, 6, 7, 9]),
                ChartDataSeries(name="S2", values=[6, 8, 5, 7]),
            ],
        )
        renderer = CHART_RENDERERS["radar"]
        renderer.render(pptx_slide, chart_data, position, theme)
        chart = pptx_slide.shapes[0].chart
        assert len(chart.series) == 2


# ── Full Registry Completeness Test ──────────────────────────────────────────


class TestFullRegistry:
    def test_registry_has_all_14_chart_types(self):
        """After Task 2, registry should have all 14 native + 10 placeholder types."""
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        all_types = [
            # Native (14)
            "bar", "stacked_bar", "grouped_bar", "horizontal_bar",
            "line", "multi_line", "area", "stacked_area",
            "pie", "donut",
            "scatter", "bubble", "combo", "radar",
            # Placeholder (10)
            "waterfall", "funnel", "treemap", "tornado",
            "football_field", "sensitivity_table",
            "heatmap", "sankey", "gantt", "sunburst",
        ]
        for chart_type in all_types:
            assert chart_type in CHART_RENDERERS, f"Missing: {chart_type}"

    def test_render_chart_dispatches_scatter(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import render_chart

        chart_data = ScatterChartData(
            x_values=[1.0, 2.0, 3.0],
            series=[ChartDataSeries(name="S1", values=[10.0, 20.0, 30.0])],
        )
        render_chart(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart

    def test_render_chart_dispatches_radar(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import render_chart

        chart_data = RadarChartData(
            axes=["A", "B", "C"],
            series=[ChartDataSeries(name="S1", values=[1, 2, 3])],
        )
        render_chart(pptx_slide, chart_data, position, theme)
        assert pptx_slide.shapes[0].has_chart
