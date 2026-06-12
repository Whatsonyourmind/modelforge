"""Tests for finance slide renderers -- all 9 types, registry, and PptxRenderer integration."""

from __future__ import annotations

import io
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from modelforge.deck.ir.elements.base import Position
from modelforge.deck.ir.elements.data import (
    ChartElement,
    TableContent,
    TableElement,
)
from modelforge.deck.ir.elements.text import (
    BodyTextContent,
    BodyTextElement,
    HeadingContent,
    HeadingElement,
)
from modelforge.deck.themes.types import (
    ResolvedTheme,
    ThemeColors,
    ThemeSpacing,
    ThemeTypography,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture()
def theme() -> ResolvedTheme:
    """Minimal resolved theme for testing."""
    return ResolvedTheme(
        name="test-finance",
        description="Test theme for finance renderers",
        colors=ThemeColors(
            primary="#0A1E3D",
            secondary="#1A3A5C",
            accent="#FF6B35",
            background="#FFFFFF",
            surface="#F5F5F5",
            text_primary="#1A1A1A",
            text_secondary="#4A4A4A",
            text_muted="#8A8A8A",
            positive="#28A745",
            negative="#DC3545",
            warning="#FFC107",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
    )


@pytest.fixture()
def pptx_slide():
    """Create a blank PPTX slide for testing."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    return slide


def _make_comp_table_slide():
    """Create a CompTableSlide with sample financial data."""
    from modelforge.deck.ir.slides.finance import CompTableSlide

    table_elem = TableElement(
        content=TableContent(
            headers=["Company", "EV/EBITDA", "P/E", "Market Cap", "Revenue Growth"],
            rows=[
                ["Company A", 12.5, 18.3, 4200000000, 0.15],
                ["Company B", 10.2, 14.7, 2800000000, 0.22],
                ["Company C", 15.8, 22.1, 6100000000, 0.08],
                ["Company D", 11.0, 16.5, 3500000000, 0.18],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=11.8, height=4.5),
    )
    heading = HeadingElement(
        content=HeadingContent(text="Comparable Companies Analysis", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    return CompTableSlide(elements=[heading, table_elem])


def _make_dcf_slide():
    """Create a DcfSummarySlide with assumptions and sensitivity data."""
    from modelforge.deck.ir.slides.finance import DcfSummarySlide

    # Assumptions as a table element
    assumptions_table = TableElement(
        content=TableContent(
            headers=["Assumption", "Value"],
            rows=[
                ["WACC", 0.10],
                ["Terminal Growth Rate", 0.025],
                ["Projection Period", 5],
                ["Terminal Value Method", "Gordon Growth"],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=5.0, height=3.0),
    )

    # Sensitivity matrix as a table element with numeric-looking headers
    sensitivity_table = TableElement(
        content=TableContent(
            headers=["Discount Rate", "1.5%", "2.0%", "2.5%", "3.0%"],
            rows=[
                ["8.0%", 145, 152, 160, 170],
                ["9.0%", 130, 136, 142, 150],
                ["10.0%", 118, 123, 128, 135],
                ["11.0%", 108, 112, 116, 122],
            ],
        ),
        position=Position(x=6.5, y=1.5, width=6.0, height=3.0),
    )

    heading = HeadingElement(
        content=HeadingContent(text="DCF Valuation Summary", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )

    return DcfSummarySlide(
        elements=[heading, assumptions_table, sensitivity_table],
        discount_rate_range=[0.08, 0.09, 0.10, 0.11],
        terminal_growth_range=[0.015, 0.020, 0.025, 0.030],
    )


def _make_waterfall_slide():
    """Create a WaterfallChartSlide with chart data."""
    from modelforge.deck.ir.charts.types import WaterfallChartData
    from modelforge.deck.ir.slides.finance import WaterfallChartSlide

    chart_elem = ChartElement(
        chart_data=WaterfallChartData(
            categories=["Revenue", "COGS", "Gross Profit", "OpEx", "EBITDA", "Net Income"],
            values=[100, -40, 60, -25, 35, 28],
            title="Bridge to Net Income",
        ),
        position=Position(x=0.75, y=1.5, width=11.8, height=5.0),
    )
    heading = HeadingElement(
        content=HeadingContent(text="Earnings Bridge", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    return WaterfallChartSlide(elements=[heading, chart_elem], show_running_total=True)


# ── CompTableRenderer Tests ──────────────────────────────────────────────────


class TestCompTableRenderer:
    def test_produces_table_shape(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Should have at least one table shape
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 1, "CompTableRenderer should produce at least one table shape"

    def test_table_has_formatted_numbers(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Find the main data table
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert table_shapes, "Should have a table"
        table = table_shapes[0].table

        # Collect all cell text
        all_text = []
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                all_text.append(table.cell(row_idx, col_idx).text)

        joined = " ".join(all_text)
        # Should have formatted multiples (e.g., "12.5x")
        assert "x" in joined, "Should contain multiple formatted values (e.g., 12.5x)"
        # Should have formatted currency (e.g., "$4.2B")
        assert "$" in joined, "Should contain currency formatted values"

    def test_median_row_highlighted(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert table_shapes
        table = table_shapes[0].table

        # The last row should be the median footer row
        last_row_idx = len(table.rows) - 1
        cell_text = table.cell(last_row_idx, 0).text
        assert "median" in cell_text.lower() or "mean" in cell_text.lower(), \
            f"Last row should be median/mean, got: {cell_text}"

    def test_numeric_columns_right_aligned(self, pptx_slide, theme):
        from pptx.enum.text import PP_ALIGN

        from modelforge.deck.rendering.slide_renderers.comp_table import CompTableRenderer

        renderer = CompTableRenderer()
        ir_slide = _make_comp_table_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert table_shapes
        table = table_shapes[0].table

        # Check a numeric cell (row 1, col 1 = EV/EBITDA value)
        cell = table.cell(1, 1)
        for para in cell.text_frame.paragraphs:
            if para.alignment is not None:
                assert para.alignment == PP_ALIGN.RIGHT, "Numeric columns should be right-aligned"


# ── DcfSummaryRenderer Tests ─────────────────────────────────────────────────


class TestDcfSummaryRenderer:
    def test_produces_at_least_two_tables(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.dcf_summary import DcfSummaryRenderer

        renderer = DcfSummaryRenderer()
        ir_slide = _make_dcf_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 2, \
            f"DcfSummaryRenderer should produce at least 2 tables (assumptions + sensitivity), got {len(table_shapes)}"

    def test_sensitivity_matrix_has_gradient_colors(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.dcf_summary import DcfSummaryRenderer

        renderer = DcfSummaryRenderer()
        ir_slide = _make_dcf_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Find the sensitivity table (the one with more columns)
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 2

        # Get the larger table (sensitivity matrix)
        sensitivity_table = max(table_shapes, key=lambda s: len(s.table.columns))
        table = sensitivity_table.table

        # Check that data cells have fill colors applied (not all the same)
        fills = set()
        for row_idx in range(1, len(table.rows)):  # Skip header
            for col_idx in range(1, len(table.columns)):  # Skip row header
                cell = table.cell(row_idx, col_idx)
                if cell.fill.type is not None:
                    fills.add(str(cell.fill.fore_color.rgb))

        assert len(fills) >= 2, "Sensitivity matrix should have varied gradient colors"


# ── WaterfallSlideRenderer Tests ─────────────────────────────────────────────


class TestWaterfallSlideRenderer:
    @patch("modelforge.deck.rendering.chart_renderers.native_ib.WaterfallChartRenderer.render")
    def test_produces_slide_with_chart(self, mock_render, pptx_slide, theme):
        """WaterfallSlideRenderer should delegate chart rendering to the registry renderer."""
        from modelforge.deck.rendering.slide_renderers.waterfall_slide import WaterfallSlideRenderer

        renderer = WaterfallSlideRenderer()
        ir_slide = _make_waterfall_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # The renderer should have called WaterfallChartRenderer.render
        assert mock_render.called, "Should delegate to the registry waterfall renderer"

    def test_has_title(self, pptx_slide, theme):
        """WaterfallSlideRenderer should add a title text box."""
        from modelforge.deck.rendering.slide_renderers.waterfall_slide import WaterfallSlideRenderer

        # Mock the chart renderer to isolate the title text-box behavior
        with patch("modelforge.deck.rendering.chart_renderers.native_ib.WaterfallChartRenderer.render"):
            renderer = WaterfallSlideRenderer()
            ir_slide = _make_waterfall_slide()
            renderer.render(pptx_slide, ir_slide, theme)

        # Should have a text shape with the title
        text_shapes = [s for s in pptx_slide.shapes if s.has_text_frame]
        title_texts = [s.text_frame.text for s in text_shapes]
        assert any("Earnings Bridge" in t or "Bridge" in t for t in title_texts), \
            f"Should have title text box, found: {title_texts}"


# ── Registry Tests ───────────────────────────────────────────────────────────


class TestFinanceSlideRegistry:
    def test_registry_has_comp_table(self):
        from modelforge.deck.rendering.slide_renderers import FINANCE_SLIDE_RENDERERS
        from modelforge.deck.rendering.slide_renderers.comp_table import CompTableRenderer

        assert "comp_table" in FINANCE_SLIDE_RENDERERS
        assert isinstance(FINANCE_SLIDE_RENDERERS["comp_table"], CompTableRenderer)

    def test_registry_has_all_9_types(self):
        from modelforge.deck.rendering.slide_renderers import FINANCE_SLIDE_RENDERERS

        expected_types = {
            "comp_table",
            "dcf_summary",
            "waterfall_chart",
            "deal_overview",
            "returns_analysis",
            "capital_structure",
            "market_landscape",
            "investment_thesis",
            "risk_matrix",
        }
        assert set(FINANCE_SLIDE_RENDERERS.keys()) == expected_types, \
            f"Missing types: {expected_types - set(FINANCE_SLIDE_RENDERERS.keys())}"

    def test_render_finance_slide_dispatches(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers import render_finance_slide

        ir_slide = _make_comp_table_slide()
        result = render_finance_slide(pptx_slide, ir_slide, theme)
        assert result is True, "render_finance_slide should return True for comp_table"

    def test_render_finance_slide_returns_false_for_unknown(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers import render_finance_slide

        # Create a mock slide with unknown type
        mock_slide = MagicMock()
        mock_slide.slide_type = "title_slide"
        result = render_finance_slide(pptx_slide, mock_slide, theme)
        assert result is False, "render_finance_slide should return False for unknown types"


# ── Helper factories for Task 2 slide types ──────────────────────────────────


def _make_deal_overview_slide():
    """Create a DealOverviewSlide with metrics and status indicators."""
    from modelforge.deck.ir.slides.finance import DealOverviewSlide

    heading = HeadingElement(
        content=HeadingContent(text="Project Alpha -- Deal Overview", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    metrics_table = TableElement(
        content=TableContent(
            headers=["Metric", "Value"],
            rows=[
                ["Enterprise Value", 2500000000],
                ["Equity Value", 1800000000],
                ["Premium", 0.25],
                ["Revenue", 750000000],
                ["EBITDA", 180000000],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=5.5, height=3.5),
    )
    status_table = TableElement(
        content=TableContent(
            headers=["Workstream", "Status"],
            rows=[
                ["Due Diligence", "green"],
                ["Financing", "yellow"],
                ["Regulatory", "red"],
            ],
        ),
        position=Position(x=7.0, y=1.5, width=5.5, height=2.5),
    )
    return DealOverviewSlide(elements=[heading, metrics_table, status_table])


def _make_returns_analysis_slide():
    """Create a ReturnsAnalysisSlide with scenario table."""
    from modelforge.deck.ir.slides.finance import ReturnsAnalysisSlide

    heading = HeadingElement(
        content=HeadingContent(text="Returns Analysis", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    returns_table = TableElement(
        content=TableContent(
            headers=["Scenario", "IRR", "MOIC", "CoC", "Equity Value"],
            rows=[
                ["Base", 0.22, 2.5, 1.8, 500000000],
                ["Upside", 0.35, 3.2, 2.4, 720000000],
                ["Downside", 0.08, 1.4, 1.1, 280000000],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=11.8, height=3.5),
    )
    return ReturnsAnalysisSlide(elements=[heading, returns_table])


def _make_capital_structure_slide():
    """Create a CapitalStructureSlide with sources & uses."""
    from modelforge.deck.ir.slides.finance import CapitalStructureSlide

    heading = HeadingElement(
        content=HeadingContent(text="Capital Structure", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    sources = TableElement(
        content=TableContent(
            headers=["Source", "Amount", "% of Total"],
            rows=[
                ["Senior Debt", 600000000, 0.40],
                ["Mezzanine", 150000000, 0.10],
                ["Equity", 750000000, 0.50],
            ],
            footer_row=["Total", 1500000000, 1.0],
        ),
        position=Position(x=0.75, y=1.5, width=5.5, height=3.5),
    )
    uses = TableElement(
        content=TableContent(
            headers=["Use", "Amount", "% of Total"],
            rows=[
                ["Enterprise Value", 1400000000, 0.933],
                ["Fees & Expenses", 50000000, 0.033],
                ["Working Capital", 50000000, 0.033],
            ],
            footer_row=["Total", 1500000000, 1.0],
        ),
        position=Position(x=7.0, y=1.5, width=5.5, height=3.5),
    )
    return CapitalStructureSlide(elements=[heading, sources, uses])


def _make_market_landscape_slide():
    """Create a MarketLandscapeSlide with TAM/SAM/SOM and market data."""
    from modelforge.deck.ir.slides.finance import MarketLandscapeSlide

    heading = HeadingElement(
        content=HeadingContent(text="Market Landscape", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    tam_table = TableElement(
        content=TableContent(
            headers=["Segment", "Size", "Growth"],
            rows=[
                ["TAM", 50000000000, 0.12],
                ["SAM", 15000000000, 0.15],
                ["SOM", 3000000000, 0.20],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=5.5, height=3.0),
    )
    competitors = TableElement(
        content=TableContent(
            headers=["Competitor", "Market Share", "Revenue"],
            rows=[
                ["Alpha Corp", 0.30, 4500000000],
                ["Beta Inc", 0.20, 3000000000],
                ["Gamma Ltd", 0.15, 2250000000],
            ],
        ),
        position=Position(x=7.0, y=1.5, width=5.5, height=3.0),
    )
    return MarketLandscapeSlide(elements=[heading, tam_table, competitors])


def _make_investment_thesis_slide():
    """Create an InvestmentThesisSlide with thesis points."""
    from modelforge.deck.ir.slides.finance import InvestmentThesisSlide

    heading = HeadingElement(
        content=HeadingContent(text="Investment Thesis", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    point1 = BodyTextElement(
        content=BodyTextContent(text="Market leader with 35% share in growing $50B market"),
        position=Position(x=0.75, y=1.5, width=11.8, height=0.6),
    )
    point2 = BodyTextElement(
        content=BodyTextContent(text="Strong recurring revenue base with 95% retention"),
        position=Position(x=0.75, y=2.3, width=11.8, height=0.6),
    )
    point3 = BodyTextElement(
        content=BodyTextContent(text="Multiple expansion levers: pricing, cross-sell, M&A"),
        position=Position(x=0.75, y=3.1, width=11.8, height=0.6),
    )
    return InvestmentThesisSlide(elements=[heading, point1, point2, point3])


def _make_risk_matrix_slide():
    """Create a RiskMatrixSlide with risk items."""
    from modelforge.deck.ir.slides.finance import RiskMatrixSlide

    heading = HeadingElement(
        content=HeadingContent(text="Risk Assessment Matrix", level="h2"),
        position=Position(x=0.75, y=0.5, width=11.8, height=0.8),
    )
    risk_table = TableElement(
        content=TableContent(
            headers=["Risk", "Likelihood", "Impact"],
            rows=[
                ["Regulatory change", "High", "High"],
                ["Key person departure", "Medium", "High"],
                ["Market downturn", "Medium", "Medium"],
                ["Technology disruption", "Low", "High"],
                ["Supply chain risk", "Low", "Low"],
            ],
        ),
        position=Position(x=0.75, y=1.5, width=11.8, height=4.0),
    )
    return RiskMatrixSlide(
        elements=[heading, risk_table],
        axes_labels={"x": "Likelihood", "y": "Impact"},
    )


# ── DealOverviewRenderer Tests ──────────────────────────────────────────────


class TestDealOverviewRenderer:
    def test_produces_title_and_table(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.deal_overview import DealOverviewRenderer

        renderer = DealOverviewRenderer()
        ir_slide = _make_deal_overview_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Should have at least one table and a title text box
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        text_shapes = [s for s in pptx_slide.shapes if s.has_text_frame]
        assert len(table_shapes) >= 1, "Should produce at least one table"
        assert len(text_shapes) >= 1, "Should produce title text"

    def test_has_traffic_light_shapes(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.deal_overview import DealOverviewRenderer

        renderer = DealOverviewRenderer()
        ir_slide = _make_deal_overview_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Should have circular shapes for traffic lights
        all_shapes = list(pptx_slide.shapes)
        # Count shapes that are auto shapes (ovals)
        auto_shapes = [s for s in all_shapes if hasattr(s, "shape_type") and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        assert len(auto_shapes) >= 1, "Should have traffic light indicator shapes"


# ── ReturnsAnalysisRenderer Tests ────────────────────────────────────────────


class TestReturnsAnalysisRenderer:
    def test_produces_returns_table(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.returns_analysis import ReturnsAnalysisRenderer

        renderer = ReturnsAnalysisRenderer()
        ir_slide = _make_returns_analysis_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 1, "Should produce returns matrix table"

        # Check the table has IRR/MOIC formatted cells
        table = table_shapes[0].table
        all_text = []
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                all_text.append(table.cell(row_idx, col_idx).text)
        joined = " ".join(all_text)
        # Should have percentage formatted IRR values
        assert "%" in joined or "x" in joined, "Should contain formatted financial values"


# ── CapitalStructureRenderer Tests ───────────────────────────────────────────


class TestCapitalStructureRenderer:
    def test_produces_sources_and_uses(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.capital_structure import CapitalStructureRenderer

        renderer = CapitalStructureRenderer()
        ir_slide = _make_capital_structure_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 2, \
            f"Should produce sources and uses tables, got {len(table_shapes)}"


# ── MarketLandscapeRenderer Tests ────────────────────────────────────────────


class TestMarketLandscapeRenderer:
    def test_produces_tam_and_market_data(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.market_landscape import MarketLandscapeRenderer

        renderer = MarketLandscapeRenderer()
        ir_slide = _make_market_landscape_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Should have shapes for TAM/SAM/SOM and/or tables
        all_shapes = list(pptx_slide.shapes)
        assert len(all_shapes) >= 3, \
            f"Should produce TAM/SAM/SOM shapes and market data, got {len(all_shapes)} shapes"


# ── InvestmentThesisRenderer Tests ───────────────────────────────────────────


class TestInvestmentThesisRenderer:
    def test_produces_numbered_points(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.investment_thesis import InvestmentThesisRenderer

        renderer = InvestmentThesisRenderer()
        ir_slide = _make_investment_thesis_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        # Should have text boxes with numbered thesis points
        text_shapes = [s for s in pptx_slide.shapes if s.has_text_frame]
        all_text = " ".join(s.text_frame.text for s in text_shapes)
        assert "1." in all_text or "1)" in all_text, \
            f"Should have numbered thesis points, got: {all_text[:200]}"


# ── RiskMatrixRenderer Tests ─────────────────────────────────────────────────


class TestRiskMatrixRenderer:
    def test_produces_colored_grid(self, pptx_slide, theme):
        from modelforge.deck.rendering.slide_renderers.risk_matrix import RiskMatrixRenderer

        renderer = RiskMatrixRenderer()
        ir_slide = _make_risk_matrix_slide()
        renderer.render(pptx_slide, ir_slide, theme)

        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) >= 1, "Should produce risk grid table"

        # The grid table should have color-coded cells
        table = table_shapes[0].table
        fills = set()
        for row_idx in range(1, len(table.rows)):
            for col_idx in range(1, len(table.columns)):
                cell = table.cell(row_idx, col_idx)
                if cell.fill.type is not None:
                    fills.add(str(cell.fill.fore_color.rgb))
        assert len(fills) >= 2, "Risk matrix should have varied color-coded cells"


# ── PptxRenderer Integration Tests ──────────────────────────────────────────


class TestPptxRendererFinanceIntegration:
    def test_dispatches_finance_slides(self, theme):
        """PptxRenderer should dispatch finance slide types to FINANCE_SLIDE_RENDERERS."""
        from modelforge.deck.ir.presentation import Presentation
        from modelforge.deck.layout.types import LayoutResult
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        ir_slide = _make_comp_table_slide()
        lr = LayoutResult(slide=ir_slide, positions={})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Finance Test"},
            "slides": [ir_slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)

        # Should produce valid PPTX bytes
        assert isinstance(result, bytes) and len(result) > 0
        prs = PptxPresentation(io.BytesIO(result))
        slide = prs.slides[0]

        # Finance renderer should have created table shapes
        table_shapes = [s for s in slide.shapes if s.has_table]
        assert len(table_shapes) >= 1, "Finance slide should have been rendered with table"

    def test_non_finance_still_works(self, theme):
        """Non-finance slides should still render via element renderers."""
        from modelforge.deck.ir.presentation import Presentation
        from modelforge.deck.ir.slides.universal import TitleSlide
        from modelforge.deck.layout.types import LayoutResult
        from modelforge.deck.rendering.pptx_renderer import PptxRenderer

        slide = TitleSlide(elements=[
            HeadingElement(
                content=HeadingContent(text="Regular Title", level="h1"),
                position=Position(x=1.0, y=2.0, width=11.0, height=2.0),
            ),
        ])
        lr = LayoutResult(slide=slide, positions={})
        ir = Presentation.model_validate({
            "schema_version": "1.0",
            "metadata": {"title": "Test"},
            "slides": [slide.model_dump()],
        })

        renderer = PptxRenderer()
        result = renderer.render(ir, [lr], theme)
        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 1
