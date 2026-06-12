"""Tests for the native IB-style chart renderers (native_ib.py).

Covers the four signature IB chart types implemented with pure python-pptx
(no Plotly/kaleido/PNG): tornado, football_field, waterfall (invisible-base
stacked-bar trick) and sensitivity_table (native PPTX table with theme-driven
conditional shading). Includes determinism tests: rendering the same IR twice
into two .pptx files must yield identical chart XML parts.
"""

from __future__ import annotations

import re
import zipfile

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.util import Inches

from modelforge.deck.ir.charts.types import (
    FootballFieldChartData,
    SensitivityTableData,
    TornadoChartData,
    WaterfallChartData,
)
from modelforge.deck.ir.elements.base import Position
from modelforge.deck.themes.types import (
    ResolvedTheme,
    ThemeColors,
    ThemeSpacing,
    ThemeTypography,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────


@pytest.fixture
def pptx_slide():
    """Create a blank PPTX presentation slide for chart testing."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


@pytest.fixture
def theme():
    """Minimal resolved theme for chart formatting."""
    return ResolvedTheme(
        name="test-theme",
        description="Test theme",
        colors=ThemeColors(
            primary="#2E86AB",
            secondary="#A23B72",
            accent="#F18F01",
            background="#1A1A2E",
            surface="#16213E",
            text_primary="#EAEAEA",
            text_secondary="#B0B0B0",
            text_muted="#808080",
            positive="#4CAF50",
            negative="#F44336",
            warning="#FF9800",
        ),
        typography=ThemeTypography(
            heading_family="Arial",
            body_family="Calibri",
            mono_family="Consolas",
        ),
        spacing=ThemeSpacing(
            margin_top=0.5,
            margin_bottom=0.5,
            margin_left=0.75,
            margin_right=0.75,
            gutter=0.3,
            element_gap=0.2,
            section_gap=0.5,
        ),
        chart_colors=["#2E86AB", "#A23B72", "#F18F01", "#4CAF50", "#FF9800", "#9C27B0"],
    )


@pytest.fixture
def position():
    """Standard chart position on a slide."""
    return Position(x=1.0, y=1.5, width=11.0, height=5.0)


def _tornado_data():
    return TornadoChartData(
        categories=["WACC", "Growth"],
        low_values=[80, 95],
        high_values=[130, 105],
        base_value=100,
        title="IRR sensitivity drivers",
    )


def _football_field_data():
    return FootballFieldChartData(
        categories=["DCF", "Trading comps", "Precedent transactions"],
        low_values=[90, 80, 110],
        high_values=[140, 120, 160],
        title="Valuation summary",
    )


def _waterfall_data():
    return WaterfallChartData(
        categories=["Revenue", "COGS", "Gross Profit", "OpEx", "EBITDA", "Net Income"],
        values=[100, -40, 60, -25, 35, 28],
        title="Bridge to Net Income",
    )


def _sensitivity_data():
    return SensitivityTableData(
        row_header="WACC",
        col_header="Exit multiple",
        row_values=[8.0, 9.0, 10.0],
        col_values=[7.0, 8.0],
        data=[[120, 140], [100, 115], [85, 95]],
        title="EV sensitivity",
    )


def _series_by_name(chart):
    return {s.name: tuple(s.values) for s in chart.series}


# ── Registry ──────────────────────────────────────────────────────────────────


class TestNativeIbRegistry:
    def test_four_ib_types_no_longer_pending(self):
        from modelforge.deck.rendering.chart_renderers import (
            CHART_RENDERERS,
            FootballFieldChartRenderer,
            PendingNativeChartRenderer,
            SensitivityTableRenderer,
            TornadoChartRenderer,
            WaterfallChartRenderer,
        )

        expected = {
            "tornado": TornadoChartRenderer,
            "football_field": FootballFieldChartRenderer,
            "waterfall": WaterfallChartRenderer,
            "sensitivity_table": SensitivityTableRenderer,
        }
        for chart_type, cls in expected.items():
            renderer = CHART_RENDERERS[chart_type]
            assert isinstance(renderer, cls), chart_type
            assert not isinstance(renderer, PendingNativeChartRenderer), chart_type

    def test_remaining_six_still_pending(self):
        from modelforge.deck.rendering.chart_renderers import (
            CHART_RENDERERS,
            PendingNativeChartRenderer,
        )

        for chart_type in ("funnel", "treemap", "heatmap", "sankey", "gantt", "sunburst"):
            assert isinstance(CHART_RENDERERS[chart_type], PendingNativeChartRenderer)

    @pytest.mark.parametrize(
        "chart_data_factory",
        [_tornado_data, _football_field_data, _waterfall_data, _sensitivity_data],
        ids=["tornado", "football_field", "waterfall", "sensitivity_table"],
    )
    def test_render_chart_dispatches(self, pptx_slide, theme, position, chart_data_factory):
        from modelforge.deck.rendering.chart_renderers import render_chart

        render_chart(pptx_slide, chart_data_factory(), position, theme)
        assert len(pptx_slide.shapes) > 0


# ── Tornado ───────────────────────────────────────────────────────────────────


class TestTornadoChartRenderer:
    def test_creates_stacked_horizontal_bar_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["tornado"].render(pptx_slide, _tornado_data(), position, theme)
        assert len(pptx_slide.shapes) == 1
        shape = pptx_slide.shapes[0]
        assert shape.has_chart
        assert shape.chart.chart_type == XL_CHART_TYPE.BAR_STACKED

    def test_series_structure_base_downside_upside(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["tornado"].render(pptx_slide, _tornado_data(), position, theme)
        chart = pptx_slide.shapes[0].chart
        assert [s.name for s in chart.series] == ["_base", "Downside", "Upside"]

    def test_invisible_base_plus_spread_math(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["tornado"].render(pptx_slide, _tornado_data(), position, theme)
        chart = pptx_slide.shapes[0].chart
        series = _series_by_name(chart)

        # Sorted ascending by spread: Growth (10) first, WACC (50) last (top).
        assert list(chart.plots[0].categories) == ["Growth", "WACC"]
        # Growth: low 95, high 105, base 100 -> base seg 95, down 5, up 5.
        # WACC:   low 80, high 130, base 100 -> base seg 80, down 20, up 30.
        assert series["_base"] == (95.0, 80.0)
        assert series["Downside"] == (5.0, 20.0)
        assert series["Upside"] == (5.0, 30.0)

    def test_base_series_is_invisible(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["tornado"].render(pptx_slide, _tornado_data(), position, theme)
        chart = pptx_slide.shapes[0].chart
        base = list(chart.series)[0]
        assert base.format.fill.type == MSO_FILL.BACKGROUND  # noFill

    def test_spread_series_use_theme_pos_neg_colors(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["tornado"].render(pptx_slide, _tornado_data(), position, theme)
        series = list(pptx_slide.shapes[0].chart.series)
        assert str(series[1].format.fill.fore_color.rgb) == "F44336"  # negative
        assert str(series[2].format.fill.fore_color.rgb) == "4CAF50"  # positive

    def test_default_base_value_is_envelope_midpoint(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        data = TornadoChartData(
            categories=["X"], low_values=[10], high_values=[30], base_value=None
        )
        CHART_RENDERERS["tornado"].render(pptx_slide, data, position, theme)
        series = _series_by_name(pptx_slide.shapes[0].chart)
        # base = (10 + 30) / 2 = 20 -> down 10, up 10, invisible 10.
        assert series["_base"] == (10.0,)
        assert series["Downside"] == (10.0,)
        assert series["Upside"] == (10.0,)

    def test_legend_disabled_and_title_set(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["tornado"].render(pptx_slide, _tornado_data(), position, theme)
        chart = pptx_slide.shapes[0].chart
        assert not chart.has_legend
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "IRR sensitivity drivers"


# ── Football field ────────────────────────────────────────────────────────────


class TestFootballFieldChartRenderer:
    def test_creates_stacked_horizontal_bar_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["football_field"].render(
            pptx_slide, _football_field_data(), position, theme
        )
        shape = pptx_slide.shapes[0]
        assert shape.has_chart
        assert shape.chart.chart_type == XL_CHART_TYPE.BAR_STACKED

    def test_floating_low_to_high_bars(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["football_field"].render(
            pptx_slide, _football_field_data(), position, theme
        )
        chart = pptx_slide.shapes[0].chart
        assert [s.name for s in chart.series] == ["_base", "Range"]
        series = _series_by_name(chart)
        assert series["_base"] == (90.0, 80.0, 110.0)  # lows
        assert series["Range"] == (50.0, 40.0, 50.0)  # high - low

    def test_category_order_preserved(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["football_field"].render(
            pptx_slide, _football_field_data(), position, theme
        )
        chart = pptx_slide.shapes[0].chart
        assert list(chart.plots[0].categories) == [
            "DCF",
            "Trading comps",
            "Precedent transactions",
        ]

    def test_base_invisible_range_uses_theme_primary(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["football_field"].render(
            pptx_slide, _football_field_data(), position, theme
        )
        series = list(pptx_slide.shapes[0].chart.series)
        assert series[0].format.fill.type == MSO_FILL.BACKGROUND  # noFill base
        assert str(series[1].format.fill.fore_color.rgb) == "2E86AB"  # primary


# ── Waterfall ─────────────────────────────────────────────────────────────────


class TestWaterfallChartRenderer:
    def test_creates_stacked_column_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["waterfall"].render(pptx_slide, _waterfall_data(), position, theme)
        shape = pptx_slide.shapes[0]
        assert shape.has_chart
        assert shape.chart.chart_type == XL_CHART_TYPE.COLUMN_STACKED

    def test_series_structure(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["waterfall"].render(pptx_slide, _waterfall_data(), position, theme)
        chart = pptx_slide.shapes[0].chart
        assert [s.name for s in chart.series] == ["_base", "Increase", "Decrease", "Total"]

    def test_earnings_bridge_floating_bricks(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["waterfall"].render(pptx_slide, _waterfall_data(), position, theme)
        series = _series_by_name(pptx_slide.shapes[0].chart)
        # Revenue=total(100); COGS=-40 flow; Gross Profit=60 subtotal (==cum);
        # OpEx=-25 flow; EBITDA=35 subtotal (==cum); Net Income=total(28).
        assert series["_base"] == (0.0, 60.0, 0.0, 35.0, 0.0, 0.0)
        assert series["Increase"] == (0.0, 0.0, 0.0, 0.0, 0.0, 0.0)
        assert series["Decrease"] == (0.0, 40.0, 0.0, 25.0, 0.0, 0.0)
        assert series["Total"] == (100.0, 0.0, 60.0, 0.0, 35.0, 28.0)

    def test_positive_flow_brick(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        data = WaterfallChartData(
            categories=["Start", "Up", "Down", "End"],
            values=[100, 20, -30, 90],
        )
        CHART_RENDERERS["waterfall"].render(pptx_slide, data, position, theme)
        series = _series_by_name(pptx_slide.shapes[0].chart)
        assert series["_base"] == (0.0, 100.0, 90.0, 0.0)
        assert series["Increase"] == (0.0, 20.0, 0.0, 0.0)
        assert series["Decrease"] == (0.0, 0.0, 30.0, 0.0)
        assert series["Total"] == (100.0, 0.0, 0.0, 90.0)

    def test_base_invisible_and_pos_neg_total_colors(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["waterfall"].render(pptx_slide, _waterfall_data(), position, theme)
        series = list(pptx_slide.shapes[0].chart.series)
        assert series[0].format.fill.type == MSO_FILL.BACKGROUND  # noFill base
        assert str(series[1].format.fill.fore_color.rgb) == "4CAF50"  # positive
        assert str(series[2].format.fill.fore_color.rgb) == "F44336"  # negative
        assert str(series[3].format.fill.fore_color.rgb) == "2E86AB"  # primary


# ── Sensitivity table ─────────────────────────────────────────────────────────


class TestSensitivityTableRenderer:
    def test_renders_table_not_chart(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["sensitivity_table"].render(
            pptx_slide, _sensitivity_data(), position, theme
        )
        assert not any(s.has_chart for s in pptx_slide.shapes)
        table_shapes = [s for s in pptx_slide.shapes if s.has_table]
        assert len(table_shapes) == 1

    def test_table_dimensions(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["sensitivity_table"].render(
            pptx_slide, _sensitivity_data(), position, theme
        )
        table = next(s for s in pptx_slide.shapes if s.has_table).table
        assert len(table.rows) == 1 + 3  # header + row_values
        assert len(table.columns) == 1 + 2  # row label col + col_values

    def test_headers_and_corner_cell(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["sensitivity_table"].render(
            pptx_slide, _sensitivity_data(), position, theme
        )
        table = next(s for s in pptx_slide.shapes if s.has_table).table
        corner = table.cell(0, 0).text
        assert "WACC" in corner and "Exit multiple" in corner
        assert table.cell(0, 1).text == "7"
        assert table.cell(0, 2).text == "8"
        assert table.cell(1, 0).text == "8"
        assert table.cell(3, 0).text == "10"
        assert table.cell(1, 2).text == "140"

    def test_conditional_shading_extremes_use_theme_colors(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["sensitivity_table"].render(
            pptx_slide, _sensitivity_data(), position, theme
        )
        table = next(s for s in pptx_slide.shapes if s.has_table).table
        # min value 85 at data[2][0] -> cell(3, 1) -> exactly theme negative
        assert str(table.cell(3, 1).fill.fore_color.rgb) == "F44336"
        # max value 140 at data[0][1] -> cell(1, 2) -> exactly theme positive
        assert str(table.cell(1, 2).fill.fore_color.rgb) == "4CAF50"

    def test_conditional_shading_is_gradient(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["sensitivity_table"].render(
            pptx_slide, _sensitivity_data(), position, theme
        )
        table = next(s for s in pptx_slide.shapes if s.has_table).table
        fills = {
            str(table.cell(r, c).fill.fore_color.rgb)
            for r in range(1, len(table.rows))
            for c in range(1, len(table.columns))
        }
        assert len(fills) >= 3, "data cells should have varied gradient shading"

    def test_title_textbox_rendered(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        CHART_RENDERERS["sensitivity_table"].render(
            pptx_slide, _sensitivity_data(), position, theme
        )
        texts = [s.text_frame.text for s in pptx_slide.shapes if s.has_text_frame]
        assert any("EV sensitivity" in t for t in texts)

    def test_no_title_means_no_textbox(self, pptx_slide, theme, position):
        from modelforge.deck.rendering.chart_renderers import CHART_RENDERERS

        data = _sensitivity_data().model_copy(update={"title": None})
        CHART_RENDERERS["sensitivity_table"].render(pptx_slide, data, position, theme)
        assert not any(s.has_text_frame for s in pptx_slide.shapes)


# ── Determinism ───────────────────────────────────────────────────────────────


def _render_to_file(chart_data, theme, path):
    """Build a fresh presentation, render one chart, save to *path*."""
    from modelforge.deck.rendering.chart_renderers import render_chart

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_chart(slide, chart_data, Position(x=1.0, y=1.5, width=11.0, height=5.0), theme)
    prs.save(str(path))


def _chart_xml_parts(path) -> list[str]:
    """Extract all chart XML parts (sorted by part name) from a .pptx file."""
    with zipfile.ZipFile(str(path)) as zf:
        names = sorted(
            n for n in zf.namelist() if re.fullmatch(r"ppt/charts/chart\d+\.xml", n)
        )
        return [zf.read(n).decode("utf-8") for n in names]


def _slide_xml(path) -> str:
    with zipfile.ZipFile(str(path)) as zf:
        return zf.read("ppt/slides/slide1.xml").decode("utf-8")


class TestDeterminism:
    @pytest.mark.parametrize(
        "chart_data_factory",
        [_tornado_data, _football_field_data, _waterfall_data],
        ids=["tornado", "football_field", "waterfall"],
    )
    def test_same_ir_twice_yields_identical_chart_xml(
        self, tmp_path, theme, chart_data_factory
    ):
        path_a = tmp_path / "a.pptx"
        path_b = tmp_path / "b.pptx"
        _render_to_file(chart_data_factory(), theme, path_a)
        _render_to_file(chart_data_factory(), theme, path_b)

        parts_a = _chart_xml_parts(path_a)
        parts_b = _chart_xml_parts(path_b)
        assert len(parts_a) == 1, "expected exactly one chart XML part"
        assert parts_a == parts_b

    def test_sensitivity_table_slide_xml_identical(self, tmp_path, theme):
        path_a = tmp_path / "a.pptx"
        path_b = tmp_path / "b.pptx"
        _render_to_file(_sensitivity_data(), theme, path_a)
        _render_to_file(_sensitivity_data(), theme, path_b)

        assert _chart_xml_parts(path_a) == [] == _chart_xml_parts(path_b)
        assert _slide_xml(path_a) == _slide_xml(path_b)
