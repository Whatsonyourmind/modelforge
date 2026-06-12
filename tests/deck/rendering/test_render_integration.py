"""Integration tests for the full render pipeline: IR -> LayoutEngine -> PptxRenderer -> PPTX.

Tests the rendering pipeline in isolation (no DB, Redis, or S3 required).
Verifies that various IR configurations produce valid, loadable PPTX files.
"""

from __future__ import annotations

import io
import time

import pytest
from pptx import Presentation as PptxPresentation

from modelforge.deck.ir.presentation import Presentation
from modelforge.deck.layout.engine import LayoutEngine
from modelforge.deck.layout.text_measurer import TextMeasurer
from modelforge.deck.rendering import PptxRenderer
from modelforge.deck.themes.registry import ThemeRegistry


# E2 extraction note: deckforge.workers (Celery task layer) and the QA pass
# were not extracted. This local helper reproduces the core of
# deckforge.workers.tasks.render_pipeline (PPTX path) used by these tests;
# the QA report slot is returned as None and ignored by the tests.
def render_pipeline(presentation, output_format="pptx"):
    """Local stand-in for the DeckForge worker render pipeline (PPTX only)."""
    theme_registry = ThemeRegistry()
    layout_engine = LayoutEngine(TextMeasurer(), theme_registry)
    layout_results = layout_engine.layout_presentation(presentation)
    theme = theme_registry.get_theme(presentation.theme, presentation.brand_kit)
    pptx_bytes = PptxRenderer().render(presentation, layout_results, theme)
    return pptx_bytes, None


# ── Helpers ──────────────────────────────────────────────────────────────────


def _make_ir(slides: list[dict], theme: str = "executive-dark") -> dict:
    """Build a minimal IR dict with given slides."""
    return {
        "schema_version": "1.0",
        "metadata": {"title": "Integration Test"},
        "theme": theme,
        "slides": slides,
    }


def _title_slide(title: str = "DeckForge Demo", subtitle: str | None = None) -> dict:
    elements = [
        {"type": "heading", "content": {"text": title, "level": "h1"}},
    ]
    if subtitle:
        elements.append(
            {"type": "subheading", "content": {"text": subtitle}},
        )
    return {
        "slide_type": "title_slide",
        "elements": elements,
        "speaker_notes": "Welcome to the demo.",
        "transition": "fade",
    }


def _bullet_slide(heading: str = "Key Features", items: list[str] | None = None) -> dict:
    items = items or ["API-first design", "Native charts", "Theme system"]
    return {
        "slide_type": "bullet_points",
        "elements": [
            {"type": "heading", "content": {"text": heading, "level": "h2"}},
            {"type": "bullet_list", "content": {"items": items}},
        ],
    }


def _chart_slide(chart_type: str = "bar") -> dict:
    return {
        "slide_type": "chart_slide",
        "elements": [
            {"type": "heading", "content": {"text": "Revenue Analysis", "level": "h2"}},
            {
                "type": "chart",
                "chart_data": {
                    "chart_type": chart_type,
                    "title": "Q1 Revenue",
                    "categories": ["Jan", "Feb", "Mar"],
                    "series": [
                        {"name": "Revenue", "values": [100, 150, 200]},
                    ],
                },
            },
        ],
    }


# ── Pipeline Integration Tests ───────────────────────────────────────────────


class TestRenderPipeline:
    """Test the full render pipeline using render_pipeline() helper."""

    def test_simple_title_slide(self):
        """A single title slide produces valid PPTX bytes."""
        ir = Presentation.model_validate(_make_ir([_title_slide()]))
        result, _qa = render_pipeline(ir)

        assert isinstance(result, bytes)
        assert len(result) > 1000

        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 1

    def test_three_slide_presentation(self):
        """Title + bullets + chart produces a 3-slide PPTX."""
        ir = Presentation.model_validate(
            _make_ir([
                _title_slide("DeckForge Demo", "Automated Presentations"),
                _bullet_slide(),
                _chart_slide("bar"),
            ])
        )
        result, _qa = render_pipeline(ir)

        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 3
        assert len(result) > 5000

    def test_speaker_notes_survive_pipeline(self):
        """Speaker notes from IR appear in the generated PPTX."""
        ir = Presentation.model_validate(_make_ir([_title_slide()]))
        result, _qa = render_pipeline(ir)

        prs = PptxPresentation(io.BytesIO(result))
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Welcome to the demo" in notes_text

    def test_transition_survives_pipeline(self):
        """Transition attribute from IR is rendered into the PPTX."""
        ir = Presentation.model_validate(
            _make_ir([_title_slide()])
        )
        result, _qa = render_pipeline(ir)

        # If the PPTX loads without error and has a slide, the transition was applied
        # (python-pptx doesn't expose a transition read API, but the XML is there)
        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 1

    def test_ten_slide_presentation_under_5_seconds(self):
        """A 10-slide presentation renders in under 5 seconds."""
        slides = [_title_slide("Slide " + str(i)) for i in range(10)]
        ir = Presentation.model_validate(_make_ir(slides))

        start = time.monotonic()
        result, _qa = render_pipeline(ir)
        elapsed = time.monotonic() - start

        assert elapsed < 5.0, f"Rendering took {elapsed:.1f}s (threshold: 5s)"
        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 10

    def test_theme_background_applied(self):
        """Theme background color is applied to slides."""
        ir = Presentation.model_validate(
            _make_ir([_title_slide()], theme="executive-dark")
        )
        result, _qa = render_pipeline(ir)

        # Verify the PPTX is valid (background application is tested in
        # test_pptx_renderer.py; here we ensure the pipeline doesn't break it)
        prs = PptxPresentation(io.BytesIO(result))
        slide = prs.slides[0]
        bg = slide.background
        assert bg is not None

    def test_bullet_list_renders_correctly(self):
        """Bullet list elements produce shapes in the PPTX."""
        ir = Presentation.model_validate(
            _make_ir([_bullet_slide()])
        )
        result, _qa = render_pipeline(ir)

        prs = PptxPresentation(io.BytesIO(result))
        slide = prs.slides[0]
        # Should have at least heading + bullet list shapes
        assert len(slide.shapes) >= 1

    def test_chart_slide_renders(self):
        """Chart element produces a chart shape in the PPTX."""
        ir = Presentation.model_validate(
            _make_ir([_chart_slide("bar")])
        )
        result, _qa = render_pipeline(ir)

        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 1
        # Chart slide should have shapes (heading + chart)
        assert len(prs.slides[0].shapes) >= 1


class TestRenderPipelineComponents:
    """Test the pipeline components individually for correctness."""

    def test_layout_engine_produces_results(self):
        """LayoutEngine produces a LayoutResult per slide."""
        registry = ThemeRegistry()
        measurer = TextMeasurer()
        engine = LayoutEngine(measurer, registry)

        ir = Presentation.model_validate(
            _make_ir([_title_slide(), _bullet_slide()])
        )

        results = engine.layout_presentation(ir)
        assert len(results) >= 2

    def test_renderer_accepts_layout_results(self):
        """PptxRenderer accepts LayoutResults and produces bytes."""
        registry = ThemeRegistry()
        measurer = TextMeasurer()
        engine = LayoutEngine(measurer, registry)

        ir = Presentation.model_validate(_make_ir([_title_slide()]))
        results = engine.layout_presentation(ir)
        theme = registry.get_theme("executive-dark")

        renderer = PptxRenderer()
        pptx_bytes = renderer.render(ir, results, theme)

        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0

    def test_pipeline_with_mixed_slide_types(self):
        """Pipeline handles different slide types without errors."""
        slides = [
            _title_slide("Welcome", "Subtitle"),
            _bullet_slide("Features", ["A", "B", "C"]),
            _chart_slide("line"),
            {
                "slide_type": "section_divider",
                "elements": [
                    {"type": "heading", "content": {"text": "Section 2", "level": "h1"}},
                ],
            },
            _bullet_slide("Summary", ["Point 1", "Point 2"]),
        ]
        ir = Presentation.model_validate(_make_ir(slides))
        result, _qa = render_pipeline(ir)

        prs = PptxPresentation(io.BytesIO(result))
        assert len(prs.slides) == 5
