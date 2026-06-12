"""Slide-content gate for the certified deck pipeline.

A certified deck must NEVER ship hollow (title-only) slides. This module
builds the real certified decks (ic_memo + teaser) from an example
sponsor_lbo spec and asserts, per slide:

    1. >= 2 visible shapes (title + content), unless the slide is on the
       explicit whitelist (cover slide only);
    2. zero "Skipping element" / "No table element" / "no position assigned"
       notices were emitted during the build (silent content loss is dead:
       the renderer now RAISES on unpositioned elements);
    3. the QA structural checker reports no CRITICAL (severity="error")
       findings on the composed presentation.

Also locks the loud-failure contract: an element without a layout position
makes the renderer raise DeckRenderError instead of dropping content.
"""

from __future__ import annotations

import logging
from pathlib import Path

import pytest
import yaml

REPO_ROOT = Path(__file__).resolve().parents[2]
SPEC_PATH = REPO_ROOT / "examples" / "sponsor_lbo_techco.yaml"

# Slides allowed to carry fewer than 2 visible shapes (1-based index).
# Only the cover slide qualifies — and with the byline now positioned it
# carries 3 shapes anyway; the whitelist is belt-and-braces.
_WHITELISTED_SLIDES = {1}

# Log fragments that indicate silent content loss.
_CONTENT_LOSS_FRAGMENTS = (
    "Skipping element",
    "No table element",
    "No table elements",
    "no position assigned",
    "Failed to render element",
)


class _RecordCollector(logging.Handler):
    """Collects every log record emitted under modelforge.deck."""

    def __init__(self) -> None:
        super().__init__(level=logging.DEBUG)
        self.records: list[logging.LogRecord] = []

    def emit(self, record: logging.LogRecord) -> None:  # pragma: no cover
        self.records.append(record)


# ─────────────────────────────────────────────────────────────────────────────
# Fixtures: one certified workbook + both decks, logs captured
# ─────────────────────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def built_workbook(tmp_path_factory) -> Path:
    """Build the sponsor_lbo_techco example exactly as `modelforge build` ships it."""
    out_dir = tmp_path_factory.mktemp("slide_content")
    from modelforge.cli import _inject_trust_moat_and_finish, _load_spec_class
    from modelforge.templates import build_model

    spec_bytes = SPEC_PATH.read_bytes()
    raw = yaml.safe_load(spec_bytes)
    spec_class = _load_spec_class(raw["model_type"])
    spec = spec_class.model_validate(raw)

    xlsx_out = out_dir / "sponsor_lbo_techco.xlsx"
    xlsx, _graph = build_model(
        spec, xlsx_out,
        spec_source_bytes=spec_bytes,
        spec_source_path=SPEC_PATH,
    )
    _inject_trust_moat_and_finish(xlsx, spec, spec_bytes, SPEC_PATH, quiet=True)
    xlsx = Path(xlsx)
    assert xlsx.exists()
    assert xlsx.with_suffix(".manifest.json").exists()
    return xlsx


@pytest.fixture(scope="module")
def built_decks(built_workbook: Path):
    """Build BOTH certified decks while capturing every deck log record."""
    from modelforge.deck.pipeline import build_deck_from_workbook

    collector = _RecordCollector()
    deck_logger = logging.getLogger("modelforge.deck")
    old_level = deck_logger.level
    deck_logger.addHandler(collector)
    deck_logger.setLevel(logging.DEBUG)
    try:
        results = {
            deck_type: build_deck_from_workbook(
                built_workbook,
                deck_type=deck_type,
                out_path=built_workbook.with_name(f"content_{deck_type}.pptx"),
            )
            for deck_type in ("ic_memo", "teaser")
        }
    finally:
        deck_logger.removeHandler(collector)
        deck_logger.setLevel(old_level)
    return results, collector.records


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────


def _visible_shape_count(slide) -> int:
    """Count shapes that actually display something."""
    count = 0
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            count += 1
        elif getattr(shape, "has_chart", False):
            count += 1
        elif shape.has_text_frame:
            if shape.text_frame.text.strip():
                count += 1
        else:
            # Pictures, auto shapes (KPI cards, traffic lights), etc.
            count += 1
    return count


def _shape_counts(pptx_path: Path) -> list[int]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    return [_visible_shape_count(s) for s in prs.slides]


# ─────────────────────────────────────────────────────────────────────────────
# The gate
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.parametrize("deck_type", ["ic_memo", "teaser"])
def test_every_slide_has_title_and_content(built_decks, deck_type):
    """No hollow slides: every slide carries >= 2 visible shapes."""
    results, _records = built_decks
    counts = _shape_counts(results[deck_type].pptx_path)

    report = ", ".join(
        f"slide {i}: {c} shapes" for i, c in enumerate(counts, start=1)
    )
    hollow = [
        i for i, c in enumerate(counts, start=1)
        if c < 2 and i not in _WHITELISTED_SLIDES
    ]
    assert not hollow, (
        f"{deck_type}: hollow slides {hollow} (need >= 2 visible shapes). "
        f"Full map: {report}"
    )


def test_zero_content_loss_notices_during_build(built_decks):
    """The build emits no skip/no-table/unpositioned notices at any level."""
    _results, records = built_decks
    offending = [
        f"{r.name}:{r.levelname}: {r.getMessage()}"
        for r in records
        if any(frag in r.getMessage() for frag in _CONTENT_LOSS_FRAGMENTS)
    ]
    assert not offending, f"Silent content-loss notices emitted: {offending}"


@pytest.mark.parametrize("deck_type", ["ic_memo", "teaser"])
def test_structural_checker_no_critical_findings(built_workbook, deck_type):
    """The QA structural checker is clean on the composed presentation."""
    from modelforge.deck.adapter import adapt_workbook
    from modelforge.deck.layout import LayoutEngine
    from modelforge.deck.layout.text_measurer import TextMeasurer
    from modelforge.deck.pipeline import build_certification_slide
    from modelforge.deck.qa.checkers.structural import StructuralChecker
    from modelforge.deck.themes.registry import ThemeRegistry

    wf = adapt_workbook(built_workbook)
    if deck_type == "ic_memo":
        from modelforge.deck.compose import compose_ic_memo

        presentation = compose_ic_memo(wf.deal_facts())
    else:
        from modelforge.deck.compose import compose_teaser

        presentation = compose_teaser(wf.teaser_facts())
    presentation = presentation.model_copy(
        update={"slides": [*presentation.slides, build_certification_slide(wf)]}
    )

    registry = ThemeRegistry()
    engine = LayoutEngine(TextMeasurer(), registry)
    layout_results = engine.layout_presentation(presentation)
    theme = registry.get_theme(presentation.theme, None)

    issues = StructuralChecker().check(presentation, layout_results, theme)
    critical = [i for i in issues if i.severity == "error"]
    assert not critical, (
        f"{deck_type}: structural checker CRITICAL findings: "
        f"{[(i.type, i.slide_index, i.message) for i in critical]}"
    )


def test_no_data_slides_are_omitted_not_hollow(built_decks):
    """The LBO workbook has no comps and no sensitivity grid — those slides
    must be absent from the IC memo, not rendered title-only."""
    results, _records = built_decks
    # 12-slide spine - comps - sensitivity + certification appendix
    assert results["ic_memo"].slide_count == 11
    assert results["teaser"].slide_count == 7


# ─────────────────────────────────────────────────────────────────────────────
# Loud-failure contract (unit level, no workbook build)
# ─────────────────────────────────────────────────────────────────────────────


def test_unpositioned_element_raises_deck_render_error():
    """An element with no layout position must raise, never be dropped."""
    from modelforge.deck.ir.elements.base import Position
    from modelforge.deck.ir.elements.text import (
        BodyTextContent,
        BodyTextElement,
        HeadingContent,
        HeadingElement,
    )
    from modelforge.deck.ir.presentation import Presentation
    from modelforge.deck.ir.slides.universal import KeyMessageSlide
    from modelforge.deck.layout.types import LayoutResult
    from modelforge.deck.rendering import DeckRenderError, PptxRenderer
    from modelforge.deck.themes.registry import ThemeRegistry

    slide = KeyMessageSlide(elements=[
        HeadingElement(
            content=HeadingContent(text="Recommendation", level="h1"),
            position=Position(x=1.0, y=1.0, width=11.0, height=1.0),
        ),
        BodyTextElement(
            content=BodyTextContent(text="This rationale must not vanish."),
            # no position -> must raise, not skip
        ),
    ])
    ir = Presentation.model_validate({
        "schema_version": "1.0",
        "metadata": {"title": "Loud failure"},
        "slides": [slide.model_dump()],
    })
    theme = ThemeRegistry().get_theme("finance-pro", None)

    with pytest.raises(DeckRenderError, match="no position assigned"):
        PptxRenderer().render(ir, [LayoutResult(slide=slide, positions={})], theme)
