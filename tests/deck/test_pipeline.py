"""End-to-end tests for the certified deck pipeline.

Covers the load-bearing chain:

    spec → build (trust/moat + deterministic finishing + manifest)
         → adapt (fail-closed: manifest verify + CERTIFIED audit)
         → compose (+ mandatory Certification & Red Flags slide)
         → render → deterministic stamp (hashes in core props)

One sponsor_lbo workbook is built ONCE per module (session-expensive:
template build + trust/moat + recalc audits) and shared by every test.
"""

from __future__ import annotations

import hashlib
import json
import shutil
import sys
from pathlib import Path
from types import SimpleNamespace

import pytest
import yaml

REPO_ROOT = Path(__file__).resolve().parents[2]
SPEC_PATH = REPO_ROOT / "examples" / "sponsor_lbo_us_saas.yaml"


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


# ─────────────────────────────────────────────────────────────────────────────
# Fixtures
# ─────────────────────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def built_workbook(tmp_path_factory) -> Path:
    """Build the sponsor_lbo example exactly as `modelforge build` ships it."""
    out_dir = tmp_path_factory.mktemp("deck_pipeline")
    from modelforge.cli import _inject_trust_moat_and_finish, _load_spec_class
    from modelforge.templates import build_model

    spec_bytes = SPEC_PATH.read_bytes()
    raw = yaml.safe_load(spec_bytes)
    spec_class = _load_spec_class(raw["model_type"])
    spec = spec_class.model_validate(raw)

    xlsx_out = out_dir / "sponsor_lbo_us_saas.xlsx"
    xlsx, _graph = build_model(
        spec, xlsx_out,
        spec_source_bytes=spec_bytes,
        spec_source_path=SPEC_PATH,
    )
    _inject_trust_moat_and_finish(xlsx, spec, spec_bytes, SPEC_PATH, quiet=True)
    xlsx = Path(xlsx)
    assert xlsx.exists()
    assert xlsx.with_suffix(".manifest.json").exists()
    return xlsx


@pytest.fixture(scope="module")
def ic_deck(built_workbook: Path):
    """One IC-memo deck built from the certified workbook."""
    from modelforge.deck.pipeline import build_deck_from_workbook

    return build_deck_from_workbook(
        built_workbook,
        deck_type="ic_memo",
        out_path=built_workbook.with_name("deck_a.pptx"),
    )


# ─────────────────────────────────────────────────────────────────────────────
# (a) sponsor_lbo example → deck builds, Certification slide, SHAs in props
# ─────────────────────────────────────────────────────────────────────────────


class TestDeckBuilds:
    def test_deck_file_exists_and_certified(self, ic_deck):
        assert ic_deck.pptx_path.exists()
        assert ic_deck.pptx_path.stat().st_size > 10_000
        assert ic_deck.audit_verdict == "CERTIFIED"
        assert ic_deck.template == "sponsor_lbo"
        # 12 composer slides, minus comps + sensitivity (no data in an LBO
        # workbook -> OMITTED, never hollow), + mandatory certification slide
        assert ic_deck.slide_count == 11

    def test_certification_slide_is_last_and_complete(self, ic_deck):
        from pptx import Presentation

        prs = Presentation(str(ic_deck.pptx_path))
        assert len(prs.slides) == 11
        last = prs.slides[-1]
        texts = []
        for shape in last.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        blob = " | ".join(texts)
        assert "Certification & Red Flags" in blob
        assert "CERTIFIED" in blob
        assert ic_deck.workbook_sha256 in blob
        assert ic_deck.spec_sha256 in blob
        # Trust-Layer entries (or the explicit none-raised line)
        assert ("none raised" in blob) or ("Red flag" in blob)

    def test_core_properties_carry_both_shas(self, ic_deck, built_workbook):
        from modelforge.deck.determinism import read_pptx_stamp

        stamp = read_pptx_stamp(ic_deck.pptx_path)
        assert stamp["creator"] == "modelforge"
        assert stamp["workbook_sha256"] == ic_deck.workbook_sha256
        assert stamp["spec_sha256"] == ic_deck.spec_sha256
        # And they match the manifest sidecar on disk
        manifest = json.loads(
            built_workbook.with_suffix(".manifest.json").read_text(encoding="utf-8")
        )
        assert stamp["workbook_sha256"] == manifest["workbook_sha256"]
        assert stamp["spec_sha256"] == manifest["spec_sha256"]
        # created/modified are derived from the workbook hash, not wall clock
        assert stamp["created"] == stamp["modified"]

    def test_facts_carry_source_cell_refs(self, ic_deck):
        refs = ic_deck.source_refs
        for key in ("irr_strategic", "moic_strategic", "entry_ev",
                    "sponsor_equity", "senior_debt", "total_uses",
                    "exit_year"):
            assert key in refs, f"missing lineage for {key}"
            assert "!" in refs[key], f"{key} ref {refs[key]!r} is not sheet!cell"

    def test_speaker_notes_carry_lineage(self, ic_deck):
        from pptx import Presentation

        prs = Presentation(str(ic_deck.pptx_path))
        # Returns-summary slide (index 4) shows IRR/MoIC → notes cite cells.
        notes = prs.slides[4].notes_slide.notes_text_frame.text
        assert "irr_strategic=SourcesUses!" in notes
        assert "moic_strategic=SourcesUses!" in notes

    def test_teaser_deck_builds(self, built_workbook):
        from modelforge.deck.pipeline import build_deck_from_workbook

        res = build_deck_from_workbook(
            built_workbook,
            deck_type="teaser",
            out_path=built_workbook.with_name("deck_teaser.pptx"),
        )
        assert res.pptx_path.exists()
        # 6 composer slides + certification slide
        assert res.slide_count == 7
        from modelforge.deck.determinism import read_pptx_stamp

        stamp = read_pptx_stamp(res.pptx_path)
        assert stamp["workbook_sha256"] == res.workbook_sha256


# ─────────────────────────────────────────────────────────────────────────────
# (b) byte-stability: same workbook → two decks → identical SHA-256
# ─────────────────────────────────────────────────────────────────────────────


class TestByteStability:
    def test_same_workbook_two_decks_byte_identical(self, built_workbook, ic_deck):
        from modelforge.deck.pipeline import build_deck_from_workbook

        second = build_deck_from_workbook(
            built_workbook,
            deck_type="ic_memo",
            out_path=built_workbook.with_name("deck_b.pptx"),
        )
        assert _sha256(ic_deck.pptx_path) == _sha256(second.pptx_path)


# ─────────────────────────────────────────────────────────────────────────────
# (c) fail-closed: missing manifest / tampered workbook → refused
# ─────────────────────────────────────────────────────────────────────────────


class TestFailClosed:
    def test_missing_manifest_refused(self, built_workbook, tmp_path):
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        orphan = tmp_path / "orphan.xlsx"
        shutil.copy(built_workbook, orphan)  # no manifest sidecar copied
        with pytest.raises(DeckAdapterError, match="[Mm]anifest"):
            build_deck_from_workbook(orphan)

    def test_tampered_workbook_refused(self, built_workbook, tmp_path):
        from openpyxl import load_workbook

        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        tampered = tmp_path / "tampered.xlsx"
        shutil.copy(built_workbook, tampered)
        shutil.copy(
            built_workbook.with_suffix(".manifest.json"),
            tampered.with_suffix(".manifest.json"),
        )
        # Flip a hardcoded input value (senior debt) — bytes no longer match
        # the manifest's workbook_sha256.
        wb = load_workbook(tampered)
        ws = wb["SourcesUses"]
        ws["D7"] = (ws["D7"].value or 0) + 1.0
        wb.save(tampered)

        with pytest.raises(DeckAdapterError,
                           match="verification FAILED|not the bytes"):
            build_deck_from_workbook(tampered)

    def test_uncertified_workbook_refused(self, built_workbook, tmp_path):
        """A workbook whose manifest verifies but whose audit is not
        CERTIFIED (cached #REF! literal) must be refused at step 2."""
        from openpyxl import load_workbook

        from modelforge.analytics.manifest import write_manifest
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        bad = tmp_path / "uncertified.xlsx"
        shutil.copy(built_workbook, bad)
        wb = load_workbook(bad)
        ws = wb["SourcesUses"]
        ws["Z2"] = "#REF!"  # cached error literal → audit FAIL
        wb.save(bad)
        # Re-manifest the tampered bytes so the manifest gate PASSES and the
        # refusal provably comes from the certification audit.
        write_manifest(bad, SimpleNamespace(), spec_source_bytes=b"tampered")

        with pytest.raises(DeckAdapterError, match="not\\s+CERTIFIED"):
            build_deck_from_workbook(bad)

    def test_unknown_deck_type_refused(self, built_workbook):
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        with pytest.raises(DeckAdapterError, match="deck type"):
            build_deck_from_workbook(built_workbook, deck_type="lp_update")

    def test_unknown_theme_refused(self, built_workbook):
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        with pytest.raises(DeckAdapterError, match="Unknown theme"):
            build_deck_from_workbook(built_workbook, theme="no-such-theme")


# ─────────────────────────────────────────────────────────────────────────────
# (d) unsupported template → friendly error
# ─────────────────────────────────────────────────────────────────────────────


class TestUnsupportedTemplate:
    def test_unsupported_template_friendly_error(self, tmp_path):
        """A verified + certifiable workbook of an unmapped template is
        refused with the friendly supported-list message."""
        from openpyxl import Workbook

        from modelforge.analytics.manifest import write_manifest
        from modelforge.deck.pipeline import (
            DeckAdapterError, build_deck_from_workbook,
        )

        xlsx = tmp_path / "re_model.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Cover"
        ws["A1"] = "Real-estate model (text only, zero numeric cells)"
        rf = wb.create_sheet("RedFlags")
        rf["A1"] = "ModelForge Trust Layer — Red Flags"
        rf["A3"] = "Template: real_estate"
        rf["B8"] = "ALL CLEAR"
        wb.save(xlsx)
        write_manifest(xlsx, SimpleNamespace(),
                       spec_source_bytes=b"model_type: real_estate\n")

        with pytest.raises(DeckAdapterError) as exc:
            build_deck_from_workbook(xlsx)
        msg = str(exc.value)
        assert "real_estate" in msg
        assert "not deck-mappable yet" in msg
        assert "sponsor_lbo" in msg

    def test_detect_template_reads_redflags_line(self, tmp_path):
        from openpyxl import Workbook, load_workbook

        from modelforge.deck.adapter import detect_template

        xlsx = tmp_path / "probe.xlsx"
        wb = Workbook()
        wb.active.title = "Cover"
        rf = wb.create_sheet("RedFlags")
        rf["A3"] = "Template: unitranche"
        wb.save(xlsx)
        assert detect_template(load_workbook(xlsx)) == "unitranche"


# ─────────────────────────────────────────────────────────────────────────────
# Determinism unit tests (no workbook build required)
# ─────────────────────────────────────────────────────────────────────────────


class TestStampUnit:
    def test_derived_instant_is_pure_function_of_hash(self):
        from modelforge.deck.determinism import derive_build_datetime

        sha = "ab" * 32
        assert derive_build_datetime(sha) == derive_build_datetime(sha)
        assert (derive_build_datetime("00" * 32)
                != derive_build_datetime("ff" * 32))

    def test_derived_instant_is_in_fixed_past_window(self):
        # Clamped into calendar year 2020: never a future created/modified
        # date in the .pptx properties, regardless of the hash value.
        from modelforge.deck.determinism import derive_build_datetime

        for sha in ("00" * 32, "ff" * 32, "ab" * 32, "deadbeef" + "00" * 28):
            assert derive_build_datetime(sha).year == 2020

    def test_stamp_is_idempotent(self, ic_deck, tmp_path):
        from modelforge.deck.determinism import stamp_pptx

        copy = tmp_path / "restamp.pptx"
        shutil.copy(ic_deck.pptx_path, copy)
        before = _sha256(copy)
        stamp_pptx(copy,
                   workbook_sha256=ic_deck.workbook_sha256,
                   spec_sha256=ic_deck.spec_sha256,
                   title=None)
        assert _sha256(copy) == before
